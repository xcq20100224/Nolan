# -*- coding: utf-8 -*-
"""
N3 真机验收脚本：实跑 make_ppt 并用 python-pptx 提取验证。
用法（jarvis/ 目录下）：
    python accept_n3_courseware.py 1   # 一元二次方程课件（课件型）
    python accept_n3_courseware.py 2   # 开学季（讲话动员型，跳过研究）
    python accept_n3_courseware.py 3   # 全球新能源汽车产业格局（报告型）
每个病例输出判定结果、大纲结构、备注口吻与禁词扫描证据。
"""
import json
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

import ppt_maker
from pptx import Presentation

# 市场分析腔禁词（课件型/讲话型内容里不该出现）
BANNED = ["市场", "规模", "格局", "趋势分析"]
# 教师口吻信号词
TEACHER_TONE = ["我们先看", "大家", "同学", "这道题", "注意", "来看", "记住"]

CASES = {
    "1": dict(topic="一元二次方程解法课件（给初三学生上课用）", pages=6,
              style="课堂讲解"),
    "2": dict(topic="开学季（站在班主任的视角，告诉同学们要开学了）", pages=3,
              style="课堂讲解"),
    "3": dict(topic="全球新能源汽车产业格局", pages=3, style="工作汇报"),
}


def extract_pptx(path: str):
    """用 python-pptx 提取每页正文与备注。"""
    prs = Presentation(path)
    pages = []
    for idx, slide in enumerate(prs.slides):
        body = "\n".join(sh.text_frame.text for sh in slide.shapes
                         if sh.has_text_frame)
        note = slide.notes_slide.notes_text_frame.text \
            if slide.has_notes_slide else ""
        pages.append({"idx": idx, "body": body, "note": note})
    return pages


def main():
    case_id = sys.argv[1]
    cfg = CASES[case_id]
    r = ppt_maker.make_ppt(cfg["topic"], pages=cfg["pages"], style=cfg["style"])
    print("=== make_ppt 返回 ===")
    print(json.dumps(r, ensure_ascii=False, indent=1))
    if not r["ok"]:
        sys.exit(1)

    run = ppt_maker.last_run
    print("\n=== 意图判定 ===")
    print(json.dumps(run["intent"], ensure_ascii=False))
    print("research_chars =", run["research_chars"])

    # 读存档拿大纲结构
    deck_json = Path(r["path"]).with_suffix(".deck.json")
    outline_titles = []
    if deck_json.exists():
        arc = json.loads(deck_json.read_text(encoding="utf-8"))
        outline_titles = [(p["layout"], p["page_title"])
                          for p in arc["outline"]["pages"]]
    print("\n=== 大纲结构（layout | page_title）===")
    for layout, title in outline_titles:
        print(f"  {layout:12s} | {title}")

    pages = extract_pptx(r["path"])
    print("\n=== 物理页数 ===", len(pages))

    print("\n=== 各页备注（前 80 字）===")
    for p in pages:
        print(f"  [页{p['idx']}] {p['note'][:80]}")

    # 口吻与禁词扫描（跳过封面页 0）
    all_note = "\n".join(p["note"] for p in pages)
    all_body = "\n".join(p["body"] for p in pages)
    tones = [w for w in TEACHER_TONE if w in all_note]
    banned_hits = [w for w in BANNED if (w in all_note or w in all_body)]
    print("\n=== 口吻信号命中 ===", tones)
    print("=== 禁词命中 ===", banned_hits)

    # 例题页证据：正文同时含题干特征与解析特征
    print("\n=== 例题页扫描（题干/解析关键词）===")
    for p in pages:
        if re.search(r"例[题0-9一二三四]|题干|已知|求[：:]", p["body"]):
            mark = [w for w in ("解析", "解：", "思路", "步骤", "答案")
                    if w in p["body"]]
            print(f"  [页{p['idx']}] 题干特征命中，解析特征={mark}")
            print("    正文摘录:", p["body"][:300].replace("\n", " / "))

    print("\n=== 文件 ===", r["path"])


if __name__ == "__main__":
    main()
