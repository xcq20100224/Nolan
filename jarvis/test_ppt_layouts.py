# -*- coding: utf-8 -*-
"""
jarvis/test_ppt_layouts.py —— ppt_layouts 版式引擎的验收测试（不走 LLM，纯本地）。

覆盖：
  1. 全 8 种 layout 的完整 deck（chart 三种类型各一页，共 10 内容页 + 1 封面）；
  2. 每页 notes_slide 物理非空（含封面）；
  3. chart 页含 GraphicFrame 原生图表（has_chart），且类型映射正确；
  4. 深色版式（封面/章节/金句）与浅色版式背景色断言；
  5. 缺字段页不抛异常、未知 layout 降级 bullets、chart 数据缺失降级；
  6. toc 超 8 条分两栏（仅冒烟，不断言坐标）；
  7. AI 配图扩展（路 B 契约）：bullets 带图左文右图、cover_image 封面背景图
     + 蒙层 z-order、无效图片路径静默降级。

运行：python jarvis/test_ppt_layouts.py（全部断言通过即全绿）
"""
from __future__ import annotations

import io
import os
import sys
import tempfile
from pathlib import Path

# 保证能 import 同目录模块（从仓库根或 jarvis/ 目录运行都行）
sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation          # noqa: E402
from pptx.util import Inches           # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE    # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.oxml.ns import qn                 # noqa: E402

from ppt_layouts import render_deck    # noqa: E402


# ---------------------------------------------------------------- 构造测试 deck

def _full_deck() -> dict:
    """覆盖全部 8 种 layout；chart 三种类型各一页，共 10 内容页。"""
    return {
        "title": "Nolan 项目季度复盘",
        "subtitle": "从语音助手到个人智能体",
        "pages": [
            {"layout": "toc",
             # 10 条 > 8，触发两栏分支
             "entries": ["季度概览", "关键指标", "产品进展", "技术架构",
                         "用户反馈", "竞品动态", "问题与风险", "下季规划",
                         "资源需求", "总结行动"],
             "speaker_note": "目录页：带听众快速过一遍十个部分的结构，说明递进关系。"},
            {"layout": "section",
             "page_title": "关键指标",
             "core_point": "三个数字定义这个季度：留存、时延、任务完成率。",
             "speaker_note": "章节过渡：宣告进入指标部分，先抛问题再亮数字。"},
            {"layout": "bullets",
             "page_title": "产品进展",
             "bullets": ["语音链路完成全双工改造，打断响应时延从 800ms 降到 230ms",
                         "记忆系统升级 v2，跨会话事实召回率由 61% 提升到 84%",
                         "新增 PPT 生成管线，大纲加逐页精写两阶段产出带演讲稿的成稿",
                         "微信发送走 UI 自动化兜底，群聊与私聊成功率稳定在 98% 以上"],
             "speaker_note": "逐条讲：先时延、再记忆、再 PPT 管线、最后发送稳定性。"},
            {"layout": "two_column",
             "page_title": "架构对比",
             "left": {"heading": "旧架构",
                      "points": ["单轮问答式交互，无法打断", "记忆基于文件检索，召回不稳",
                                 "技能硬编码，扩展靠改代码"]},
             "right": {"heading": "新架构",
                       "points": ["全双工语音流，支持随时打断", "向量加图双路记忆，可解释召回",
                                  "技能即插即用，运行时热加载"]},
             "speaker_note": "左右对照讲，强调每一项旧痛点对应的新方案。"},
            {"layout": "big_number",
             "page_title": "季度关键数字",
             "stats": [{"number": "84%", "caption": "跨会话记忆召回率"},
                       {"number": "230ms", "caption": "全双工打断响应时延"},
                       {"number": "98.6%", "caption": "消息发送成功率"}],
             "speaker_note": "三个数字逐个亮：召回率、时延、成功率，每个配一句解释。"},
            {"layout": "quote",
             "quote": "工具的价值不在于功能多，而在于关键时刻靠得住。",
             "attribution": "团队复盘共识",
             "speaker_note": "金句页：放慢语速读一遍，停顿两秒再展开。"},
            {"layout": "chart",
             "page_title": "周活跃任务量",
             "chart": {"type": "bar", "title": "近四周任务量",
                       "categories": ["第1周", "第2周", "第3周", "第4周"],
                       "series": [{"name": "任务数", "values": [42, 55, 61, 78]}]},
             "bullets": ["第 4 周环比提升 28%，全双工上线拉动明显",
                         "语音类任务占比首次超过文本类"],
             "speaker_note": "柱状图先讲趋势，再点出第 4 周的拐点原因。"},
            {"layout": "chart",
             "page_title": "任务类型分布",
             "chart": {"type": "pie", "title": "类型占比",
                       "categories": ["信息查询", "文件处理", "消息发送", "日程提醒", "其他"],
                       "series": [{"name": "占比", "values": [35, 25, 18, 12, 10]}]},
             "bullets": ["信息查询与文件处理合计六成，是核心场景",
                         "提醒类占比低但留存最高，值得加码"],
             "speaker_note": "饼图讲结构：头部两个场景占六成，尾部长尾不追。"},
            {"layout": "chart",
             "page_title": "响应时延趋势",
             "chart": {"type": "line", "title": "P50 时延（ms）",
                       "categories": ["4月", "5月", "6月"],
                       "series": [{"name": "语音链路", "values": [820, 460, 230]},
                                  {"name": "文本链路", "values": [300, 260, 210]}]},
             "bullets": ["语音链路三个月降 72%，已贴近文本链路",
                         "6 月两条曲线收敛，体验差距基本抹平"],
             "speaker_note": "折线图讲收敛：两条曲线的差距就是体验差距。"},
            {"layout": "closing",
             "page_title": "总结与行动",
             "bullets": ["把全双工能力推广到所有语音入口，7 月底前完成",
                         "记忆系统开放用户可视化编辑，降低纠错成本",
                         "PPT 管线接入更多版式与模板，覆盖周报与课件场景"],
             "speaker_note": "收尾：重申三个行动项的负责人与时间点，致谢退场。"},
        ],
    }


# ---------------------------------------------------------------- 辅助断言

def _new_prs() -> Presentation:
    """新建 16:9 演示文稿（模拟调用方已设好画幅）。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _reopen(prs: Presentation) -> Presentation:
    """存进内存再重新打开，模拟真实落盘后的读取断言。"""
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _bg_hex(slide) -> str:
    return str(slide.background.fill.fore_color.rgb)


def _note_text(slide) -> str:
    return slide.notes_slide.notes_text_frame.text.strip()


def _charts(slide):
    return [sh.chart for sh in slide.shapes if sh.has_chart]


# ---------------------------------------------------------------- 测试图现场生成（PIL）

def _make_test_images() -> dict:
    """用 PIL 现场生成两张暖色纯色块测试图，存临时目录。
    返回 {"img_4_3": 1024x768 路径, "img_16_9": 1920x1080 路径}。"""
    from PIL import Image
    tmp = tempfile.mkdtemp(prefix="nolan_ppt_img_")
    p43 = os.path.join(tmp, "warm_block_4x3.png")
    p169 = os.path.join(tmp, "warm_block_16x9.png")
    Image.new("RGB", (1024, 768), (192, 96, 74)).save(p43)     # 赭红暖色块
    Image.new("RGB", (1920, 1080), (168, 123, 95)).save(p169)  # 陶棕暖色块
    return {"img_4_3": p43, "img_16_9": p169}


def _pictures(slide):
    """页面上所有 Picture 形状（按 z-order 顺序）。"""
    return [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]


def _fill_hex(shape) -> str:
    """取形状实心填充色；非实心/取不到返回空串。"""
    try:
        return str(shape.fill.fore_color.rgb)
    except Exception:
        return ""


# ---------------------------------------------------------------- 测试项

def test_full_deck():
    """完整 deck：页数、备注、图表、背景色。"""
    deck = _full_deck()
    prs = _new_prs()
    render_deck(prs, deck, style="工作汇报")
    prs = _reopen(prs)

    # 1) 页数 = 10 内容页 + 1 封面
    assert len(prs.slides) == len(deck["pages"]) + 1 == 11, \
        f"页数不符：{len(prs.slides)}"

    # 2) 每页（含封面）备注物理非空
    for i, slide in enumerate(prs.slides, start=1):
        assert _note_text(slide), f"第 {i} 页备注为空"

    # 3) 深色版式背景：封面(1)、章节(3)、金句(7) 为深棕
    for idx in (0, 2, 6):
        assert _bg_hex(prs.slides[idx]) == "3B322C", \
            f"第 {idx + 1} 页应为深色底，实际 {_bg_hex(prs.slides[idx])}"
    # 浅色内容页（要点页 4）为米白
    assert _bg_hex(prs.slides[3]) == "FAF7F2"

    # 4) 三个 chart 页（物理第 8/9/10 页）各含一个原生图表，且类型映射正确
    expected = {
        7: XL_CHART_TYPE.COLUMN_CLUSTERED,   # bar -> 纵向簇状柱
        8: XL_CHART_TYPE.PIE,                # pie -> 饼图
        9: XL_CHART_TYPE.LINE,               # line -> 折线
    }
    for idx, ctype in expected.items():
        charts = _charts(prs.slides[idx])
        assert len(charts) == 1, f"第 {idx + 1} 页应有 1 个原生图表，实际 {len(charts)}"
        assert charts[0].chart_type == ctype, \
            f"第 {idx + 1} 页图表类型应为 {ctype}，实际 {charts[0].chart_type}"

    # 5) 折线图两条系列都在（数据真的进了图表）
    line_chart = _charts(prs.slides[9])[0]
    assert len(list(line_chart.series)) == 2, "折线图应有 2 条系列"

    print("[PASS] test_full_deck: 11 页 / 备注全非空 / 3 个原生图表类型正确 / 深浅背景正确")


def test_fault_tolerance():
    """容错：缺字段、未知 layout、chart 数据缺失，均不抛异常且备注非空。"""
    deck = {
        "title": "",          # 缺主标题 -> 默认「未命名汇报」
        "subtitle": None,
        "pages": [
            {},                                            # 全缺 -> bullets 兜底
            {"layout": "bullets"},                         # 缺 page_title / bullets / note
            {"layout": "mystery_layout", "page_title": "未知版式"},  # 未知 -> 降级 bullets
            {"layout": "chart", "page_title": "无数据的图表页"},      # chart 数据缺 -> 降级
            {"layout": "big_number", "stats": "不是列表"},            # 非法 stats
            {"layout": "two_column"},                      # 缺 left/right
            "不是字典的页",                                  # 非 dict 页
        ],
    }
    prs = _new_prs()
    render_deck(prs, deck)   # 不抛异常即为第一断言
    prs = _reopen(prs)

    assert len(prs.slides) == len(deck["pages"]) + 1 == 8, "降级渲染后页数不能缺斤短两"
    for i, slide in enumerate(prs.slides, start=1):
        assert _note_text(slide), f"容错用例第 {i} 页备注为空"

    # chart 数据缺失的那页应被降级为要点页：不含图表，但有文本内容
    chartless = prs.slides[4]
    assert not _charts(chartless), "chart 数据缺失页不应残留图表"
    texts = [sh.text_frame.text for sh in chartless.shapes if sh.has_text_frame]
    assert any(t.strip() for t in texts), "降级页应有文本内容"

    print("[PASS] test_fault_tolerance: 7 种异常输入全部安全降级，无异常抛出，备注全非空")


def test_unknown_layout_falls_back_to_bullets():
    """未知 layout 明确按 bullets 版式渲染（标题 + 要点符号）。"""
    deck = {"title": "降级测试", "subtitle": "",
            "pages": [{"layout": "weird_type",
                       "page_title": "这一页用了未知版式",
                       "bullets": ["第一条要点", "第二条要点"],
                       "speaker_note": "未知版式降级验证。"}]}
    prs = _new_prs()
    render_deck(prs, deck)
    prs = _reopen(prs)
    slide = prs.slides[1]
    assert _bg_hex(slide) == "FAF7F2", "未知 layout 降级后应走浅色 bullets 版式"
    all_text = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "这一页用了未知版式" in all_text, "降级页应保留页标题"
    assert "第一条要点" in all_text and "▪" in all_text, "降级页应按要点版式渲染符号与内容"

    print("[PASS] test_unknown_layout_falls_back_to_bullets: 未知版式正确降级为要点页")


def test_bullets_with_image():
    """bullets 带图页：左文右图 —— 有 Picture、浅金衬底在图下、文字仍在；
    two_column 携带 image 字段应被忽略（保持纯文字）。"""
    imgs = _make_test_images()
    deck = {
        "title": "配图版式测试", "subtitle": "",
        "pages": [
            {"layout": "bullets",
             "page_title": "左文右图要点页",
             "bullets": ["全双工语音链路时延降到 230ms，体验接近真人对话",
                         "记忆系统跨会话召回率提升到 84%",
                         "PPT 管线新增 AI 配图版式，要点页自动左文右图"],
             "image": imgs["img_4_3"],
             "speaker_note": "带图要点页：先讲左侧三条要点，再点右侧配图。"},
            {"layout": "two_column",
             "page_title": "带 image 字段的两栏页（应忽略）",
             "left": {"heading": "左栏", "points": ["要点一"]},
             "right": {"heading": "右栏", "points": ["要点二"]},
             "image": imgs["img_4_3"],   # two_column 不响应 image
             "speaker_note": "两栏页应保持纯文字。"},
        ],
    }
    prs = _new_prs()
    render_deck(prs, deck)
    prs = _reopen(prs)

    # 1) bullets 带图页：恰好 1 张配图
    slide = prs.slides[1]
    pics = _pictures(slide)
    assert len(pics) == 1, f"带图要点页应有 1 张配图，实际 {len(pics)}"

    # 2) 文字仍在：标题与要点都渲染出来了
    all_text = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "左文右图要点页" in all_text, "带图页应保留页标题"
    assert "全双工语音链路时延降到 230ms" in all_text, "带图页应保留要点文字"

    # 3) 浅金衬底色块存在，且 z-order 在图片之下
    shapes = list(slide.shapes)
    pic_idx = shapes.index(pics[0])
    gold_idx = next((i for i, sh in enumerate(shapes)
                     if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                     and _fill_hex(sh) == "D9C9A3"), None)
    assert gold_idx is not None, "配图外应有浅金衬底色块"
    assert gold_idx < pic_idx, "衬底色块应垫在图片之下（z-order）"

    # 4) 图高不超过正文区（5.1 英寸）
    assert pics[0].height <= Inches(5.11), f"图高超限：{pics[0].height}"

    # 5) two_column 携带 image 字段被忽略：无 Picture
    tc_slide = prs.slides[2]
    assert not _pictures(tc_slide), "two_column 应忽略 image 字段，保持纯文字"
    tc_text = "\n".join(sh.text_frame.text for sh in tc_slide.shapes if sh.has_text_frame)
    assert "左栏" in tc_text and "右栏" in tc_text

    print("[PASS] test_bullets_with_image: 左文右图（配图+浅金衬底 z-order 正确）/ "
          "two_column 忽略 image")


def test_cover_image():
    """封面背景图：Picture + 深棕蒙层（30% 不透明）+ 标题文字框，z-order 依次递增。"""
    imgs = _make_test_images()
    deck = {
        "title": "封面背景图测试",
        "subtitle": "蒙层之上的副标题",
        "cover_image": imgs["img_16_9"],
        "pages": [{"layout": "bullets", "page_title": "占位页",
                   "bullets": ["占位要点"], "speaker_note": "占位备注。"}],
    }
    prs = _new_prs()
    render_deck(prs, deck)
    prs = _reopen(prs)

    slide = prs.slides[0]
    shapes = list(slide.shapes)

    # 1) 封面有背景图，且在形状列表最前（先图）
    pics = [(i, sh) for i, sh in enumerate(shapes)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics, "封面应有背景图 Picture"
    pic_idx = pics[0][0]

    # 2) 蒙层：图片之后第一个深棕 #3B322C 矩形（后蒙层）
    overlay_idx = next((i for i, sh in enumerate(shapes)
                        if i > pic_idx
                        and sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                        and _fill_hex(sh) == "3B322C"), None)
    assert overlay_idx is not None, "背景图之上应有深棕蒙层矩形"

    # 3) 蒙层透明度：OOXML alpha = 30000（30% 不透明 / 70% 透明）
    srgb = shapes[overlay_idx]._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgb.find(qn("a:alpha")) if srgb is not None else None
    assert alpha is not None and alpha.get("val") == "30000", \
        f"蒙层 alpha 应为 30000，实际 {alpha.get('val') if alpha is not None else None}"

    # 4) 标题文字框在蒙层之上（后文字）
    title_idx = next((i for i, sh in enumerate(shapes)
                      if sh.has_text_frame and "封面背景图测试" in sh.text_frame.text), None)
    assert title_idx is not None, "封面应有标题文字框"
    assert pic_idx < overlay_idx < title_idx, \
        f"z-order 应为 图({pic_idx}) < 蒙层({overlay_idx}) < 文字({title_idx})"

    # 5) 副标题也在（文字层完整）
    all_text = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "蒙层之上的副标题" in all_text

    print("[PASS] test_cover_image: 封面 图->蒙层(alpha 30000)->文字 z-order 正确")


def test_missing_image_degrades():
    """image / cover_image 指向不存在路径：静默降级无图版式，绝不抛异常。"""
    bogus = os.path.join(tempfile.gettempdir(), "nolan_绝不存在的配图_9f3k2.png")
    assert not os.path.isfile(bogus), "前置条件：测试用 bogus 路径不能真的存在"
    deck = {
        "title": "无效图降级测试", "subtitle": "",
        "cover_image": bogus,   # 封面图不存在 -> 纯色封面
        "pages": [
            {"layout": "bullets", "page_title": "无效配图页",
             "bullets": ["要点甲", "要点乙"],
             "image": bogus,   # 配图不存在 -> 无图 bullets 版式
             "speaker_note": "降级备注。"},
            {"layout": "closing", "page_title": "结尾页",
             "bullets": ["行动项一"], "image": bogus,
             "speaker_note": "结尾备注。"},
        ],
    }
    prs = _new_prs()
    render_deck(prs, deck)   # 不抛异常即第一断言
    prs = _reopen(prs)

    # 封面与内容页都不应出现 Picture
    for i, slide in enumerate(prs.slides, start=1):
        assert not _pictures(slide), f"第 {i} 页不应有图片（无效路径应降级）"

    # 降级后文字内容完整
    s1_text = "\n".join(sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
    assert "无效配图页" in s1_text and "要点甲" in s1_text
    assert _bg_hex(prs.slides[1]) == "FAF7F2", "降级后应为标准米白 bullets 版式"

    print("[PASS] test_missing_image_degrades: 无效图片路径静默降级，无异常，文字完整")


# ---------------------------------------------------------------- 入口

if __name__ == "__main__":
    test_full_deck()
    test_fault_tolerance()
    test_unknown_layout_falls_back_to_bullets()
    test_bullets_with_image()
    test_cover_image()
    test_missing_image_degrades()
    print("\n全部测试通过 ✔")
