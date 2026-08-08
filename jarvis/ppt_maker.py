# -*- coding: utf-8 -*-
"""
Nolan · PPT 引擎（ppt_maker.py）——两阶段精写 + 自动选版式版
一句话生成带演讲稿的真 .pptx 文件：
  - 阶段 1 · 大纲：1 次 LLM 调用产出 总标题 + 副标题 + 每页
    {layout, page_title, core_point, keywords, ...版式草稿字段}；
    标题规则为结论式标题（action title）：page_title 必须是带数字/比较级的断言句，
    全篇按「情境 -> 冲突 -> 分析 -> 行动」叙事弧线组织页面顺序；
    LLM 按内容为每页从八种版式（toc/section/bullets/two_column/big_number/quote/chart/closing）
    中自动选型，并给出 two_column 栏题、big_number 数字草稿、chart 数据草稿等；
  - 阶段 2 · 逐页精写：每页 1 次 LLM 调用（串行），按版式要内容——
    bullets/closing 出 4-6 条要点；two_column 左右各 3-4 点；big_number 校准数字；
    chart 出最终图表数据 + 2-3 条解读；quote 出金句与出处；
    toc/section 不烧调用，大纲字段直接用、演讲稿模板合成；
    每页过对应版式的质量闸，不达标重写（最多 2 次），仍不达标取历次最好；
    某页彻底失败时非常规版式降级为 bullets 兜底页，页数永不缺斤短两；
  - 生图（路 B）：大纲阶段 LLM 同时产出 cover_image_prompt（封面背景）与各页可选
    image_prompt（内容页配图，全篇最多 4 页，只给真正需要画面感的页）；
    排版前统一批量串行调 CogView 生图并下载落盘 files/ppt_assets/，
    成功才把绝对路径写进 page["image"]/deck["cover_image"]；
    单张任何失败只置 None 降级无图版式，配置缺失整条链静默跳过，绝不整单失败；
    总开关 with_images=False 时零 HTTP 调用；
  - 排版：路 A 的 jarvis/ppt_layouts.py 就位时走 render_deck（防御式导入，缺席置 None），
    不可用时回退内置排版（要点式渲染 + 正文字号三档自适应 18/16/14pt），模块任何时刻可用；
  - 文件：python-pptx 造 16:9 幻灯片，每页演讲稿写进 notes_slide（PowerPoint/WPS 备注窗格可见）；
  - 落盘：jarvis/files/ 下，文件名 {主题安全名}_{YYYYMMDD-HHMM}.pptx；
  - 降级：大纲阶段失败 -> ok=False；逐页阶段任何单页失败 -> 兜底页顶上，绝不整单失败。

集成契约（签名冻结，主控按此接线）：
    make_ppt(topic: str, pages: int = 8, style: str = "工作汇报", llm_caller=None,
             with_images: bool = True) -> dict
    成功 {"ok": True, "path": 绝对路径, "file_name": "xxx.pptx", "pages": 物理总页数(含封面), "title": "..."}
    失败 {"ok": False, "error": "人话原因"}
    with_images 为可选关键字参数（路 B 追加）：False 时整条生图链跳过、零 HTTP 调用。

pages 参数语义：向 LLM 请求的内容页数量（钳制 3-20）；返回值里的 pages 是物理总页数（内容页 + 1 封面）。

耗时预期：N 内容页 = 1 + N 次串行 LLM 调用（toc/section 页不耗调用），10 页约 60-120 秒。
"""
from __future__ import annotations

import json
import re
import time
import urllib.request
from pathlib import Path

# ---- 版式渲染引擎（路 A）：防御式导入。ppt_layouts 完工前此模块不存在，
#      _render_deck 置 None，_build_pptx 回退内置排版路径，保证本模块任何时刻可用。
try:
    from ppt_layouts import render_deck as _render_deck
except Exception:
    _render_deck = None

# ---- 联网研究模块（R1）：防御式导入。缺席/失败时研究材料为空串，
#      管线自动降级为纯模型记忆生成，主流程无感。
try:
    from ppt_research import research_topic as _research_topic
except Exception:
    _research_topic = None

# ---- 进度总线（实时进度推送）：防御式导入。缺席时 _emit 为空操作，
#      生成管线行为与接入前完全一致；任何异常静默，绝不影响主流程。
try:
    import progress as _progress
except Exception:
    _progress = None


def _emit(step: str, i=None, n=None) -> None:
    """往进度总线发一条进度（人话短文案 ≤30 字；任何失败静默）。"""
    if _progress is None:
        return
    try:
        _progress.emit(step, i, n)
    except Exception:
        pass

MODULE_DIR = Path(__file__).resolve().parent
FILES_DIR = MODULE_DIR / "files"

MIN_PAGES = 3
MAX_PAGES = 20

# Nolan 视觉规范：低饱和暖色；禁蓝紫渐变、禁高饱和背景
COLOR_TITLE = "3B322C"    # 深棕灰（标题）
COLOR_BODY = "4A4440"     # 正文
COLOR_ACCENT = "C0604A"   # 陶土（强调）
COLOR_BG = "FAF7F2"       # 米白（背景）
FONT_NAME = "微软雅黑"

MAX_BULLETS = 6
MAX_BULLET_LEN = 80       # 防御性截断（要求 LLM 30-60 字，留余量）

# ---- CogView AI 配图（路 B·生图篇）----
MAX_IMAGE_PAGES = 4       # 全篇内容页配图上限：按页序保留前 4 页
IMG_API_TIMEOUT = 30      # 生图 API 单张预算（秒）
IMG_DL_TIMEOUT = 20       # 图片下载单张预算（秒）
# 全套配图统一风格后缀：由代码拼接到每条画面描述之后，保证整套图风格一致
IMAGE_STYLE_SUFFIX = "，极简扁平插画风格，暖色调，赭红与米白配色，无文字无水印"

# ---- 版式词汇表（路 A 交接契约，八种）----
LAYOUTS = {"toc", "section", "bullets", "two_column",
           "big_number", "quote", "chart", "closing"}
NO_DETAIL_LAYOUTS = {"toc", "section"}      # 不需要精写：大纲字段直接用
CHART_TYPES = {"bar", "pie", "line"}        # chart.type 合法值

# ---- 质量闸（逐页验收）----
GATE_MIN_BULLETS = 4                 # bullets 系每页至少 4 条要点
GATE_MIN_CHARS_DEFAULT = 180         # 正文总字数基线（科普分享等）
GATE_MIN_CHARS_STRICT = 220          # 工作汇报 / 课堂讲解 更严
GATE_STRICT_STYLES = {"工作汇报", "课堂讲解"}
MAX_REWRITES = 2                     # 每页最多重写 2 次（即最多 3 次尝试）
# 版式专属闸线
GATE_TC_MIN_POINTS = 3               # two_column 左右各至少 3 点
GATE_CHART_MIN_BULLETS = 2           # chart 解读至少 2 条
GATE_CHART_MIN_CHARS = 60            # chart 解读总字数下限
GATE_QUOTE_MIN_CHARS = 20            # quote 金句至少 20 字

# 最近一次 make_ppt 的运行统计（供验收/调试读取，不属于公开契约）
last_run: dict = {}


# ---------------------------------------------------------------- 通用：LLM 输出解析

def _default_llm_caller():
    from brain import glm_one_shot  # 延迟导入，避免模块加载即拖起大脑
    return glm_one_shot


def _extract_json(text: str):
    """从 LLM 输出中提取 JSON 对象：先去 markdown 围栏直解，失败再做花括号平衡扫描。"""
    if not text:
        return None
    s = text.strip()
    # 拆围栏：```json ... ``` / ``` ... ```
    fence = re.search(r"```(?:json)?\s*(.*?)```", s, re.S)
    candidates = [s]
    if fence:
        candidates.insert(0, fence.group(1).strip())
    for cand in candidates:
        # 直解
        try:
            obj = json.loads(cand)
            if isinstance(obj, dict):
                return obj
        except ValueError:
            pass
        # 花括号平衡扫描（尊重字符串与转义）
        start = cand.find("{")
        if start < 0:
            continue
        depth = 0
        in_str = False
        esc = False
        for i in range(start, len(cand)):
            ch = cand[i]
            if in_str:
                if esc:
                    esc = False
                elif ch == "\\":
                    esc = True
                elif ch == '"':
                    in_str = False
                continue
            if ch == '"':
                in_str = True
            elif ch == "{":
                depth += 1
            elif ch == "}":
                depth -= 1
                if depth == 0:
                    try:
                        obj = json.loads(cand[start:i + 1])
                        if isinstance(obj, dict):
                            return obj
                    except ValueError:
                        pass
                    break
    return None


_REPAIR_PROMPT = """你上一次的输出无法被解析为 JSON。请重新输出：只输出一个合法 JSON 对象本身——不要 markdown 围栏、不要解释、不要任何多余文字。第一个字符必须是 {{，最后一个字符必须是 }}。

原始需求：
{prompt}"""


def _call_json(prompt: str, caller, repair_once: bool = True):
    """调 LLM 并解析 JSON；解析失败可按需修复重试 1 次。返回 dict 或 None。"""
    try:
        raw = caller(prompt)
    except Exception:
        raise                      # 调用异常上交，由调用方决定兜底策略
    data = _extract_json(str(raw)) if raw and str(raw).strip() else None
    if data is None and repair_once:
        try:
            raw2 = caller(_REPAIR_PROMPT.format(prompt=prompt))
        except Exception:
            raw2 = None
        if raw2 and str(raw2).strip():
            data = _extract_json(str(raw2))
    return data


# ---------------------------------------------------------------- 阶段 1：大纲

_OUTLINE_PROMPT = """你是资深演示文稿策划。请为主题「{topic}」设计一份 {pages} 页的{style} PPT 大纲。

只输出一个 JSON 对象，不要输出任何其他文字、解释或 markdown 围栏。格式严格如下：
{{
  "title": "整套 PPT 主标题：一句断言，一句话说清全篇核心论点（≤20字，见下方标题规则）",
  "subtitle": "副标题：补充本套 PPT 的场景与受众（≤30字）",
  "cover_image_prompt": "封面背景图的画面描述（必需，见下方配图规则）",
  "pages": [
    {{
      "layout": "本页版式（八种之一，见下方选版式规则）",
      "page_title": "本页核心断言（≤22字，见下方标题规则）",
      "core_point": "本页核心论点，一句话（30-50字）：必须具体、有信息量，是本页正文要论证的靶心",
      "keywords": ["本页 3-5 个关键词：具体的技术名词、数据点、案例名或机制名"],
      "image_prompt": "本页配图的画面描述（可选，仅按下方配图规则给，不配图的页不要输出此字段）"
    }}
  ]
}}

【标题规则：结论式标题（action title）】
- 每页 "page_title" 必须是一句「断言」而非「话题词」：读完标题就知道本页结论，
  尽量带数字或比较级；标题即论点，正文要点是它的论据；
- 反例 ✗「处理法与发酵控制」→ 正例 ✓「处理法决定 60% 的风味基调」；
  反例 ✗「全球供需格局演变」→ 正例 ✓「中国供给全球 60% 的动力电池」；
- 顶层 "title" 同样是断言句：一句话说清整套 PPT 的核心论点；
  "subtitle" 负责补充场景与受众（给谁看、什么场合）。

【配图规则】
- 顶层 "cover_image_prompt" 必填：封面背景图，画面要宏大或抽象、有氛围感，
  暖色调，绝不包含任何文字、字母、数字，也不含人脸特写；
- 内容页的 "image_prompt"：本套 PPT 应配 3 到 4 页图（上限 4 页），凡是
  bullets 版式且有画面感的页（产品、场景、工艺、概念可视化）默认都应该配图；
  只有数据页、对比页、金句页才明确不配图；
- 画面描述必须具体：写清主体 + 场景 + 色调，中英文均可；只描述画面内容，
  统一风格由系统拼接，不要自己写风格词或「插画」「无文字」这类指令。

【选版式规则】为每页从以下八种版式中选最合适的一种，写进 "layout"，并按需附加字段：
- "bullets"：默认形态，绝大多数内容页用它，无需附加字段；
- "two_column"：内容天然有对比/对立结构（利弊、前后、中外）时选用，
  附加 "left": {{"heading": "左栏标题"}} 和 "right": {{"heading": "右栏标题"}}；
- "big_number"：本页有 1-3 个震撼数字（市场规模、增长率、占比）时选用，
  附加 "stats": [{{"number": "数字（带单位，如 42.7%）", "caption": "口径与来源说明"}}]（1-3 个）；
- "chart"：本页有随类别/时间变化的数据系列时选用，
  附加 "chart": {{"type": "bar 或 pie 或 line", "title": "图表标题",
    "categories": ["类别1", "类别2", "类别3"],
    "series": [{{"name": "系列名", "values": [数值1, 数值2, 数值3]}}]}}，
  values 先给合理估计值（后续精写阶段会校准），数量与 categories 一致；
- "quote"：金句/名言点题页，附加 "quote": "金句原文" 和 "attribution": "出处"，
  全篇至多 1 页，可以没有；
- "section"：章节分隔页，仅当总页数 ≥12 时才允许插入，可以没有；
- "toc"：目录页，当总页数 ≥8 时在 pages 数组首部插入一页，
  附加 "entries": ["各页标题", ...]，不计入 {pages} 页正文页数；
- "closing"：收尾页（行动建议/总结要点），最后一页用它，无需附加字段。
- 版式节奏：同一种版式不得连续出现超过 2 页（bullets 除外但也应穿插变化），
  用 two_column / big_number / chart / quote 打散连续感，全篇读起来有呼吸。

要求：
- pages 数组恰好 {pages} 项正文页（toc 另算），对应 {pages} 页内容；
- 全篇按叙事弧线组织页面顺序：情境（现状/背景）→ 冲突（矛盾/挑战）→
  分析（原因/机制）→ 行动（建议/展望）；第一页内容页必须回答
  「为什么现在该关心这件事」，页与页之间分工明确、互不重复；
- 每页 core_point 必须是可论证的具体论点，禁止「介绍XX」「XX很重要」式空话；
- keywords 要能给后续写作提供弹药（名词、数字、案例），不要空泛形容词；
- 内容贴合「{style}」场景。"""


def _norm_layout(raw) -> str:
    """版式名归一：未知/缺失一律降级 bullets（大纲阶段的防御闸门）。"""
    layout = str(raw or "bullets").strip().lower()
    return layout if layout in LAYOUTS else "bullets"


def _gen_outline(topic: str, pages: int, style: str, caller, research: str = ""):
    """阶段 1：拿大纲。返回 (规范化大纲 dict, None) 或 (None, 人话错误)。
    大纲页除 page_title/core_point/keywords 外，还带 layout 与版式草稿字段。
    research：R1 联网研究材料（空串=无，纯模型记忆）。"""
    prompt = _OUTLINE_PROMPT.format(topic=topic, pages=pages, style=style) \
        + _research_block(research)
    try:
        data = _call_json(prompt, caller, repair_once=True)
    except Exception:
        return None, "大脑调用时出错了，请稍后再试"
    if data is None:
        return None, "大脑返回的内容格式乱了（不是合法 JSON），重试后仍失败，稍后再试"

    title = str(data.get("title") or "").strip()[:30]
    subtitle = str(data.get("subtitle") or "").strip()[:40]
    pages_raw = data.get("pages")
    if not title or not isinstance(pages_raw, list) or not pages_raw:
        return None, "大脑返回的大纲结构不完整（缺标题或缺页面），稍后再试"

    outline_pages = []
    content_count = 0          # toc/section 不占正文页配额
    for i, item in enumerate(pages_raw):
        if not isinstance(item, dict):
            continue
        layout = _norm_layout(item.get("layout"))
        if layout not in NO_DETAIL_LAYOUTS:
            if content_count >= pages:
                continue                       # 超出请求页数的正文页丢弃
            content_count += 1
        page_title = str(item.get("page_title") or "").strip()[:22]   # 断言句上限放宽到 22 字
        if not page_title:
            page_title = "目录" if layout == "toc" else f"第 {content_count or i + 1} 节"
        core_point = str(item.get("core_point") or "").strip()
        if not core_point:
            core_point = f"围绕「{page_title}」展开本页的核心内容与关键结论。"
        keywords_raw = item.get("keywords")
        keywords = []
        if isinstance(keywords_raw, list):
            for k in keywords_raw[:6]:
                k = str(k).strip()
                if k:
                    keywords.append(k[:20])
        entry = {
            "layout": layout,
            "page_title": page_title,
            "core_point": core_point,
            "keywords": keywords,
        }
        # ---- 配图草稿（路 B）：内容页可选 image_prompt；只进大纲，不进契约 page，
        #      生图阶段再按最终版式决定去留（非 bullets 页静默丢弃）
        image_prompt = str(item.get("image_prompt") or "").strip()[:200]
        if image_prompt:
            entry["image_prompt"] = image_prompt
        # ---- 版式草稿字段：精写阶段的弹药/上下文
        if layout == "toc":
            entries = []
            if isinstance(item.get("entries"), list):
                for e in item["entries"][:24]:
                    e = str(e).strip()
                    if e:
                        entries.append(e[:22])
            entry["entries"] = entries
        elif layout == "two_column":
            left = item.get("left") if isinstance(item.get("left"), dict) else {}
            right = item.get("right") if isinstance(item.get("right"), dict) else {}
            entry["left_heading"] = str(left.get("heading") or "左栏").strip()[:12] or "左栏"
            entry["right_heading"] = str(right.get("heading") or "右栏").strip()[:12] or "右栏"
        elif layout == "big_number":
            entry["stats"] = _norm_stats(item.get("stats"))
        elif layout == "chart":
            entry["chart"] = _norm_chart(item.get("chart"))
        elif layout == "quote":
            entry["quote"] = str(item.get("quote") or "").strip()[:150]
            entry["attribution"] = str(item.get("attribution") or "").strip()[:40]
        outline_pages.append(entry)

    if not outline_pages:
        return None, "大脑返回的大纲结构不完整（缺标题或缺页面），稍后再试"

    # ---- 选版式规则的规范化兜底（LLM 不守规矩时的防御）----
    seen_quote = False
    for e in outline_pages:
        if e["layout"] == "quote":
            if seen_quote:                       # 全篇至多 1 页 quote，多余的降级
                e["layout"] = "bullets"
            else:
                seen_quote = True
    if len(outline_pages) < 12:                  # section 仅当总页数 ≥12 才允许
        for e in outline_pages:
            if e["layout"] == "section":
                e["layout"] = "bullets"
    # 收尾页用 closing：最后一页是普通要点页时升级为 closing（精写同路径，零成本）
    if outline_pages and outline_pages[-1]["layout"] == "bullets":
        outline_pages[-1]["layout"] = "closing"

    # 封面背景图画面描述（路 B）：必需但 LLM 可能漏给，漏给则无封面图、不阻塞
    cover_image_prompt = str(data.get("cover_image_prompt") or "").strip()[:300]

    return {"title": title, "subtitle": subtitle, "pages": outline_pages,
            "cover_image_prompt": cover_image_prompt}, None


# ---------------------------------------------------------------- 阶段 2：逐页精写

# 风格写作要求：注入逐页 prompt，让口吻与密度匹配场景
_STYLE_HINTS = {
    "工作汇报": "专业克制、数据导向：多用结论与数字说话，每条要点给出指标、对比或明确行动",
    "课堂讲解": "循序渐进、善用类比：从直觉到机制，类比要贴切，术语首次出现时用一句话解释",
    "科普分享": "通俗但不低幼、保证事实密度：用大白话讲清机制，但每条都要落到具体事实、数字或案例上",
}

_PAGE_PROMPT = """你正在为一份{style} PPT 逐页精写正文，现在写第 {idx}/{total} 页。

【全篇背景】
- 总主题：「{topic}」
- 整套 PPT 标题：「{deck_title}」
- 前一页标题：「{prev_title}」；后一页标题：「{next_title}」（内容不得与它们重复撞车）

【本页任务】
- 本页核心断言（标题即论点，正文是它的论据）：「{page_title}」
- 本页核心论点（必须围绕它展开论证）：{core_point}
- 可用素材关键词：{keywords}

只输出一个 JSON 对象，不要输出任何其他文字、解释或 markdown 围栏。格式严格如下：
{{
  "bullets": ["要点1", "要点2", "要点3", "要点4"],
  "speaker_note": "本页演讲稿"
}}

硬性要求（验收线，达不到会被退回重写）：
- bullets 必须 4 到 6 条；每条 30 到 60 字；
- 正文要点必须与本页断言式标题形成论证关系：标题是论点，
  每条要点都是支撑它的论据（数据、案例、机制或推导），不允许与断言无关的凑数内容；
- 每条要点必须承载具体信息：真实数据、具体案例、机制解释或可操作结论，
  禁止「AI很强大」「前景广阔」这类空话套话；
- 要点之间互补不重复，合起来要能论证本页核心论点；
- 风格要求：{style_hint}；
- speaker_note 为 150-250 字的口语化演讲稿：开场第一句要点出本页断言，
  然后讲这页怎么展开、要点之间怎么串、如何自然过渡到下一页，
  是真人上台能照着讲的稿子，不是要点的复读。"""

# 非常规版式的共享上下文块（与 _PAGE_PROMPT 背景部分一致）
_PAGE_CONTEXT = """你正在为一份{style} PPT 逐页精写正文，现在写第 {idx}/{total} 页。

【全篇背景】
- 总主题：「{topic}」
- 整套 PPT 标题：「{deck_title}」
- 前一页标题：「{prev_title}」；后一页标题：「{next_title}」（内容不得与它们重复撞车）

【本页任务】
- 本页核心断言（标题即论点，正文是它的论据）：「{page_title}」
- 本页核心论点（必须围绕它展开论证）：{core_point}
- 可用素材关键词：{keywords}

"""

_TWO_COLUMN_TASK = """本页版式：双栏对比页（two_column）。左栏主题「{left_heading}」，右栏主题「{right_heading}」。

只输出一个 JSON 对象，不要输出任何其他文字、解释或 markdown 围栏。格式严格如下：
{{
  "left": {{"heading": "{left_heading}", "points": ["要点1", "要点2", "要点3"]}},
  "right": {{"heading": "{right_heading}", "points": ["要点1", "要点2", "要点3"]}},
  "speaker_note": "本页演讲稿"
}}

硬性要求（验收线，达不到会被退回重写）：
- 左右两栏各 3 到 4 条要点；每条 20 到 50 字；
- 两栏要形成鲜明对照：同一维度上左说左的、右说右的，不要各说各话；
- 每条要点必须承载具体信息：真实数据、具体案例、机制解释或可操作结论，禁止空话套话；
- 两栏合计总字数不少于 {min_chars} 字；
- 两栏内容必须共同论证本页断言式标题（标题是论点，两栏对照是它的论据）；
- 风格要求：{style_hint}；
- speaker_note 为 150-250 字的口语化演讲稿：怎么开场、两栏怎么对照着讲、如何过渡到下一页。"""

_BIG_NUMBER_TASK = """本页版式：大数字页（big_number）。大纲草稿数字（可校准改写）：{draft}

只输出一个 JSON 对象，不要输出任何其他文字、解释或 markdown 围栏。格式严格如下：
{{
  "stats": [{{"number": "42.7%", "caption": "2024 年国内市场渗透率，数据来源：行业白皮书"}}],
  "speaker_note": "本页演讲稿"
}}

硬性要求（验收线，达不到会被退回重写）：
- stats 1 到 3 个，每个 number 非空：必须是带单位或百分号的震撼数字，要真实、有出处感；
- caption 20-50 字：说清数字的口径、年份与来源类型；
- 几个数字之间要有分工（规模/增速/占比），不要同义反复；
- 数字必须直接论证本页断言式标题（标题是论点，数字是它的论据）；
- 风格要求：{style_hint}；
- speaker_note 为 150-250 字的口语化演讲稿：数字怎么抛、怎么解释意义、如何过渡。"""

_CHART_TASK = """本页版式：图表页（chart）。大纲草稿（数据需校准）：{draft}

只输出一个 JSON 对象，不要输出任何其他文字、解释或 markdown 围栏。格式严格如下：
{{
  "chart": {{
    "type": "bar",
    "title": "图表标题",
    "categories": ["类别1", "类别2", "类别3"],
    "series": [{{"name": "系列名", "values": [100, 200, 300]}}]
  }},
  "bullets": ["解读1", "解读2"],
  "speaker_note": "本页演讲稿"
}}

硬性要求（验收线，达不到会被退回重写）：
- type 只能是 "bar"、"pie" 或 "line"，沿用大纲草稿的类型，除非它明显不合适；
- values 必须是纯数字（不要带单位/百分号），每个系列的 values 长度与 categories 完全一致；
- 数据要合理可信：量级、趋势符合真实世界常识，在大纲草稿基础上校准；
- bullets 是图解读，2 到 3 条，每条 30 到 60 字：说趋势、说拐点、说含义，禁止复述数字；
- 图表与解读必须直接论证本页断言式标题（标题是论点，数据是它的论据）；
- 风格要求：{style_hint}；
- speaker_note 为 150-250 字的口语化演讲稿：图怎么看、结论是什么、如何过渡。"""

_QUOTE_TASK = """本页版式：金句页（quote）。

只输出一个 JSON 对象，不要输出任何其他文字、解释或 markdown 围栏。格式严格如下：
{{
  "quote": "金句原文",
  "attribution": "—— 作者/出处",
  "speaker_note": "本页演讲稿"
}}

硬性要求（验收线，达不到会被退回重写）：
- quote 至少 20 字：必须与本页核心论点严丝合缝，是真金句而非口号；
- attribution 给出处（人名/作品/场合），存疑就写「佚名」；
- 风格要求：{style_hint}；
- speaker_note 为 150-250 字的口语化演讲稿：为什么在这里引用它、它与全篇论证的关系、如何过渡。"""

# 重写 prompt：带上一次不达标的具体原因，换措辞再要一次
_PAGE_REWRITE_PROMPT = """你上一次为这一页写的内容没有通过验收，需要重写。

【不达标原因】{feedback}

请换一种写法重新撰写：补充更具体的数据、案例或机制细节，把篇幅写到线以上。

{original_prompt}"""


def _gate_min_chars(style: str) -> int:
    """质量闸字数线：工作汇报/课堂讲解更严。"""
    return GATE_MIN_CHARS_STRICT if style in GATE_STRICT_STYLES else GATE_MIN_CHARS_DEFAULT


def _norm_bullets(raw) -> list:
    """归一要点列表：清洗、截断、去空。"""
    bullets = []
    if isinstance(raw, list):
        for b in raw[:MAX_BULLETS]:
            b = str(b).strip()
            if b:
                bullets.append(b[:MAX_BULLET_LEN])
    return bullets


def _norm_column(raw, fallback_heading: str) -> dict:
    """归一 two_column 的单栏：heading ≤12 字，points 最多 4 条。"""
    d = raw if isinstance(raw, dict) else {}
    heading = str(d.get("heading") or "").strip()[:12] or fallback_heading
    points = []
    if isinstance(d.get("points"), list):
        for p in d["points"][:4]:
            p = str(p).strip()
            if p:
                points.append(p[:MAX_BULLET_LEN])
    return {"heading": heading, "points": points}


def _norm_stats(raw) -> list:
    """归一 big_number 数字列表：钳 1-3 个（此处先截到 3，空壳丢弃）。"""
    stats = []
    if isinstance(raw, list):
        for s in raw[:3]:
            if not isinstance(s, dict):
                continue
            number = str(s.get("number") or "").strip()[:20]
            caption = str(s.get("caption") or "").strip()[:MAX_BULLET_LEN]
            if number or caption:
                stats.append({"number": number, "caption": caption})
    return stats


def _coerce_float(v) -> float:
    """chart values 强制转 float：字符串数字、千分位逗号、百分号都能吃，转不动补 0.0。"""
    try:
        return float(str(v).replace(",", "").replace("%", "").strip())
    except (TypeError, ValueError):
        return 0.0


def _norm_chart(raw) -> dict:
    """归一 chart 结构：type 合法化；categories 清洗；series.values 转 float 并截齐补齐到
    categories 等长（长了截断、短了补 0.0）。"""
    d = raw if isinstance(raw, dict) else {}
    ctype = str(d.get("type") or "bar").strip().lower()
    if ctype not in CHART_TYPES:
        ctype = "bar"
    title = str(d.get("title") or "").strip()[:30]
    categories = []
    if isinstance(d.get("categories"), list):
        for c in d["categories"][:12]:
            c = str(c).strip()
            if c:
                categories.append(c[:15])
    series = []
    if isinstance(d.get("series"), list):
        for s in d["series"][:4]:
            if not isinstance(s, dict):
                continue
            name = str(s.get("name") or "").strip()[:20] or f"系列{len(series) + 1}"
            values = []
            if isinstance(s.get("values"), list):
                values = [_coerce_float(v) for v in s["values"][:len(categories)]]
            while len(values) < len(categories):
                values.append(0.0)
            series.append({"name": name, "values": values})
    return {"type": ctype, "title": title, "categories": categories, "series": series}


def _extract_content(layout: str, data: dict, item: dict) -> dict:
    """按版式从 LLM 的精写 JSON 中抽取并归一内容字段。"""
    if layout == "two_column":
        return {
            "left": _norm_column(data.get("left"), item.get("left_heading") or "左栏"),
            "right": _norm_column(data.get("right"), item.get("right_heading") or "右栏"),
        }
    if layout == "big_number":
        return {"stats": _norm_stats(data.get("stats"))}
    if layout == "chart":
        return {"chart": _norm_chart(data.get("chart")),
                "bullets": _norm_bullets(data.get("bullets"))[:3]}
    if layout == "quote":
        return {"quote": str(data.get("quote") or "").strip()[:150],
                "attribution": str(data.get("attribution") or "").strip()[:40]}
    # bullets / closing 及任何漏网版式
    return {"bullets": _norm_bullets(data.get("bullets"))}


def _quality_check(bullets: list, style: str):
    """bullets 系质量闸：要点数 >= 4 且正文总字数达标。返回 (是否通过, 不达标原因)。"""
    total_chars = sum(len(b) for b in bullets)
    min_chars = _gate_min_chars(style)
    if len(bullets) < GATE_MIN_BULLETS:
        return False, f"要点只有 {len(bullets)} 条（要求至少 {GATE_MIN_BULLETS} 条），正文共 {total_chars} 字（要求 ≥{min_chars} 字）"
    if total_chars < min_chars:
        return False, f"正文总字数只有 {total_chars} 字（要求 ≥{min_chars} 字），要点普遍太短、缺少具体信息"
    return True, ""


def _quality_check_page(layout: str, content: dict, style: str):
    """按版式分发质量闸。返回 (是否通过, 不达标原因)。"""
    if layout == "two_column":
        lp = content["left"]["points"]
        rp = content["right"]["points"]
        total = sum(len(p) for p in lp) + sum(len(p) for p in rp)
        min_chars = _gate_min_chars(style)
        if len(lp) < GATE_TC_MIN_POINTS or len(rp) < GATE_TC_MIN_POINTS:
            return False, (f"左栏 {len(lp)} 条、右栏 {len(rp)} 条要点"
                           f"（要求左右各至少 {GATE_TC_MIN_POINTS} 条），两栏合计 {total} 字")
        if total < min_chars:
            return False, f"两栏合计总字数只有 {total} 字（要求 ≥{min_chars} 字），要点普遍太短、缺少具体信息"
        return True, ""
    if layout == "big_number":
        stats = content["stats"]
        if not stats:
            return False, "缺少 stats 数字列表，大数字页必须给出 1-3 个震撼数字"
        empty = [i + 1 for i, s in enumerate(stats) if not s["number"]]
        if empty:
            return False, f"第 {empty} 个数字的 number 为空，大数字页的数字必须非空、有出处感"
        return True, ""
    if layout == "chart":
        chart = content["chart"]
        if not chart["categories"] or not chart["series"]:
            return False, "图表数据不完整（缺 categories 或 series），无法成图"
        b = content["bullets"]
        chars = sum(len(x) for x in b)
        if len(b) < GATE_CHART_MIN_BULLETS:
            return False, f"图解读只有 {len(b)} 条（要求 {GATE_CHART_MIN_BULLETS}-3 条）"
        if chars < GATE_CHART_MIN_CHARS:
            return False, f"图解读总字数只有 {chars} 字（要求 ≥{GATE_CHART_MIN_CHARS} 字），解读太浅"
        return True, ""
    if layout == "quote":
        q = content["quote"]
        if len(q) < GATE_QUOTE_MIN_CHARS:
            return False, f"金句只有 {len(q)} 字（要求 ≥{GATE_QUOTE_MIN_CHARS} 字），太短、撑不起一页"
        return True, ""
    # bullets / closing
    return _quality_check(content["bullets"], style)


def _content_nonempty(layout: str, content: dict) -> bool:
    """该版式内容是否物理非空（决定能否成为候选）。"""
    if layout == "two_column":
        return bool(content["left"]["points"] or content["right"]["points"])
    if layout == "big_number":
        return bool(content["stats"])
    if layout == "chart":
        return bool(content["chart"]["categories"] or content["bullets"])
    if layout == "quote":
        return bool(content["quote"])
    return bool(content["bullets"])


def _content_chars(layout: str, content: dict) -> int:
    """内容总字数：重写仍不达标时取历次最好的比较基准。"""
    if layout == "two_column":
        return (sum(len(p) for p in content["left"]["points"])
                + sum(len(p) for p in content["right"]["points"]))
    if layout == "big_number":
        return sum(len(s["number"]) + len(s["caption"]) for s in content["stats"])
    if layout == "chart":
        return sum(len(b) for b in content["bullets"])
    if layout == "quote":
        return len(content["quote"])
    return sum(len(b) for b in content["bullets"])


def _synthesize_note(page_title: str, bullets: list) -> str:
    """LLM 没给演讲稿时兜底合成，保证 notes 物理非空。"""
    return (f"这一页讲「{page_title}」。开场先用一句话点明本页核心，"
            f"然后围绕这几点逐条展开：" + "；".join(bullets) +
            "。每条要点讲透一个事实再往下走，全部讲完后用一句小结收住，"
            "自然过渡到下一页。建议用时约 2 分钟。")


def _synthesize_note_for(layout: str, item: dict, content: dict) -> str:
    """按版式合成兜底演讲稿：把内容摊平成要点式语句复用模板。"""
    title = item["page_title"]
    if layout == "two_column":
        flat = ([f"{content['left']['heading']}：{p}" for p in content["left"]["points"]]
                + [f"{content['right']['heading']}：{p}" for p in content["right"]["points"]])
        return _synthesize_note(title, flat) if flat else _synthesize_note(title, [item["core_point"]])
    if layout == "big_number":
        flat = [f"{s['number']}（{s['caption']}）" for s in content["stats"]]
        return _synthesize_note(title, flat) if flat else _synthesize_note(title, [item["core_point"]])
    if layout == "chart":
        flat = content["bullets"] or [f"图表「{content['chart']['title']}」的数据走势"]
        return _synthesize_note(title, flat)
    if layout == "quote":
        return (f"这一页是金句页。先平稳念出这句话：「{content['quote']}」（{content['attribution']}），"
                f"停顿两秒让它沉下去，再用一两句话点明它与「{item['core_point'][:30]}」的关系，"
                "然后自然过渡到下一页。建议用时约 1 分钟。")
    return _synthesize_note(title, content.get("bullets") or [item["core_point"]])


def _synthesize_struct_note(layout: str, item: dict) -> str:
    """toc/section 的演讲稿：不烧 LLM 调用，模板合成。"""
    if layout == "toc":
        entries = "、".join(item.get("entries") or []) or "后续各章节"
        return (f"这一页是目录。开场后用 30 秒报一遍整体结构：{entries}。"
                "告诉听众每个部分回答什么问题，让大家的预期对齐，然后直接进入第一页。")
    return (f"这一页是章节页「{item['page_title']}」。先用一句话收住上一部分，"
            f"再抛出本章节的核心：{item['core_point']}。"
            "停顿一拍，给听众一个翻篇的信号。建议用时约 30 秒。")


def _fallback_page(item: dict) -> dict:
    """某页 LLM 彻底失败时的兜底内容：用大纲 core_point + keywords 扩成要点，
    保证页数完整、内容贴题，绝不整单失败。"""
    cp = item["core_point"]
    kws = item.get("keywords") or []
    # 按句切 core_point，尽量保留原话
    parts = [p.strip() for p in re.split(r"[。；;！!？?]", cp) if p.strip()]
    bullets = []
    for p in parts[:MAX_BULLETS]:
        bullets.append(p[:MAX_BULLET_LEN])
    # 句数不够时用关键词补弹药
    for kw in kws:
        if len(bullets) >= GATE_MIN_BULLETS:
            break
        bullets.append(f"关于「{kw}」：{cp[:40]}，细节以现场讲解为准。"[:MAX_BULLET_LEN])
    while len(bullets) < 3:
        bullets.append(f"本页核心：{cp[:MAX_BULLET_LEN - 6]}")
    note = (f"这一页讲「{item['page_title']}」。核心论点是：{cp}"
            f"围绕{'、'.join(kws) if kws else '上述要点'}逐条展开说明，"
            "讲清事实与机制后自然过渡到下一页。建议用时约 2 分钟。")
    return {"bullets": bullets, "speaker_note": note}


def _page_prompt(layout: str, style: str, idx: int, total: int, topic: str,
                 deck_title: str, prev_title: str, next_title: str,
                 item: dict, style_hint: str) -> str:
    """按版式组装逐页精写 prompt。"""
    if layout in ("bullets", "closing"):
        prompt = _PAGE_PROMPT.format(
            style=style, idx=idx, total=total, topic=topic, deck_title=deck_title,
            prev_title=prev_title, next_title=next_title,
            page_title=item["page_title"], core_point=item["core_point"],
            keywords="、".join(item["keywords"]) if item["keywords"] else "（无）",
            style_hint=style_hint)
        if layout == "closing":
            prompt += ("\n- 本页是全篇收尾页：要点以总结结论与可执行的行动建议为主，"
                       "每条行动建议要有主语、有抓手，不要空喊口号。")
        return prompt
    ctx = _PAGE_CONTEXT.format(
        style=style, idx=idx, total=total, topic=topic, deck_title=deck_title,
        prev_title=prev_title, next_title=next_title,
        page_title=item["page_title"], core_point=item["core_point"],
        keywords="、".join(item["keywords"]) if item["keywords"] else "（无）")
    if layout == "two_column":
        return ctx + _TWO_COLUMN_TASK.format(
            left_heading=item.get("left_heading") or "左栏",
            right_heading=item.get("right_heading") or "右栏",
            min_chars=_gate_min_chars(style), style_hint=style_hint)
    if layout == "big_number":
        draft = json.dumps(item.get("stats") or [], ensure_ascii=False) or "（无草稿）"
        return ctx + _BIG_NUMBER_TASK.format(draft=draft, style_hint=style_hint)
    if layout == "chart":
        draft = json.dumps(item.get("chart") or {}, ensure_ascii=False) or "（无草稿）"
        return ctx + _CHART_TASK.format(draft=draft, style_hint=style_hint)
    # quote
    return ctx + _QUOTE_TASK.format(style_hint=style_hint)


def _research_block(research: str) -> str:
    """研究材料注入段：非空时拼在 prompt 尾部（引文编号无意义，提示忽略）；
    空串返回 ""——管线降级为纯模型记忆，行为与接入前完全一致。"""
    research = (research or "").strip()
    if not research:
        return ""
    return (
        "\n\n【真实研究材料】以下是联网检索到的与主题相关的最新事实、数据与案例"
        "（附来源与年份；方括号引文编号请忽略）：大纲设计、标题断言、正文数据与"
        "图表数值必须优先采用这些材料，且不得与材料中的事实矛盾：\n" + research + "\n")


def _gen_page(topic: str, deck_title: str, style: str, item: dict,
              idx: int, total: int, prev_title: str, next_title: str,
              caller, research: str = "") -> dict:
    """阶段 2 单页：按版式精写 -> 版式质量闸 -> 不达标重写（最多 2 次）-> 取历次最好 -> 兜底。
    返回 {layout, content, speaker_note, rewrites, fallback}。
    toc/section 不烧调用；彻底失败时非常规版式降级为 bullets 兜底页。"""
    layout = item.get("layout", "bullets")

    # toc/section：大纲字段直接用，演讲稿模板合成，零 LLM 调用
    if layout in NO_DETAIL_LAYOUTS:
        return {"layout": layout, "content": {},
                "speaker_note": _synthesize_struct_note(layout, item),
                "rewrites": 0, "fallback": False}

    style_hint = _STYLE_HINTS.get(style, _STYLE_HINTS["科普分享"])
    prompt = _page_prompt(layout, style, idx, total, topic, deck_title,
                          prev_title, next_title, item, style_hint) \
        + _research_block(research)   # 研究材料拼一次，重写 prompt 引用 original_prompt 自动带上

    candidates = []   # 历次有效候选：(content, note, 是否达标)
    feedback = ""
    for attempt in range(1 + MAX_REWRITES):
        if attempt > 0:
            # 质量闸未过、进入重写：让前端看到「还在打磨」而不是卡死
            _emit(f"第 {idx}/{total} 页没过质量闸，重写第 {attempt} 次",
                  idx, total)
        this_prompt = prompt if attempt == 0 else _PAGE_REWRITE_PROMPT.format(
            feedback=feedback, original_prompt=prompt)
        try:
            data = _call_json(this_prompt, caller, repair_once=False)
        except Exception:
            # 调用异常：有候选就见好就收，没有就直接兜底，不再烧调用
            break
        if data is None:
            feedback = "输出不是合法 JSON，无法解析"
            continue
        content = _extract_content(layout, data, item)
        note = str(data.get("speaker_note") or "").strip()
        ok, reason = _quality_check_page(layout, content, style)
        if _content_nonempty(layout, content):
            if not note:
                note = _synthesize_note_for(layout, item, content)
            candidates.append((content, note, ok))
        if ok:
            return {"layout": layout, "content": content, "speaker_note": note,
                    "rewrites": attempt, "fallback": False}
        feedback = reason if _content_nonempty(layout, content) else ("输出内容结构缺失。" + reason)

    if candidates:
        # 重写仍不达标：取历次最好（内容字数最多者），绝不整单失败
        content, note, _ = max(candidates, key=lambda c: _content_chars(layout, c[0]))
        return {"layout": layout, "content": content, "speaker_note": note,
                "rewrites": MAX_REWRITES, "fallback": False}
    # 彻底失败：非常规版式降级为 bullets 兜底页，页数与内容完整
    fb = _fallback_page(item)
    return {"layout": "bullets", "content": {"bullets": fb["bullets"]},
            "speaker_note": fb["speaker_note"], "rewrites": 0, "fallback": True}


# ---------------------------------------------------------------- 归一化：路 A 契约结构

def _normalize_page(item: dict, page: dict) -> dict:
    """把 (大纲项, 精写内容) 归一成路 A 契约的 page 结构；每页保证 speaker_note 非空。"""
    layout = page["layout"]
    content = page["content"]
    title = item["page_title"]

    if layout == "toc":
        entries = [str(e).strip()[:20] for e in (item.get("entries") or []) if str(e).strip()]
        pg = {"layout": "toc", "entries": entries}
    elif layout == "section":
        pg = {"layout": "section", "page_title": title, "core_point": item["core_point"]}
    elif layout == "two_column":
        pg = {"layout": "two_column", "page_title": title,
              "left": content["left"], "right": content["right"]}
    elif layout == "big_number":
        stats = content["stats"][:3]
        if not stats:                        # 防御：空数字页降级 bullets
            fb = _fallback_page(item)
            pg = {"layout": "bullets", "page_title": title, "bullets": fb["bullets"]}
            content = {"bullets": fb["bullets"]}
        else:
            pg = {"layout": "big_number", "page_title": title, "stats": stats}
    elif layout == "quote":
        if not content["quote"]:             # 防御：空金句页降级 bullets
            fb = _fallback_page(item)
            pg = {"layout": "bullets", "page_title": title, "bullets": fb["bullets"]}
            content = {"bullets": fb["bullets"]}
        else:
            pg = {"layout": "quote", "quote": content["quote"],
                  "attribution": content["attribution"] or "佚名"}
    elif layout == "chart":
        chart = content["chart"]
        if not chart["categories"] or not chart["series"]:   # 防御：空图表降级 bullets
            fb = _fallback_page(item)
            pg = {"layout": "bullets", "page_title": title, "bullets": fb["bullets"]}
            content = {"bullets": fb["bullets"]}
        else:
            pg = {"layout": "chart", "page_title": title,
                  "chart": chart, "bullets": content["bullets"][:3]}
    elif layout == "closing":
        pg = {"layout": "closing", "page_title": title, "bullets": content["bullets"]}
    else:                                    # bullets 及未知降级
        pg = {"layout": "bullets", "page_title": title,
              "bullets": content.get("bullets") or []}

    note = (page.get("speaker_note") or "").strip()
    if not note:
        note = _synthesize_note_for(pg["layout"], item, content)
    pg["speaker_note"] = note
    return pg


def _normalize_deck(outline: dict, contents: list, style: str) -> dict:
    """归一化整副 deck：路 A render_deck 的入参结构。
    {title, subtitle, pages:[...契约 page...]}"""
    pages = [_normalize_page(item, page)
             for item, page in zip(outline["pages"], contents)]
    # toc 缺 entries 时用全篇标题兜底
    titles = [p["page_title"] for p in pages if p.get("page_title")]
    for p in pages:
        if p["layout"] == "toc" and not p["entries"]:
            p["entries"] = titles
    subtitle = outline.get("subtitle") or f"{style} · {time.strftime('%Y年%m月%d日')}"
    return {"title": outline["title"], "subtitle": subtitle, "pages": pages}


# ---------------------------------------------------------------- 生图阶段（路 B：CogView 配图）

def _load_image_config():
    """读 llm_config.json 的 base_url/api_key（与大脑同一把 key）。
    配置缺失/读失败返回 None —— 整条生图链静默跳过。"""
    try:
        cfg = json.loads((MODULE_DIR / "llm_config.json").read_text(encoding="utf-8"))
        base = str(cfg.get("base_url") or "").strip().rstrip("/")
        key = str(cfg.get("api_key") or "").strip()
        if base and key:
            return base, key
    except Exception:
        pass
    return None


def _cogview_generate(prompt: str, base: str, key: str, size: str = "1024x1024"):
    """调 CogView 生图 API，返回图片 url；任何异常/结构缺失返回 None。
    size：封面用 1344x768（近 16:9，铺屏不形变），内容页默认 1024x1024。"""
    body = json.dumps({
        "model": "cogview-3",
        "prompt": prompt + IMAGE_STYLE_SUFFIX,   # 统一风格后缀在此拼接
        "size": size,
    }).encode("utf-8")
    req = urllib.request.Request(
        base + "/images/generations", data=body,
        headers={"Content-Type": "application/json",
                 "Authorization": "Bearer " + key})
    try:
        with urllib.request.urlopen(req, timeout=IMG_API_TIMEOUT) as resp:
            payload = json.loads(resp.read().decode("utf-8"))
        data = payload.get("data")
        if isinstance(data, list) and data and isinstance(data[0], dict):
            url = str(data[0].get("url") or "").strip()
            if url:
                return url
    except Exception:
        pass
    return None


def _download_image(url: str, assets_dir: Path, seq: int):
    """下载图片落盘，按 content-type 定后缀（默认 .png）；成功返回绝对路径，失败 None。"""
    try:
        req = urllib.request.Request(
            url, headers={"User-Agent": "Nolan-PPT/1.0"})
        with urllib.request.urlopen(req, timeout=IMG_DL_TIMEOUT) as resp:
            ctype = ""
            try:
                ctype = str(resp.headers.get("Content-Type") or "")
            except Exception:
                ctype = ""
            blob = resp.read()
        if not blob:
            return None
        ctype = ctype.split(";")[0].strip().lower()
        ext = ".jpg" if ("jpeg" in ctype or "jpg" in ctype) else ".png"
        assets_dir.mkdir(parents=True, exist_ok=True)
        path = assets_dir / f"pptimg_{time.strftime('%Y%m%d-%H%M%S')}_{seq}{ext}"
        path.write_bytes(blob)
        return str(path.resolve())
    except Exception:
        return None


def _gen_one_image(prompt: str, base: str, key: str, assets_dir: Path, seq: int,
                   size: str = "1024x1024"):
    """单张全流程：API 生图 -> 下载落盘。任何一层失败返回 None（调用方置字段 None）。"""
    url = _cogview_generate(prompt, base, key, size)
    if not url:
        return None
    return _download_image(url, assets_dir, seq)


def _attach_images(deck: dict, outline: dict, enabled: bool = True) -> int:
    """生图阶段总入口（排版前统一批量生成，串行）。
    先把契约键初始化成 None（路 A 可依赖键存在），enabled=False 或配置缺失直接返回。
    内容页配图候选：最终版式为 bullets 且大纲给了 image_prompt 的页，按页序钳前 4。
    返回成功落盘的图片张数。绝不抛异常。"""
    deck["cover_image"] = None
    for pg in deck["pages"]:
        if pg["layout"] == "bullets":
            pg["image"] = None
    if not enabled:
        return 0
    try:
        cfg = _load_image_config()
        if cfg is None:
            return 0
        base, key = cfg
        assets_dir = FILES_DIR / "ppt_assets"
        made = 0
        seq = 0
        # 封面背景图
        cover_prompt = str(outline.get("cover_image_prompt") or "").strip()
        # 内容页配图候选先算出来：进度总线需要「第 j/M 张」的总数 M
        cands = []
        for item, pg in zip(outline["pages"], deck["pages"]):
            prompt = str(item.get("image_prompt") or "").strip()
            if prompt and pg["layout"] == "bullets":
                cands.append((pg, prompt))
        todo = cands[:MAX_IMAGE_PAGES]
        total_img = (1 if cover_prompt else 0) + len(todo)
        j = 0
        if cover_prompt:
            seq += 1
            j += 1
            _emit(f"正在生成配图 {j}/{total_img}：封面", j, total_img)
            path = _gen_one_image(cover_prompt, base, key, assets_dir, seq,
                                  size="1344x768")  # 近 16:9，封面铺屏不形变
            if path:
                deck["cover_image"] = path
                made += 1
        # 内容页配图：zip 大纲与契约页（同序），非 bullets 页的 image_prompt 静默丢弃
        for pg, prompt in todo:
            seq += 1
            j += 1
            _emit(f"正在生成配图 {j}/{total_img}", j, total_img)
            path = _gen_one_image(prompt, base, key, assets_dir, seq)
            if path:
                pg["image"] = path
                made += 1
        return made
    except Exception:
        return 0


def _stat_counts(pg: dict):
    """运行统计用：从契约 page 折算 (要点条数, 内容字数)。"""
    layout = pg["layout"]
    if layout in ("bullets", "closing"):
        return len(pg["bullets"]), sum(len(b) for b in pg["bullets"])
    if layout == "two_column":
        pts = pg["left"]["points"] + pg["right"]["points"]
        return len(pts), sum(len(p) for p in pts)
    if layout == "big_number":
        return len(pg["stats"]), sum(len(s["number"]) + len(s["caption"]) for s in pg["stats"])
    if layout == "chart":
        return len(pg["bullets"]), sum(len(b) for b in pg["bullets"])
    if layout == "quote":
        return 1, len(pg["quote"])
    if layout == "section":
        return 0, len(pg["core_point"])
    return 0, sum(len(e) for e in pg.get("entries") or [])   # toc


# ---------------------------------------------------------------- PPTX 造文件

def _rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str)


def _set_run_font(run, size_pt, color_hex, bold=False):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    f = run.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = _rgb(color_hex)
    f.name = FONT_NAME
    # 中文字形需显式设置 East Asian typeface
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_NAME)


def _set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(COLOR_BG)


def _add_textbox(slide, left, top, width, height):
    from pptx.util import Inches
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _add_accent_bar(slide, left, top, width, height):
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(COLOR_ACCENT)
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


# 正文区物理参数：宽 11.2 英寸、高 5.2 英寸；16:9 下每行约 40 个全角字（18pt）
_BODY_CHARS_PER_LINE = 40


def _fit_body_font(bullets: list):
    """行数自适应：估算排版行数与总字数，超阈值就缩字号（18 -> 16 -> 14 三档），
    连带压缩段后距，保证高密度要点不溢出文本框。"""
    total = sum(len(b) for b in bullets)
    lines = sum(max(1, (len(b) + 3 + _BODY_CHARS_PER_LINE - 1) // _BODY_CHARS_PER_LINE)
                for b in bullets)   # +3 是「•  」前缀
    if total <= 240 and lines <= 7:
        return 18, 0.18
    if total <= 320 and lines <= 9:
        return 16, 0.14
    return 14, 0.10


def _flatten_page(pg: dict) -> list:
    """内置兜底排版用：把任意版式的契约 page 摊平成要点行（标题行用【】标出）。"""
    layout = pg["layout"]
    if layout in ("bullets", "closing"):
        return list(pg["bullets"])
    if layout == "two_column":
        return ([f"【{pg['left']['heading']}】"] + list(pg["left"]["points"])
                + [f"【{pg['right']['heading']}】"] + list(pg["right"]["points"]))
    if layout == "big_number":
        return [f"{s['number']}  ——  {s['caption']}" for s in pg["stats"]]
    if layout == "chart":
        c = pg["chart"]
        lines = [f"【图表 · {c['title'] or pg['page_title']}（{c['type']}）】",
                 " / ".join(c["categories"])]
        return lines + list(pg["bullets"])
    if layout == "quote":
        return [f"「{pg['quote']}」", pg["attribution"]]
    if layout == "section":
        return [pg["core_point"]]
    if layout == "toc":
        return [f"{i + 1}. {e}" for i, e in enumerate(pg["entries"])]
    return []


def _legacy_heading(pg: dict) -> str:
    """内置兜底排版的页标题。"""
    return pg.get("page_title") or ("目录" if pg["layout"] == "toc" else "")


def _build_pptx_legacy(deck: dict, style: str, out_path: Path) -> None:
    """内置兜底排版：要点式渲染（ppt_layouts 缺席时的保底线）。"""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN

    pages = deck["pages"]
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- 封面
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_accent_bar(s, 5.667, 2.55, 2.0, 0.07)
    tf = _add_textbox(s, 1.0, 2.8, 11.333, 1.6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run_font(p.add_run(), 40, COLOR_TITLE, bold=True)
    p.runs[0].text = deck["title"]
    tf2 = _add_textbox(s, 1.0, 4.4, 11.333, 0.8)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = deck["subtitle"]
    _set_run_font(r2, 16, COLOR_BODY)
    s.notes_slide.notes_text_frame.text = (
        f"开场白：各位好，今天汇报的主题是「{deck['title']}」。"
        f"先用一句话点明这次分享的核心价值，再交代整体结构（共 {len(pages)} 个部分），"
        "语速放慢，与听众做一次眼神交流。建议用时约 30 秒。")

    # ---- 内容页（任意版式摊平成要点行渲染）
    for idx, pg in enumerate(pages, start=1):
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _add_accent_bar(s, 0.6, 0.62, 0.09, 0.62)
        tf_h = _add_textbox(s, 0.9, 0.5, 11.8, 0.9)
        ph = tf_h.paragraphs[0]
        rh = ph.add_run()
        rh.text = f"{idx:02d}  {_legacy_heading(pg)}"
        _set_run_font(rh, 28, COLOR_TITLE, bold=True)

        bullets = _flatten_page(pg)
        font_pt, gap_in = _fit_body_font(bullets)
        tf_b = _add_textbox(s, 1.1, 1.7, 11.2, 5.2)
        for j, bullet in enumerate(bullets):
            pb = tf_b.paragraphs[0] if j == 0 else tf_b.add_paragraph()
            pb.space_after = Inches(gap_in)
            rb = pb.add_run()
            rb.text = "•  " + bullet
            _set_run_font(rb, font_pt, COLOR_BODY)

        # 演讲稿物理写入备注（PowerPoint/WPS 备注窗格可见）
        s.notes_slide.notes_text_frame.text = pg["speaker_note"]

    prs.save(str(out_path))


def _build_pptx(deck: dict, style: str, out_path: Path) -> None:
    """排版总入口：路 A 的 render_deck 可用时走版式引擎，否则走内置兜底排版。
    落盘/文件名/返回契约不变。"""
    if _render_deck is not None:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        prs.slide_width = Inches(13.333)   # 16:9，由本侧创建
        prs.slide_height = Inches(7.5)
        _render_deck(prs, deck, style)     # 封面由 deck["title"]/["subtitle"] 自动生成
        prs.save(str(out_path))
        return
    _build_pptx_legacy(deck, style, out_path)


# ---------------------------------------------------------------- 落盘

_SAFE_KEEP = re.compile(r"[0-9A-Za-z一-鿿_-]+")


def _safe_name(topic: str) -> str:
    parts = _SAFE_KEEP.findall(topic)
    name = "_".join(p for p in parts if p)[:30].strip("_")
    return name or "未命名"


def _alloc_path(topic: str) -> Path:
    FILES_DIR.mkdir(parents=True, exist_ok=True)
    stamp = time.strftime("%Y%m%d-%H%M")
    base = f"{_safe_name(topic)}_{stamp}"
    path = FILES_DIR / f"{base}.pptx"
    n = 2
    while path.exists():
        path = FILES_DIR / f"{base}-{n}.pptx"
        n += 1
    return path


# ---------------------------------------------------------------- 公开 API

def make_ppt(topic: str, pages: int = 8, style: str = "工作汇报", llm_caller=None,
             with_images: bool = True, with_research: bool = True) -> dict:
    """一句话生成带演讲稿的 .pptx。契约见模块 docstring。
    with_images：CogView AI 配图总开关（路 B），False 时整条生图链跳过、零 HTTP。
    with_research：联网研究总开关（R1），False 时跳过阶段 0、纯模型记忆生成。"""
    global last_run
    topic = (topic or "").strip()
    if not topic:
        return {"ok": False, "error": "没说 PPT 主题，巧妇难为无米之炊"}
    style = (style or "工作汇报").strip() or "工作汇报"
    try:
        pages = int(pages)
    except (TypeError, ValueError):
        pages = 8
    pages = max(MIN_PAGES, min(MAX_PAGES, pages))

    caller = llm_caller or _default_llm_caller()
    stats = {"outline_retries": 0, "page_stats": [], "llm_calls": 0}
    last_run = stats

    # ---- 阶段 0：联网研究（R1，25s 硬预算；失败/关闭返回空串，
    #      大纲与精写自动降级为纯模型记忆，行为与接入前一致）
    research = ""
    if with_research and _research_topic is not None:
        _emit("正在联网查资料…")
        try:
            research = _research_topic(topic) or ""
        except Exception:
            research = ""
        _emit(f"资料就绪（{len(research)} 字）" if research
              else "没查到资料，凭储备来写")
    stats["research_chars"] = len(research)

    # ---- 阶段 1：大纲（失败即整单失败，没有大纲就没有弹药）
    _emit("正在设计大纲…")
    outline, err = _gen_outline(topic, pages, style, caller, research)
    if outline is None:
        return {"ok": False, "error": err}
    _emit(f"大纲好了，共 {len(outline['pages'])} 页，逐页精写")

    # ---- 阶段 2：逐页精写（串行；单页失败只影响单页，兜底页顶上）
    deck_pages = outline["pages"]
    titles = [p["page_title"] for p in deck_pages]
    contents = []
    for i, item in enumerate(deck_pages):
        _emit(f"正在精写第 {i + 1}/{len(deck_pages)} 页："
              f"{item['page_title'][:12]}", i + 1, len(deck_pages))
        prev_t = titles[i - 1] if i > 0 else "（封面）"
        next_t = titles[i + 1] if i + 1 < len(titles) else "（结束页）"
        page = _gen_page(topic, outline["title"], style, item,
                         i + 1, len(deck_pages), prev_t, next_t, caller, research)
        contents.append(page)

    # ---- 归一化成路 A 契约 deck，并登记运行统计
    deck = _normalize_deck(outline, contents, style)
    for i, (item, pg) in enumerate(zip(deck_pages, deck["pages"])):
        n_bullets, n_chars = _stat_counts(pg)
        stats["page_stats"].append({
            "page": i + 1, "title": item["page_title"], "layout": pg["layout"],
            "bullets": n_bullets, "chars": n_chars,
            "rewrites": contents[i]["rewrites"], "fallback": contents[i]["fallback"],
        })

    # ---- 生图阶段（路 B）：排版前统一批量串行生成；任何单张失败只降级无图，
    #      with_images=False 或配置缺失时整条链静默跳过，绝不整单失败
    stats["images"] = _attach_images(deck, outline, enabled=bool(with_images))

    try:
        out_path = _alloc_path(topic)
        _emit("正在排版渲染…")
        _build_pptx(deck, style, out_path)
    except Exception as e:
        return {"ok": False, "error": f"PPT 文件生成失败：{e}"}

    # ---- 存档（R3 对话式编辑的弹药库）：渲染成功后把完整中间态落在 pptx 旁边，
    #      之后的「说改就改」靠它读档局部重生成。存档失败不影响主交付。
    try:
        archive = {
            "topic": topic, "style": style, "pages": pages,
            "research": research, "outline": outline, "deck": deck,
            "pptx_name": out_path.name,
        }
        (FILES_DIR / (out_path.stem + ".deck.json")).write_text(
            json.dumps(archive, ensure_ascii=False, indent=1), encoding="utf-8")
    except Exception:
        pass

    _emit("PPT 做好了，存档完成")
    return {
        "ok": True,
        "path": str(out_path.resolve()),
        "file_name": out_path.name,
        "pages": len(deck["pages"]) + 1,   # 物理总页数 = 内容页 + 封面
        "title": outline["title"],
    }
