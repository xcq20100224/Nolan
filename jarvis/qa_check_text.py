# -*- coding: utf-8 -*-
"""验证 LaTeX 残留与要点截断。"""
import re
import sys
from pptx import Presentation

p = Presentation(sys.argv[1])
for i, s in enumerate(p.slides, 1):
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            if "$" in t or "\\" in t:
                print(f"--- page {i} LaTeX残留 ---")
                print(t[:500])
print("=== 疑似截断要点（>=60字且结尾非正常收尾） ===")
for i, s in enumerate(p.slides, 1):
    for sh in s.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = para.text.strip()
                if len(t) >= 60 and not re.search(r"[。；：！？…）」』\)》\"'\d]$", t):
                    print(f"p{i} [{len(t)}字] 尾部: ...{t[-15:]}")
