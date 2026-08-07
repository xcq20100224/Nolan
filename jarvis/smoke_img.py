# -*- coding: utf-8 -*-
"""真机冒烟：配图版 PPT 端到端生成与解剖（一次性脚本）"""
import time
import ppt_maker

t0 = time.time()
r = ppt_maker.make_ppt("咖啡：从产地到一杯手冲", pages=8, style="科普分享")
print("ok:", r.get("ok"), "| 耗时: %.1fs" % (time.time() - t0))
print("文件:", r.get("file_name"), "| 页数:", r.get("pages"))
if r.get("ok"):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    p = Presentation(r["path"])
    pics = charts = 0
    for i, s in enumerate(p.slides, 1):
        np_ = sum(1 for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
        nc = sum(1 for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.CHART)
        pics += np_
        charts += nc
        texts = [
            sh.text_frame.text.strip()
            for sh in s.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()
        ]
        first = texts[0][:22] if texts else "(无)"
        print(f"  p{i:02d} | 图{np_} 表{nc} | {first}")
    print("总图片:", pics, "| 总图表:", charts)
