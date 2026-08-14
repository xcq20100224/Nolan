# -*- coding: utf-8 -*-
"""
jarvis/test_ppt_diagrams.py —— ppt_diagrams 图示引擎的验收测试（不走 LLM，纯本地）。

覆盖：timeline / process / pyramid 三种图示渲染函数——
  正常数据、缺字段、空列表、超长文本、脏数据，断言：
  - shapes 数量 > 0（页面物理非空）；
  - 任何输入不抛异常（模块级容错契约）；
  - 文字内容真实出现在 shape 文本里；
  - DIAGRAM_RENDERERS 注册表键名与可调用性（主控接线契约）。

运行：cd jarvis && python test_ppt_diagrams.py（全部断言通过即全绿）
"""
import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation            # noqa: E402
from pptx.util import Inches              # noqa: E402

import ppt_diagrams                       # noqa: E402


# ---------------------------------------------------------------- 测试工具

def _new_slide():
    """造一页 16:9 空白页（与 render_deck 的页面创建方式一致）。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _all_text(slide):
    """收集一页上所有 shape 的文本（textbox + 自选形状文本框）。"""
    chunks = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            chunks.append(shp.text_frame.text)
    return "\n".join(chunks)


def _items(n, label_prefix="节点", desc="这是一条具体的事实性说明，含数据与机制"):
    """造 n 个正常节点。"""
    return [{"label": f"{label_prefix}{i + 1}", "desc": f"{desc}（第{i + 1}条）"}
            for i in range(n)]


# ---------------------------------------------------------------- 注册表契约

class TestRendererRegistry(unittest.TestCase):
    """DIAGRAM_RENDERERS：主控接 _LAYOUTS 的交接契约。"""

    def test_registry_keys_and_callables(self):
        self.assertEqual(set(ppt_diagrams.DIAGRAM_RENDERERS),
                         {"timeline", "process", "pyramid"})
        for fn in ppt_diagrams.DIAGRAM_RENDERERS.values():
            self.assertTrue(callable(fn))

    def test_registry_entries_are_public_renderers(self):
        self.assertIs(ppt_diagrams.DIAGRAM_RENDERERS["timeline"],
                      ppt_diagrams.render_timeline)
        self.assertIs(ppt_diagrams.DIAGRAM_RENDERERS["process"],
                      ppt_diagrams.render_process)
        self.assertIs(ppt_diagrams.DIAGRAM_RENDERERS["pyramid"],
                      ppt_diagrams.render_pyramid)


# ---------------------------------------------------------------- timeline 时间轴

class TestTimeline(unittest.TestCase):
    RENDER = staticmethod(ppt_diagrams.render_timeline)
    KEY = "events"

    def test_normal_data_renders_all_text(self):
        prs, slide = _new_slide()
        page = {"layout": "timeline", "page_title": "技术演进时间轴",
                "speaker_note": "备注",
                "events": _items(4)}
        self.RENDER(slide, page, 1, "课堂讲解")
        self.assertGreater(len(slide.shapes), 0)
        text = _all_text(slide)
        self.assertIn("技术演进时间轴", text)          # 页眉标题
        for i in range(4):
            self.assertIn(f"节点{i + 1}", text)        # 每个 label 都在
            self.assertIn(f"第{i + 1}条", text)        # 每个 desc 都在
        self.assertIn("4", text)                       # 序号徽章

    def test_min_two_events(self):
        prs, slide = _new_slide()
        page = {"page_title": "两节点", "events": _items(2, "阶段")}
        self.RENDER(slide, page, 2, "工作汇报")
        self.assertGreater(len(slide.shapes), 0)
        self.assertIn("阶段2", _all_text(slide))

    def test_max_six_events(self):
        prs, slide = _new_slide()
        page = {"page_title": "六节点", "events": _items(8, "年份")}  # 超上限应截到 6
        self.RENDER(slide, page, 3, "科普分享")
        text = _all_text(slide)
        self.assertIn("年份6", text)
        self.assertNotIn("年份7", text)                # 第 7、8 个被截掉

    def test_missing_key_degrades_safely(self):
        prs, slide = _new_slide()
        self.RENDER(slide, {"page_title": "缺字段页"}, 1, "工作汇报")
        self.assertGreater(len(slide.shapes), 0)
        self.assertIn("内容筹备中", _all_text(slide))  # 占位节点，不空页

    def test_empty_list_degrades_safely(self):
        prs, slide = _new_slide()
        self.RENDER(slide, {"page_title": "空列表", "events": []}, 1, "工作汇报")
        self.assertGreater(len(slide.shapes), 0)
        self.assertIn("内容筹备中", _all_text(slide))

    def test_dirty_items_never_raise(self):
        prs, slide = _new_slide()
        page = {"page_title": "脏数据",
                "events": [None, "字符串不是dict", {"label": ""}, {},
                           {"label": "正常", "desc": None}]}
        self.RENDER(slide, page, 1, "工作汇报")        # 不抛异常即通过
        self.assertGreater(len(slide.shapes), 0)
        self.assertIn("正常", _all_text(slide))

    def test_overlong_text_truncated_and_renders(self):
        prs, slide = _new_slide()
        page = {"page_title": "超长文本",
                "events": [{"label": "标" * 60, "desc": "长" * 200},
                           {"label": "乙", "desc": "短"}]}
        self.RENDER(slide, page, 1, "工作汇报")
        text = _all_text(slide)
        self.assertIn("标" * 16, text)                 # label 截断后仍在
        self.assertNotIn("标" * 60, text)              # 完整超长串不会出现

    def test_page_not_dict_never_raise(self):
        prs, slide = _new_slide()
        self.RENDER(slide, None, 1, "工作汇报")        # page 整个是 None 也不抛
        self.assertGreater(len(slide.shapes), 0)


# ---------------------------------------------------------------- process 流程步骤

class TestProcess(unittest.TestCase):
    RENDER = staticmethod(ppt_diagrams.render_process)
    KEY = "steps"

    def test_normal_data_renders_all_text(self):
        prs, slide = _new_slide()
        page = {"layout": "process", "page_title": "四步上手",
                "speaker_note": "备注",
                "steps": _items(4, "步骤", "做什么、做到什么程度算完成")}
        self.RENDER(slide, page, 1, "课堂讲解")
        self.assertGreater(len(slide.shapes), 0)
        text = _all_text(slide)
        self.assertIn("四步上手", text)
        for i in range(4):
            self.assertIn(f"步骤{i + 1}", text)
            self.assertIn(f"第{i + 1}条", text)
        self.assertIn("1 · 步骤1", text)               # 块内序号+label 连排

    def test_max_five_steps(self):
        prs, slide = _new_slide()
        page = {"page_title": "五步法", "steps": _items(7, "环节")}
        self.RENDER(slide, page, 1, "工作汇报")
        text = _all_text(slide)
        self.assertIn("环节5", text)
        self.assertNotIn("环节6", text)                # 超上限截到 5

    def test_missing_and_empty_degrade(self):
        for page in ({"page_title": "缺steps"},
                     {"page_title": "空steps", "steps": []}):
            prs, slide = _new_slide()
            self.RENDER(slide, page, 1, "工作汇报")
            self.assertGreater(len(slide.shapes), 0)
            self.assertIn("内容筹备中", _all_text(slide))

    def test_overlong_and_dirty_never_raise(self):
        prs, slide = _new_slide()
        page = {"page_title": "脏流程",
                "steps": [{"label": "长" * 80, "desc": "述" * 300},
                          42, {"desc": "只有desc没label"}]}
        self.RENDER(slide, page, 1, "工作汇报")
        self.assertGreater(len(slide.shapes), 0)


# ---------------------------------------------------------------- pyramid 金字塔

class TestPyramid(unittest.TestCase):
    RENDER = staticmethod(ppt_diagrams.render_pyramid)
    KEY = "levels"

    def test_normal_data_renders_all_text(self):
        prs, slide = _new_slide()
        page = {"layout": "pyramid", "page_title": "能力三层模型",
                "speaker_note": "备注",
                "levels": _items(3, "层级", "该层包含什么、为什么在这个位置")}
        self.RENDER(slide, page, 1, "课堂讲解")
        self.assertGreater(len(slide.shapes), 0)
        text = _all_text(slide)
        self.assertIn("能力三层模型", text)
        for i in range(3):
            self.assertIn(f"层级{i + 1}", text)
            self.assertIn(f"第{i + 1}条", text)

    def test_max_five_levels(self):
        prs, slide = _new_slide()
        page = {"page_title": "五层", "levels": _items(7, "层")}
        self.RENDER(slide, page, 1, "工作汇报")
        text = _all_text(slide)
        self.assertIn("层5", text)
        self.assertNotIn("层6", text)

    def test_missing_and_empty_degrade(self):
        for page in ({"page_title": "缺levels"},
                     {"page_title": "空levels", "levels": []}):
            prs, slide = _new_slide()
            self.RENDER(slide, page, 1, "工作汇报")
            self.assertGreater(len(slide.shapes), 0)
            self.assertIn("内容筹备中", _all_text(slide))

    def test_overlong_and_dirty_never_raise(self):
        prs, slide = _new_slide()
        page = {"page_title": "脏层级",
                "levels": [{"label": "顶" * 50, "desc": "述" * 200},
                           "oops", {"label": "底层", "desc": ""}]}
        self.RENDER(slide, page, 1, "工作汇报")
        text = _all_text(slide)
        self.assertIn("底层", text)


# ---------------------------------------------------------------- 整页落盘冒烟

class TestSmokeSave(unittest.TestCase):
    """三种图示各渲一页 + 保存成真实 .pptx，验证文件物理可写。"""

    def test_three_diagram_pages_save(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        pages = [
            ("timeline", {"layout": "timeline", "page_title": "演进",
                          "speaker_note": "n", "events": _items(3, "节点")}),
            ("process", {"layout": "process", "page_title": "流程",
                         "speaker_note": "n", "steps": _items(3, "步骤")}),
            ("pyramid", {"layout": "pyramid", "page_title": "分层",
                         "speaker_note": "n", "levels": _items(3, "层")}),
        ]
        for i, (name, page) in enumerate(pages, start=1):
            slide = prs.slides.add_slide(blank)
            ppt_diagrams.DIAGRAM_RENDERERS[name](slide, page, i, "工作汇报")
        with tempfile.TemporaryDirectory() as td:
            out = os.path.join(td, "diagrams_smoke.pptx")
            prs.save(out)
            self.assertTrue(os.path.isfile(out))
            self.assertGreater(os.path.getsize(out), 10000)   # 非空文件


if __name__ == "__main__":
    unittest.main(verbosity=2)
