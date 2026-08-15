# -*- coding: utf-8 -*-
"""
Nolan · PPT 版式引擎（ppt_layouts.py）—— 成熟产品级版式渲染

职责：把「归一化后的 deck 字典」渲染进 python-pptx 的 Presentation。
本模块是独立版式层：不 import ppt_maker，不做任何 LLM 调用，只负责排版与绘制。

公开接口（与内容管线路 B 的交接契约，签名冻结，一个字不改）：
    render_deck(prs, deck: dict, style: str = "工作汇报") -> None

    deck = {
      "title": str, "subtitle": str,
      "pages": [ page, ... ]   # 每个 page 必含 "layout" 与 "speaker_note"
    }
    page 按 layout 分八种：toc / section / bullets / two_column /
                          big_number / quote / chart / closing
    （各 layout 的字段约定见契约 docstring，本文件头部不再重复）。

    AI 配图扩展契约（与生图管线路 B 的交接）：
      - page 可选键 "image"：图片绝对路径，仅 bullets 版式响应（左文右图），
        two_column / closing 等其余版式一律忽略；渲染前 os.path.isfile 检查，
        文件不存在就当无图降级，绝不抛异常；
      - deck 顶层可选键 "cover_image"：封面背景图绝对路径，全屏铺图 + 深色蒙层。

约定：
  - 封面由 render_deck 自动用 deck["title"]/["subtitle"] 生成，pages 里不含封面；
  - 每页（含封面）都把 speaker_note 物理写进 notes_slide（PPT/WPS 备注窗格可见）；
  - 缺字段给安全默认，未知 layout 按 bullets 降级渲染，绝不抛异常；
  - chart 版式用 python-pptx 原生图表（GraphicFrame，PPT 内可编辑），不是贴图。

精修层（视觉细节，不改契约、不改公开签名）—— Gamma 式精品默认主题：
  - 设计 tokens：全部版式统一引用顶部 token 常量区（主色/辅助/强调/纸面/墨色/
    浅灰分隔线），禁止散落硬编码；验收冻结值（深棕底 3B322C、米白纸面 FAF7F2、
    符号轮转三色、节奏缘带浅金 D9C9A3 alpha 12000、封面蒙层 alpha 30000）保持不变；
  - 字体层级：标题/副题/正文/注释四级字号字重体系，注释级更灰更轻
    （页脚 10pt 浅暖灰），封面主标加字距（rPr spc）；
  - 封面：无图时为编辑部式左对齐构图——风格胶囊 + 大号加字距主标 + 细线装饰 +
    底部细信息行；有图时全出血 + 30% 均匀蒙层（验收契约）+ 底部二次渐变蒙层
    （上浅下深）+ 左下角标题块；
  - 要点符号：页内统一实心圆点「●」，颜色按 赭红→暖灰→陶棕 三态轮转
    （形状全局统一而非轮转——形状差异会暗示不存在的语义层级，颜色轮转只贡献
    「细节的密度」，符合克制哲学）；
  - bullets 配图与浅金衬底色块做 8% 圆角裁切（a:prstGeom roundRect），失败静默
    回退直角；封面全屏出血图保持直角——出血版式圆角会露出底色边角，破坏满铺感；
  - 连续 bullets 段内偶数张（无配图时）在正文右缘加 12% 不透明浅金竖缘带，
    奇数张无，形成「素-饰-素」的翻页节奏防疲劳；段被其他版式打断即重新计数；
  - closing 页居中仪式感：短细线 + 谢谢聆听 + 落款，取代大色带收尾。
"""
from __future__ import annotations

import os
import time

# ================================================================ 设计 tokens
# 一套精品默认主题的 token 常量区：全部版式统一引用此处，禁止散落硬编码。
# 色系：NEGA 暖色系·低饱和（无蓝紫渐变、无高饱和背景），深色/浅色版式交替出节奏。
# 注：以下五行十六进制值为验收测试冻结值（test_ppt_layouts.py 硬断言），只准引用、不准改值。
COLOR_DARK = "3B322C"        # token·墨色底：深色版式底 / 浅色版式标题
COLOR_INK = "4A4440"         # token·墨灰：正文
COLOR_ACCENT = "C0604A"      # token·赭红（主强调）：色块、要点符号、大数字、图表主色
COLOR_BG = "FAF7F2"          # token·纸面：浅色版式底 / 深色版式上的浅字
COLOR_WARMGREY = "8A8578"    # token·暖灰（辅助）：次要文字、图表辅色
COLOR_GOLD = "D9C9A3"        # token·浅金（装饰）：细线、衬底、深色底上的点缀
COLOR_TAN = "A87B5F"         # token·陶棕：要点符号轮转第三色（取自图表色板，色系自洽）

# ---- 派生 tokens（精修新增，非冻结）
COLOR_DEEP = "2A231E"        # token·深墨：封面渐变蒙层底部收深色（比墨色底再深一档）
COLOR_CAPTION = "ABA394"     # token·浅暖灰：页脚/落款/注释（比暖灰更轻，四级层级最弱档）
COLOR_HAIRLINE = "E3DACA"    # token·浅灰分隔线：页眉下细线、页脚细线（纸面上的呼吸缝）

# ---- 要点符号（页内形状统一为实心圆点，颜色按要点序号三态轮转；值为验收冻结）
MARKER_GLYPH = "●  "                                        # 实心圆点 + 两空格
MARKER_ROTATION = (COLOR_ACCENT, COLOR_WARMGREY, COLOR_TAN)  # 赭红 → 暖灰 → 陶棕

# ---- 图表系列配色（暖色系轮转：赭红 -> 暖灰 -> 浅金 -> 墨灰 -> 陶棕）
CHART_PALETTE = [COLOR_ACCENT, COLOR_WARMGREY, COLOR_GOLD, COLOR_INK, COLOR_TAN]

# ---- 字体与四级字号阶梯（标题大胆 / 副题 / 正文三档缩 / 注释更灰更轻）
FONT_NAME = "微软雅黑"
SIZE_COVER_TITLE = 44        # 封面主标（40 -> 44，配合字距更有排面）
SIZE_COVER_SUB = 17          # 封面副标（副题级）
SIZE_SECTION_TITLE = 34      # 章节页标题
SIZE_PAGE_TITLE = 26         # 页标题
SIZE_BODY = 18               # 正文基准（按密度三档缩：18 -> 16 -> 14）
SIZE_BODY_MID = 16
SIZE_BODY_MIN = 14
SIZE_NOTE = 12               # 注释 / 图表标签 / 胶囊标签
SIZE_CAPTION = 10            # 页脚 / 落款（注释级再降一档，更灰更轻）
SIZE_BIG_NUMBER = 66         # 大数字页数字（60pt+）
SIZE_SECTION_NO = 150        # 章节页超大半透明序号
SIZE_QUOTE_MARK = 150        # 金句页大号引号装饰
SPACING_COVER_TITLE = 2.0    # 封面主标字距（磅；OOXML rPr spc 单位 1/100 磅）
SPACING_KICKER = 1.5         # 小标签/章节导引字距

# ---- 网格（16:9：13.333 x 7.5 英寸）
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.8                     # 页边距：所有元素左对齐到这一条线，禁止漂移
CONTENT_W = SLIDE_W - 2 * MARGIN  # 内容区宽 11.733
CONTENT_LEFT = MARGIN
HEADER_Y = 0.5                   # 页标题区顶
BODY_Y = 1.7                     # 正文区顶
FOOTER_Y = 7.05                  # 页脚细线

# ---- 深色版式集合（封面单独处理，不在 pages 里）
DARK_LAYOUTS = {"section", "quote"}

# ---- 图表类型映射（python-pptx 原生图表，PPT 里可编辑）
def _chart_type_map():
    from pptx.enum.chart import XL_CHART_TYPE
    return {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,   # 纵向簇状柱
        "hbar": XL_CHART_TYPE.BAR_CLUSTERED,     # 横向簇状条（扩展支持）
        "pie": XL_CHART_TYPE.PIE,                # 饼图
        "line": XL_CHART_TYPE.LINE,              # 折线
    }

# 正文字符宽度估算：16:9 内容区 18pt 下每行约 40 个全角字（沿用 ppt_maker 的经验值）
_BODY_CHARS_PER_LINE = 40

# ---- bullets 图文分栏参数（page["image"] 有效时启用「左文右图」）
IMG_TEXT_W = 6.45                                # 左文区宽：内容区 11.733 的 55%
IMG_GAP = 0.4                                    # 图文间距
IMG_ZONE_W = CONTENT_W - IMG_TEXT_W - IMG_GAP    # 右图区宽 ≈ 4.88（约 42%）
IMG_MAX_H = 5.1                                  # 图高上限：不超过正文区
IMG_PAD = 0.09                                   # 浅金衬底色块比图大一圈的边距
_IMG_CHARS_PER_LINE = 22                         # 半宽（约 5.8 英寸）下每行约 22 个全角字


# ================================================================ 基础手法（取经自 ppt_maker，独立实现）

def _rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str)


def _set_run_font(run, size_pt, color_hex, bold=False, spacing=None):
    """设置 run 字号/颜色/字体/字重/字距；中文字形必须显式写 East Asian typeface。
    spacing 为可选字距（磅），OOXML rPr spc 属性单位是 1/100 磅——用于封面主标、
    小标签这类「排面文字」，拉开字距立刻脱离程序生成感。"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    f = run.font
    f.size = Pt(size_pt)
    f.bold = bold
    f.color.rgb = _rgb(color_hex)
    f.name = FONT_NAME
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_NAME)
    if spacing:
        rPr.set("spc", str(int(spacing * 100)))   # 字距：2pt -> spc="200"


def _set_run_alpha(run, alpha_pct):
    """给 run 字体颜色加透明度（alpha_pct: 0-100，越小越透明）。
    用于章节页超大半透明序号、金句页大号引号这类「水印式」装饰。
    需在 _set_run_font 之后调用（依赖已存在的 solidFill）。"""
    from pptx.oxml.ns import qn
    try:
        rPr = run._r.get_or_add_rPr()
        solid = rPr.find(qn("a:solidFill"))
        if solid is None:
            return
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            return
        alpha = srgb.makeelement(qn("a:alpha"), {})
        alpha.set("val", str(int(alpha_pct * 1000)))   # OOXML：100% = 100000
        srgb.append(alpha)
    except Exception:
        pass   # 装饰性效果，失败不致命


def _set_fill_alpha(shape, alpha_pct):
    """给形状实心填充加透明度（alpha_pct: 0-100，越小越透明）。
    用于封面背景图上的深色蒙层。OOXML：100% = 100000。"""
    from pptx.oxml.ns import qn
    try:
        spPr = shape._element.spPr
        solid = spPr.find(qn("a:solidFill"))
        if solid is None:
            return
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            return
        alpha = srgb.makeelement(qn("a:alpha"), {})
        alpha.set("val", str(int(alpha_pct * 1000)))
        srgb.append(alpha)
    except Exception:
        pass   # 蒙层透明失败退化为实色，不致命


def _set_gradient_fill(shape, color_hex, top_alpha_pct, bottom_alpha_pct):
    """把形状填充改为垂直渐变蒙层：顶部 top_alpha% 不透明 -> 底部 bottom_alpha% 不透明。
    用于封面背景图「上浅下深」的二次收深——均匀蒙层保住全图可读性，
    底部渐变再压暗一层，给左下角标题块腾出稳定的明度对比。
    OOXML：gradFill/gsLst，lin ang 5400000 = 90°（自上而下）。失败静默，不致命。"""
    from pptx.oxml.ns import qn
    try:
        spPr = shape._element.spPr
        solid = spPr.find(qn("a:solidFill"))
        if solid is not None:
            spPr.remove(solid)   # 渐变取代实色（调用方先用 _add_rect 建实色形状）
        grad = spPr.makeelement(qn("a:gradFill"), {})
        gsLst = grad.makeelement(qn("a:gsLst"), {})
        for pos, alpha_pct in ((0, top_alpha_pct), (100000, bottom_alpha_pct)):
            gs = grad.makeelement(qn("a:gs"), {"pos": str(pos)})
            clr = grad.makeelement(qn("a:srgbClr"), {"val": color_hex})
            a = grad.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
            clr.append(a)
            gs.append(clr)
            gsLst.append(gs)
        grad.append(gsLst)
        grad.append(grad.makeelement(qn("a:lin"), {"ang": "5400000", "scaled": "1"}))
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            ln.addprevious(grad)   # gradFill 必须在 a:ln 之前（OOXML 序列约束）
        else:
            spPr.append(grad)
    except Exception:
        pass   # 渐变失败则保留原实色，不致命


def _set_bg(slide, color_hex):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(color_hex)


def _add_textbox(slide, left, top, width, height):
    """统一入口：所有文本框 word_wrap，防溢出第一道闸。"""
    from pptx.util import Inches
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _add_rect(slide, left, top, width, height, color_hex):
    """实心矩形（色条/装饰线/色带），去边框去阴影。"""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_accent_bar(slide, left, top, width, height):
    """页标题左侧的赭红色条 —— 全套版式统一的视觉锚点。"""
    return _add_rect(slide, left, top, width, height, COLOR_ACCENT)


def _add_round_rect(slide, left, top, width, height, color_hex, adj=0.08):
    """圆角实心矩形（衬底色块用），圆角半径 adj 为短边比例（默认 8%）。
    adjustments 失败静默回退默认圆角，不致命。"""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    try:
        shp.adjustments[0] = adj
    except Exception:
        pass   # 圆角半径设置失败用默认值，不致命
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _round_picture(pic, adj_pct=8):
    """把 Picture 的裁切几何改为圆角矩形：a:prstGeom prst="roundRect"，adj 8%。
    OOXML 形状调整值单位是千分之一百分比（8% = 8000）。
    任何失败静默回退直角（不改 prstGeom 就是默认 rect），绝不抛异常。"""
    from pptx.oxml.ns import qn
    try:
        spPr = pic._element.spPr
        prst = spPr.find(qn("a:prstGeom"))
        if prst is None:
            return
        prst.set("prst", "roundRect")
        avLst = prst.find(qn("a:avLst"))
        if avLst is None:
            avLst = prst.makeelement(qn("a:avLst"), {})
            prst.append(avLst)
        for gd in list(avLst):   # 清掉 rect 默认参数，只留我们的 adj
            avLst.remove(gd)
        gd = avLst.makeelement(qn("a:gd"), {})
        gd.set("name", "adj")
        gd.set("fmla", f"val {int(adj_pct * 1000)}")
        avLst.append(gd)
    except Exception:
        pass   # 圆角失败静默回退直角，不致命


# ================================================================ 防溢出：正文三档缩字号

def _fit_body_font(bullets, chars_per_line=None):
    """按总字数与估算行数三档缩字号（18 -> 16 -> 14），连带压缩段后距。
    chars_per_line 传半宽估值（图文模式约 22）时，按半宽重估阈值：
    行数阈值按宽度比放宽，总字数阈值同步下调，保证左文右图下不溢出。"""
    cpl = chars_per_line or _BODY_CHARS_PER_LINE
    total = sum(len(b) for b in bullets)
    lines = sum(max(1, (len(b) + 3 + cpl - 1) // cpl)
                for b in bullets)   # +3 是「●  」符号前缀
    if cpl < _BODY_CHARS_PER_LINE:        # 半宽（图文）模式阈值
        if total <= 180 and lines <= 12:
            return SIZE_BODY, 0.18
        if total <= 260 and lines <= 16:
            return SIZE_BODY_MID, 0.14
        return SIZE_BODY_MIN, 0.10
    if total <= 240 and lines <= 7:
        return SIZE_BODY, 0.18
    if total <= 320 and lines <= 9:
        return SIZE_BODY_MID, 0.14
    return SIZE_BODY_MIN, 0.10


# ================================================================ 容错：安全取值

def _safe_str(v, default=""):
    s = str(v).strip() if v is not None else ""
    return s or default


def _safe_str_list(v, max_items=None, max_len=120):
    """把任意输入归一成干净的字符串列表。"""
    out = []
    if isinstance(v, (list, tuple)):
        for item in v:
            s = _safe_str(item)
            if s:
                out.append(s[:max_len])
            if max_items and len(out) >= max_items:
                break
    return out


def _safe_bullets(page):
    """要点兜底：空 bullets 渲染「本页内容筹备中」，绝不空页。"""
    bullets = _safe_str_list(page.get("bullets"), max_items=8)
    return bullets or ["本页内容筹备中，详细要点以现场讲解为准。"]


def _valid_image_path(v):
    """配图路径有效性检查（路 B 交接契约）：必须是真实存在的文件，
    否则返回 None 当无图降级，绝不抛异常。"""
    if not isinstance(v, str) or not v.strip():
        return None
    p = v.strip()
    try:
        return p if os.path.isfile(p) else None
    except Exception:
        return None


# ================================================================ 共享装饰件

def _add_capsule(slide, left, top, text, fill_hex, text_hex, size=None):
    """风格标签胶囊：圆角拉满（adj 50% 成药丸）的小色块内嵌加粗小字。
    封面/章节页的「排面小件」——一个胶囊就能把「程序默认模板」感拉低一档。"""
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    text = _safe_str(text, "标签")
    w = 0.36 + len(text) * 0.19          # 中文字宽经验估算（12pt 约 0.17"/字 + 两侧留白）
    cap = _add_round_rect(slide, left, top, w, 0.36, fill_hex, adj=0.5)
    tf = cap.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_run_font(r, size or SIZE_NOTE, text_hex, bold=True, spacing=0.5)
    return cap


def _add_page_header(slide, page_title):
    """浅色内容页统一头部：左侧赭红小色块 + 26pt 墨色标题 + 页眉下浅灰细线。
    精修：色条改色块（比重更克制）、标题下加一道呼吸缝细线，版面骨架立刻清晰。"""
    _add_accent_bar(slide, MARGIN, 0.68, 0.16, 0.16)      # 小色块（原为 0.09x0.62 长条）
    tf = _add_textbox(slide, MARGIN + 0.32, HEADER_Y, CONTENT_W - 0.32, 0.9)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = _safe_str(page_title, "（无标题）")
    _set_run_font(r, SIZE_PAGE_TITLE, COLOR_DARK, bold=True)
    # 页眉下浅灰分隔细线：内容区从这条线以下开始，层级一目了然
    _add_rect(slide, MARGIN, HEADER_Y + 0.82, CONTENT_W, 0.008, COLOR_HAIRLINE)


def _add_footer(slide, page_no, total, dark=False):
    """页脚：细线 + 右侧页码「03 / 11」。深色版式换浅金线、浅金字。
    精修：页码降到 10pt 浅暖灰——保留导航功能，视觉上退到最弱一档。"""
    line_color = COLOR_GOLD if dark else COLOR_HAIRLINE
    text_color = COLOR_GOLD if dark else COLOR_CAPTION
    _add_rect(slide, MARGIN, FOOTER_Y, CONTENT_W, 0.012, line_color)
    tf = _add_textbox(slide, SLIDE_W - MARGIN - 2.0, FOOTER_Y + 0.06, 2.0, 0.32)
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page_no:02d} / {total:02d}"
    _set_run_font(r, SIZE_CAPTION, text_color)


def _write_note(slide, page, page_no):
    """演讲稿物理写入 notes_slide；缺失时按页型合成兜底稿，保证物理非空。"""
    note = _safe_str(page.get("speaker_note"))
    if not note:
        layout = _safe_str(page.get("layout"), "bullets")
        if layout == "quote":
            note = ("金句页：把这句话完整、放慢语速读一遍，停顿两秒，"
                    "再点出它与全篇主题的关系，然后自然翻页。建议用时约 30 秒。")
        elif layout == "section":
            note = (f"章节过渡：向听众宣告进入「{_safe_str(page.get('page_title'), '下一部分')}」，"
                    "用一句话概括本部分要回答的问题，语速放缓。建议用时约 20 秒。")
        elif layout == "toc":
            note = ("目录页：带着听众快速过一遍整体结构，说明每一部分之间的递进关系，"
                    "不要逐条展开细节。建议用时约 30 秒。")
        else:
            title = _safe_str(page.get("page_title"), f"第 {page_no} 页")
            note = (f"这一页讲「{title}」。开场先点明本页核心，再逐条展开，"
                    "讲完用一句小结收住并自然过渡。建议用时约 2 分钟。")
    try:
        slide.notes_slide.notes_text_frame.text = note
    except Exception:
        pass   # 备注写入失败不致命，绝不中断渲染


def _write_bullet_paras(tf, bullets, color_hex=COLOR_INK,
                        base_size=None, gap_in=None, chars_per_line=None):
    """写一串要点段落：实心圆点符号 + 正文文字（符号与文字分 run，符号单独上色）。
    符号精修：页内形状统一为「●」，颜色按要点序号在 MARKER_ROTATION 三态轮转；
    符号字号比正文小 2pt，圆点视觉比重更克制。base_size 为 None 时自动三档缩字号；
    chars_per_line 传半宽估值时按半宽重估。"""
    if base_size is None:
        size, gap = _fit_body_font(bullets, chars_per_line=chars_per_line)
    else:
        size, gap = base_size, (gap_in if gap_in is not None else 0.14)
    from pptx.util import Inches
    marker_size = max(size - 2, SIZE_NOTE)   # 圆点略小于正文，比重更精致
    for j, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Inches(gap)
        rm = p.add_run()
        rm.text = MARKER_GLYPH
        _set_run_font(rm, marker_size,
                      MARKER_ROTATION[j % len(MARKER_ROTATION)], bold=True)
        rt = p.add_run()
        rt.text = bullet
        _set_run_font(rt, size, color_hex)


# ================================================================ 版式 1：封面（render_deck 自动生成）

def _render_cover(prs, title, subtitle, style, cover_image=None, page_total=None,
                  content_type=""):
    """封面双构图（共用一套 tokens，按有无背景图分流）：

    无图 · 编辑部式左对齐：风格胶囊 + 赭红小方块 + 44pt 加字距主标 + 副标 +
    顶部/底部浅金细线 + 底部细信息行（风格 · 日期 · 页数）。
    有图 · 全出血：满铺背景图（直角，出血版式不圆角）+ 30% 均匀深棕蒙层
    （验收契约，保全图可读性）+ 底部二次渐变蒙层（上浅下深，给标题块压出明度差）
    + 左下角标题块（色条 + 主标 + 副标 + 信息行）。
    content_type：教学课件型时封面备注换教师开课口吻，其余保持通用分享口吻。
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _set_bg(s, COLOR_DARK)

    # 底部细信息行文案（两种构图共用）：风格 · 日期 · 总页数
    date_str = time.strftime("%Y年%m月%d日")
    info = f"{style} · {date_str}"
    if page_total:
        info += f" · 共 {page_total} 页"

    bg_img = _valid_image_path(cover_image)
    if bg_img:
        # ---- 有图构图：先铺图、再均匀蒙层、再底部渐变，文字层最上
        try:
            s.shapes.add_picture(
                bg_img, Inches(0), Inches(0),
                width=Inches(SLIDE_W), height=Inches(SLIDE_H))   # 16:9 全屏铺满
            overlay = _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, COLOR_DARK)
            _set_fill_alpha(overlay, 30)   # 均匀蒙层 30% 不透明（验收冻结契约）
            # 底部二次渐变：上浅下深，左下角标题块的明度底座
            grad = _add_rect(s, 0, SLIDE_H * 0.35, SLIDE_W, SLIDE_H * 0.65, COLOR_DEEP)
            _set_gradient_fill(grad, COLOR_DEEP, 0, 82)
        except Exception:
            pass   # 铺图失败退化为纯色封面，不致命

        # 左下角标题块：赭红色条 + 主标 + 副标 + 细信息行
        _add_accent_bar(s, MARGIN, 4.42, 1.6, 0.06)
        tf = _add_textbox(s, MARGIN, 4.62, CONTENT_W, 1.5)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_run_font(r, SIZE_COVER_TITLE, COLOR_BG, bold=True,
                      spacing=SPACING_COVER_TITLE)
        if subtitle:
            tf2 = _add_textbox(s, MARGIN, 6.0, CONTENT_W, 0.5)
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = subtitle
            _set_run_font(r2, SIZE_COVER_SUB, COLOR_GOLD)
        tf3 = _add_textbox(s, MARGIN, 6.62, CONTENT_W, 0.4)
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = info
        _set_run_font(r3, SIZE_NOTE, COLOR_CAPTION)
    else:
        # ---- 无图构图：编辑部式左对齐，细线与胶囊撑出设计感
        # 顶部与底部浅金细线（极细，框出呼吸边界而非「边框」）
        _add_rect(s, MARGIN, 0.72, CONTENT_W, 0.01, COLOR_GOLD)
        _add_rect(s, MARGIN, SLIDE_H - 0.72, CONTENT_W, 0.01, COLOR_GOLD)

        _add_capsule(s, MARGIN, 2.05, style, COLOR_ACCENT, COLOR_BG)

        # 赭红小方块 + 44pt 加字距主标（左对齐到网格线）
        _add_accent_bar(s, MARGIN, 2.72, 0.16, 0.16)
        tf = _add_textbox(s, MARGIN, 3.0, CONTENT_W, 1.7)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_run_font(r, SIZE_COVER_TITLE, COLOR_BG, bold=True,
                      spacing=SPACING_COVER_TITLE)

        if subtitle:
            tf2 = _add_textbox(s, MARGIN, 4.85, CONTENT_W, 0.6)
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = subtitle
            _set_run_font(r2, SIZE_COVER_SUB, COLOR_GOLD)

        # 底部细信息行：贴底线上方，最弱一档
        tf3 = _add_textbox(s, MARGIN, SLIDE_H - 1.25, CONTENT_W, 0.4)
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run()
        r3.text = info
        _set_run_font(r3, SIZE_NOTE, COLOR_CAPTION)

    if content_type == "教学课件型":
        s.notes_slide.notes_text_frame.text = (
            f"开场白：同学们好，我们开始上课。今天我们学习「{title}」"
            f"{('，' + subtitle) if subtitle else ''}。"
            "先用一句话点明这节课的核心目标，再带大家看一眼目录，"
            "交代本节课的知识脉络与节奏。建议用时约 1 分钟。")
    else:
        s.notes_slide.notes_text_frame.text = (
            f"开场白：各位好，今天分享的主题是「{title}」"
            f"{('，' + subtitle) if subtitle else ''}。"
            "先用一句话点明这次分享的核心价值，再交代整体结构，"
            "语速放慢，与听众做一次眼神交流。建议用时约 30 秒。")
    return s


# ================================================================ 版式 2：toc 目录页

def _render_toc(slide, page, page_no, style):
    """米白底：页标题 + 编号目录行；超过 8 条自动分两栏。"""
    from pptx.util import Inches
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "目录"))

    entries = _safe_str_list(page.get("entries"), max_items=20, max_len=40) \
        or ["（目录内容筹备中）"]
    two_col = len(entries) > 8
    if two_col:
        half = (len(entries) + 1) // 2
        columns = [entries[:half], entries[half:]]
        col_w = (CONTENT_W - 0.6) / 2
        col_x = [CONTENT_LEFT, CONTENT_LEFT + col_w + 0.6]
    else:
        columns = [entries]
        col_w = CONTENT_W
        col_x = [CONTENT_LEFT]

    for ci, col_entries in enumerate(columns):
        tf = _add_textbox(slide, col_x[ci], BODY_Y, col_w, 5.1)
        for j, entry in enumerate(col_entries):
            idx_global = j + 1 if ci == 0 else (len(columns[0]) + j + 1)
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Inches(0.16)
            rn = p.add_run()
            rn.text = f"{idx_global:02d}  "
            _set_run_font(rn, SIZE_BODY, COLOR_ACCENT, bold=True)
            rt = p.add_run()
            rt.text = entry
            _set_run_font(rt, SIZE_BODY, COLOR_INK)


# ================================================================ 版式 3：section 章节过渡页

def _render_section(slide, page, page_no, style):
    """深棕底：右侧超大半透明序号水印 + 导引小标签 + 34pt 浅标题 + 赭红色条 + 浅金核心句。
    精修：序号挪到右侧避免压字，新增「PART 02」导引行（字距拉开），层级更像画册章节页。"""
    _set_bg(slide, COLOR_DARK)

    # 超大半透明序号水印：靠右侧，15% 透明度，不与标题抢左侧视觉起点
    tf_no = _add_textbox(slide, SLIDE_W - MARGIN - 5.2, 0.9, 5.2, 3.0)
    p_no = tf_no.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p_no.alignment = PP_ALIGN.RIGHT
    r_no = p_no.add_run()
    r_no.text = f"{page_no:02d}"
    _set_run_font(r_no, SIZE_SECTION_NO, COLOR_BG, bold=True)
    _set_run_alpha(r_no, 15)

    # 导引小标签（kicker）：拉开字距的小号浅金字，画册式章节引子
    tf_k = _add_textbox(slide, MARGIN, 4.0, CONTENT_W, 0.4)
    p_k = tf_k.paragraphs[0]
    r_k = p_k.add_run()
    r_k.text = f"PART {page_no:02d}"
    _set_run_font(r_k, SIZE_NOTE, COLOR_GOLD, bold=True, spacing=SPACING_KICKER)

    # 赭红色条 + 章节标题
    _add_accent_bar(slide, MARGIN, 4.42, 1.6, 0.06)
    tf_t = _add_textbox(slide, MARGIN, 4.62, CONTENT_W, 1.0)
    p_t = tf_t.paragraphs[0]
    r_t = p_t.add_run()
    r_t.text = _safe_str(page.get("page_title"), f"第 {page_no} 部分")
    _set_run_font(r_t, SIZE_SECTION_TITLE, COLOR_BG, bold=True)

    core = _safe_str(page.get("core_point"))
    if core:
        tf_c = _add_textbox(slide, MARGIN, 5.72, CONTENT_W, 0.8)
        p_c = tf_c.paragraphs[0]
        r_c = p_c.add_run()
        r_c.text = core
        _set_run_font(r_c, SIZE_BODY_MID, COLOR_GOLD)


# ================================================================ 版式 4：bullets 标准要点页

def _add_rhythm_band(slide):
    """连续 bullets 段偶数页的节奏微变化：正文区右缘一条 12% 不透明浅金竖缘带。
    设计意图：翻页时「素-饰-素」交替，给眼睛一个极轻的节拍点，防止多页要点
    连续出现的视觉疲劳；12% 不透明度做到「说不出哪里变了，但感觉不闷」。"""
    band = _add_rect(slide, SLIDE_W - 0.24, BODY_Y - 0.25, 0.24,
                     FOOTER_Y - BODY_Y + 0.25, COLOR_GOLD)
    _set_fill_alpha(band, 12)   # 12% 不透明：存在感极低的浅金呼吸带


def _render_bullets(slide, page, page_no, style, _alt=False):
    """米白底：色条标题 + 圆点符号要点列，正文三档缩字号防溢出。
    page["image"] 有效时切换「左文右图」：文字区收窄到左 55%（半宽阈值重估
    缩档），右侧放等比圆角配图 + 浅金圆角衬底色块；无图或图片缺失时与原版式一致。
    _alt=True（连续 bullets 段偶数张、由 render_deck 按页序注入）且无配图时，
    右缘加浅金节奏缘带；有配图页本身已有视觉变化，不再叠加。"""
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "（无标题）"))
    bullets = _safe_bullets(page)

    img = _valid_image_path(page.get("image"))
    if not img:
        # 无图：与既有版式一致；连续段偶数页加节奏缘带（唯一的微变化点）
        if _alt:
            _add_rhythm_band(slide)
        tf = _add_textbox(slide, CONTENT_LEFT + 0.2, BODY_Y, CONTENT_W - 0.2, 5.1)
        _write_bullet_paras(tf, bullets)
        return

    # 有图：左文右图。文字区收窄，缩档阈值按半宽（约 22 全角字/行）重估
    tf = _add_textbox(slide, CONTENT_LEFT + 0.2, BODY_Y, IMG_TEXT_W - 0.2, IMG_MAX_H)
    _write_bullet_paras(tf, bullets, chars_per_line=_IMG_CHARS_PER_LINE)
    _place_side_image(slide, img, CONTENT_LEFT + IMG_TEXT_W + IMG_GAP, BODY_Y)


def _place_side_image(slide, img_path, left, top):
    """右图区配图：先按区宽等比铺图，超高则按高度回缩并水平居中；
    图下垫一圈浅金衬底色块（成熟产品手法，z-order 衬底在图之下）。
    精修：图片与衬底同步 8% 圆角，观感从「贴上去」变「嵌进去」；失败静默回退直角。
    任何失败都静默降级为纯文字页，绝不抛异常。"""
    from pptx.util import Inches, Emu
    try:
        pic = slide.shapes.add_picture(
            img_path, Inches(left), Inches(top), width=Inches(IMG_ZONE_W))  # 只给宽度，等比
        if pic.height > Inches(IMG_MAX_H):        # 图高不超过正文区：按高度回缩
            ratio = Inches(IMG_MAX_H) / pic.height
            pic.height = Inches(IMG_MAX_H)
            pic.width = Emu(int(pic.width * ratio))
        # 在图区内水平居中
        pic.left = Emu(int(Inches(left) + (Inches(IMG_ZONE_W) - pic.width) / 2))
        _round_picture(pic, 8)                    # 图片圆角裁切（失败静默回退直角）
        # 浅金衬底色块：比图大一圈，垫在图下，同步圆角
        pad = Inches(IMG_PAD)
        backing = _add_round_rect(slide,
                                  (pic.left - pad) / 914400, (pic.top - pad) / 914400,
                                  (pic.width + 2 * pad) / 914400, (pic.height + 2 * pad) / 914400,
                                  COLOR_GOLD, adj=0.08)
        pic._element.addprevious(backing._element)   # z-order：衬底挪到图片之下
    except Exception:
        pass   # 配图失败静默降级，绝不中断渲染


# ================================================================ 版式 5：two_column 两栏对比页

def _render_two_column(slide, page, page_no, style):
    """米白底：左右两栏，栏头赭红小条 + 18pt 栏题，栏间浅金竖线分隔。
    page["image"] 有效时启用「双栏 + 底部横图」：栏正文区压高、分隔线同步缩短，
    底部居中放等比圆角横图 + 浅金衬底；无图或图片缺失时与原版式一致。"""
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "（无标题）"))

    img = _valid_image_path(page.get("image"))
    # 有图时正文区与分隔线压高，给底部横图让位（无图时与原版式逐像素一致）
    pts_h = 2.45 if img else 4.1
    line_h = 3.35 if img else 4.9

    col_w = (CONTENT_W - 0.8) / 2
    col_x = [CONTENT_LEFT, CONTENT_LEFT + col_w + 0.8]

    # 栏间竖分隔线
    _add_rect(slide, CONTENT_LEFT + col_w + 0.4, BODY_Y + 0.1, 0.012, line_h, COLOR_GOLD)

    for ci, key in enumerate(("left", "right")):
        col = page.get(key)
        if not isinstance(col, dict):
            col = {}
        heading = _safe_str(col.get("heading"), "要点")
        points = _safe_str_list(col.get("points"), max_items=6) \
            or ["本栏内容筹备中"]

        # 栏头：赭红小色条 + 栏题
        _add_accent_bar(slide, col_x[ci], BODY_Y + 0.08, 0.5, 0.05)
        tf_h = _add_textbox(slide, col_x[ci], BODY_Y + 0.22, col_w, 0.55)
        ph = tf_h.paragraphs[0]
        rh = ph.add_run()
        rh.text = heading
        _set_run_font(rh, SIZE_BODY, COLOR_DARK, bold=True)

        tf_p = _add_textbox(slide, col_x[ci], BODY_Y + 0.9, col_w, pts_h)
        size, gap = _fit_body_font(points)
        # 栏宽只有一半，字号上限压到 16，避免半栏宽下放 18pt 溢出
        _write_bullet_paras(tf_p, points, base_size=min(size, SIZE_BODY_MID), gap_in=gap)

    if img:
        _place_bottom_image(slide, img, BODY_Y + 3.0, 2.15)


def _place_bottom_image(slide, img_path, top, max_h):
    """底部横图（two_column 配图位）：先按限高等比铺图、水平居中；
    图下垫一圈浅金衬底色块，与右图区同款 8% 圆角；任何失败静默降级。"""
    from pptx.util import Inches, Emu
    try:
        pic = slide.shapes.add_picture(
            img_path, Inches(CONTENT_LEFT), Inches(top), height=Inches(max_h))  # 只给高度，等比
        if pic.width > Inches(CONTENT_W):         # 超宽则按内容区宽回缩
            ratio = Inches(CONTENT_W) / pic.width
            pic.width = Inches(CONTENT_W)
            pic.height = Emu(int(pic.height * ratio))
        pic.left = Emu(int(Inches(CONTENT_LEFT) + (Inches(CONTENT_W) - pic.width) / 2))
        _round_picture(pic, 8)                    # 与右图区一致的圆角裁切（失败静默回退直角）
        pad = Inches(IMG_PAD)
        backing = _add_round_rect(slide,
                                  (pic.left - pad) / 914400, (pic.top - pad) / 914400,
                                  (pic.width + 2 * pad) / 914400, (pic.height + 2 * pad) / 914400,
                                  COLOR_GOLD, adj=0.08)
        pic._element.addprevious(backing._element)   # z-order：衬底挪到图片之下
    except Exception:
        pass   # 配图失败静默降级，绝不中断渲染


# ================================================================ 版式 6：big_number 大数字页

def _render_big_number(slide, page, page_no, style):
    """米白底：1-3 个 66pt 赭红大数字横排，上方装饰短线，下方 14pt 注释（限 40 字）。"""
    from pptx.enum.text import PP_ALIGN
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "关键数字"))

    stats_raw = page.get("stats")
    stats = []
    if isinstance(stats_raw, (list, tuple)):
        for it in stats_raw[:3]:
            if not isinstance(it, dict):
                continue
            number = _safe_str(it.get("number"), "—")
            caption = _safe_str(it.get("caption"), "")[:40]   # 注释限 40 字
            stats.append((number, caption))
    if not stats:
        stats = [("—", "数据筹备中")]

    n = len(stats)
    col_w = CONTENT_W / n
    for i, (number, caption) in enumerate(stats):
        cx = CONTENT_LEFT + i * col_w
        # 数字上方居中的装饰短线
        _add_accent_bar(slide, cx + col_w / 2 - 0.4, 2.5, 0.8, 0.05)

        tf_n = _add_textbox(slide, cx, 2.75, col_w, 1.5)
        pn = tf_n.paragraphs[0]
        pn.alignment = PP_ALIGN.CENTER
        rn = pn.add_run()
        rn.text = number
        _set_run_font(rn, SIZE_BIG_NUMBER, COLOR_ACCENT, bold=True)

        if caption:
            tf_c = _add_textbox(slide, cx + 0.3, 4.45, col_w - 0.6, 1.0)
            pc = tf_c.paragraphs[0]
            pc.alignment = PP_ALIGN.CENTER
            rc = pc.add_run()
            rc.text = caption
            _set_run_font(rc, SIZE_BODY_MIN, COLOR_WARMGREY)   # 注释退到暖灰，不抢数字


# ================================================================ 版式 7：quote 金句引用页

def _render_quote(slide, page, page_no, style):
    """深棕底：大号半透明引号装饰 + 26pt 居中浅字金句 + 浅金署名。"""
    from pptx.enum.text import PP_ALIGN
    _set_bg(slide, COLOR_DARK)

    # 大号引号装饰（20% 透明度，水印式）
    tf_q = _add_textbox(slide, MARGIN + 0.3, 0.9, 3.0, 2.4)
    pq = tf_q.paragraphs[0]
    rq = pq.add_run()
    rq.text = "“"
    _set_run_font(rq, SIZE_QUOTE_MARK, COLOR_ACCENT, bold=True)
    _set_run_alpha(rq, 25)

    quote = _safe_str(page.get("quote"), "金句内容筹备中。")
    tf = _add_textbox(slide, MARGIN + 0.8, 2.9, CONTENT_W - 1.6, 2.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = quote
    _set_run_font(r, SIZE_PAGE_TITLE, COLOR_BG, bold=True)

    attribution = _safe_str(page.get("attribution"))
    if attribution:
        _add_accent_bar(slide, SLIDE_W / 2 - 0.6, 5.35, 1.2, 0.05)
        tf_a = _add_textbox(slide, MARGIN, 5.55, CONTENT_W, 0.6)
        pa = tf_a.paragraphs[0]
        pa.alignment = PP_ALIGN.CENTER
        ra = pa.add_run()
        ra.text = f"—— {attribution}"
        _set_run_font(ra, SIZE_BODY_MID, COLOR_GOLD)


# ================================================================ 版式 8：chart 图表页（原生可编辑图表）

def _render_chart(slide, page, page_no, style):
    """米白底：左侧原生图表（柱/条/饼/线）+ 右侧 2-3 条解读要点。
    图表为 GraphicFrame（PPT 内双击可编辑数据），不是贴图。"""
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "数据图表"))

    chart = page.get("chart")
    if not isinstance(chart, dict):
        raise ValueError("chart 页缺 chart 数据")   # 交上层降级为 bullets 页
    ctype = _safe_str(chart.get("type"), "bar").lower()
    type_map = _chart_type_map()
    if ctype not in type_map:
        ctype = "bar"

    categories = _safe_str_list(chart.get("categories"), max_items=12, max_len=20)
    series_list = []
    raw_series = chart.get("series")
    if isinstance(raw_series, (list, tuple)):
        for s in raw_series[:5]:
            if not isinstance(s, dict):
                continue
            name = _safe_str(s.get("name"), f"系列 {len(series_list) + 1}")
            values = []
            raw_v = s.get("values")
            if isinstance(raw_v, (list, tuple)):
                for v in raw_v[:len(categories) or 12]:
                    try:
                        values.append(float(v))
                    except (TypeError, ValueError):
                        values.append(0.0)
            series_list.append((name, values))
    if not categories or not series_list:
        raise ValueError("chart 数据不完整")   # 交上层降级

    # 系列值长度对齐类目数：少了补 0，多了截断，防空列报错
    for i, (name, values) in enumerate(series_list):
        if len(values) < len(categories):
            values = values + [0.0] * (len(categories) - len(values))
        series_list[i] = (name, values[:len(categories)])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_list:
        chart_data.add_series(name, values)

    # 左图右文：图表占左 7.3 英寸，右栏 3.9 英寸放解读
    gf = slide.shapes.add_chart(
        type_map[ctype],
        Inches(CONTENT_LEFT), Inches(BODY_Y),
        Inches(7.3), Inches(5.0),
        chart_data)
    ch = gf.chart

    # 全局字体
    try:
        ch.font.size = Pt(11)
        ch.font.name = FONT_NAME
        ch.font.color.rgb = _rgb(COLOR_INK)
    except Exception:
        pass

    # 系列配色（暖色系轮转）；单系列柱/条图按类目逐点上色，更耐看
    palette = CHART_PALETTE
    if ctype == "pie":
        serie = ch.series[0]
        for pi, point in enumerate(serie.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(palette[pi % len(palette)])
    elif len(series_list) == 1 and ctype in ("bar", "hbar"):
        serie = ch.series[0]
        for pi, point in enumerate(serie.points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(palette[pi % len(palette)])
    else:
        for si, serie in enumerate(ch.series):
            color = palette[si % len(palette)]
            if ctype == "line":
                serie.format.line.color.rgb = _rgb(color)
                serie.format.line.width = Pt(2.25)
            else:
                serie.format.fill.solid()
                serie.format.fill.fore_color.rgb = _rgb(color)

    # 数据标签：饼图显示百分比，柱/条/线显示数值
    try:
        plot = ch.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        if ctype == "pie":
            dl.show_percentage = True
            dl.show_value = False
            dl.number_format = "0%"
            dl.number_format_is_linked = False
        else:
            dl.show_value = True
        dl.font.size = Pt(10)
        dl.font.name = FONT_NAME
    except Exception:
        pass

    # 图例：饼图放右侧（看类目），多系列放底部，单系列柱/条不显示
    try:
        from pptx.enum.chart import XL_LEGEND_POSITION
        if ctype == "pie":
            ch.has_legend = True
            ch.legend.position = XL_LEGEND_POSITION.RIGHT
        elif len(series_list) > 1:
            ch.has_legend = True
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        else:
            ch.has_legend = False
        if ch.has_legend:
            ch.legend.include_in_layout = False
    except Exception:
        pass

    # 右侧解读栏：小标题 + 2-3 条要点
    bullets = _safe_str_list(page.get("bullets"), max_items=3) \
        or ["图表解读筹备中"]
    tf_h = _add_textbox(slide, CONTENT_LEFT + 7.6, BODY_Y + 0.05, 3.9, 0.5)
    ph = tf_h.paragraphs[0]
    rh = ph.add_run()
    rh.text = _safe_str(chart.get("title"), "解读")
    _set_run_font(rh, SIZE_BODY_MID, COLOR_DARK, bold=True)
    tf_b = _add_textbox(slide, CONTENT_LEFT + 7.6, BODY_Y + 0.6, 3.9, 4.3)
    _write_bullet_paras(tf_b, bullets, base_size=SIZE_BODY_MIN, gap_in=0.12)


# ================================================================ 版式 9：closing 结尾行动页

def _render_closing(slide, page, page_no, style, content_type=""):
    """米白底：色块标题 + 行动要点 + 居中仪式感收尾（短细线 + 谢谢聆听 + 落款）。
    精修：弃用底部大色带（工程感重），改用画册尾页式的居中三段式。
    落款随 content_type 调口吻：教学课件型「温故知新」，讲话动员型「期待行动」，
    分析报告型「有据可依」，缺省「期待行动」。"""
    from pptx.enum.text import PP_ALIGN
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "总结与行动"))

    bullets = _safe_bullets(page)
    tf = _add_textbox(slide, CONTENT_LEFT + 0.2, BODY_Y, CONTENT_W - 0.2, 3.3)
    _write_bullet_paras(tf, bullets)

    # 居中仪式感三段：短细线 -> 谢谢聆听 -> 落款（风格 · 口吻词，随内容类型变化）
    _CLOSING_TONE = {"教学课件型": "温故知新", "分析报告型": "有据可依"}
    tone = _CLOSING_TONE.get(content_type, "期待行动")
    _add_rect(slide, SLIDE_W / 2 - 0.7, 5.35, 1.4, 0.014, COLOR_GOLD)
    tf_t = _add_textbox(slide, MARGIN, 5.55, CONTENT_W, 0.7)
    pt = tf_t.paragraphs[0]
    pt.alignment = PP_ALIGN.CENTER
    rt = pt.add_run()
    rt.text = "谢谢聆听"
    _set_run_font(rt, SIZE_PAGE_TITLE, COLOR_DARK, bold=True, spacing=SPACING_KICKER)
    tf_b = _add_textbox(slide, MARGIN, 6.35, CONTENT_W, 0.4)
    pb = tf_b.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = f"{style} · {tone}"
    _set_run_font(rb, SIZE_NOTE, COLOR_CAPTION)


# ================================================================ 降级兜底页

def _render_safe_fallback(slide, page, page_no, style):
    """单页渲染异常时的保底：按 bullets 版式渲染安全内容，绝不中断整单。"""
    try:
        _render_bullets(slide, page, page_no, style)
    except Exception:
        try:
            _set_bg(slide, COLOR_BG)
            tf = _add_textbox(slide, CONTENT_LEFT, BODY_Y, CONTENT_W, 3.0)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "本页内容筹备中"
            _set_run_font(r, SIZE_BODY, COLOR_INK)
        except Exception:
            pass


_LAYOUTS = {
    "toc": _render_toc,
    "section": _render_section,
    "bullets": _render_bullets,
    "two_column": _render_two_column,
    "big_number": _render_big_number,
    "quote": _render_quote,
    "chart": _render_chart,
    "closing": _render_closing,
}

# ---- 图示版式（ppt_diagrams：timeline/process/pyramid）：防御式导入，
# 模块缺席时这三种 layout 经「未知 layout 降级 bullets」路径自动兜底
try:
    from ppt_diagrams import DIAGRAM_RENDERERS as _DIAGRAM_RENDERERS
except Exception:  # noqa: BLE001 - 缺席不阻断，降级即可
    _DIAGRAM_RENDERERS = {}
_LAYOUTS.update(_DIAGRAM_RENDERERS)


# ================================================================ 公开 API

def render_deck(prs, deck: dict, style: str = "工作汇报") -> None:
    """把归一化后的 deck 字典渲染进 python-pptx 的 Presentation（16:9 已设好）。

    deck = {
      "title": str, "subtitle": str,
      "cover_image": str | None,   # 可选：封面背景图绝对路径，缺失/无效当无图
      "pages": [ page, ... ]   # 每个 page 必含 "layout" 与 "speaker_note"
    }
    page 按 layout 分八种：toc / section / bullets / two_column /
                          big_number / quote / chart / closing。
    page 可选键 "image"：配图绝对路径，bullets（左文右图）与 two_column
    （底部横图）版式响应，文件不存在时自动降级为无图版式，绝不抛异常。
    deck 可选键 "content_type"：讲话动员型/教学课件型/分析报告型，
    影响封面备注口吻与结尾落款措辞，缺省保持通用口吻。
    每页都要把 speaker_note 写进该页 notes_slide（物理写入，非注释）。
    封面由 render_deck 自动用 deck["title"]/["subtitle"] 生成，pages 里不含封面。

    容错契约：page 缺字段给安全默认；未知 layout 按 bullets 渲染；单页渲染
    异常降级为保底要点页；任何情况下不向调用方抛异常。
    """
    deck = deck if isinstance(deck, dict) else {}
    style = _safe_str(style, "工作汇报")
    title = _safe_str(deck.get("title"), "未命名汇报")
    subtitle = _safe_str(deck.get("subtitle"))
    cover_image = deck.get("cover_image")   # 可选：封面背景图绝对路径（路 B 契约）
    pages_raw = deck.get("pages")
    pages = [p if isinstance(p, dict) else {} for p in pages_raw] \
        if isinstance(pages_raw, list) else []

    blank = prs.slide_layouts[6]

    # ---- 封面（无图编辑部式 / 有图全出血双构图；cover_image 无效自动降级无图）
    content_type = _safe_str(deck.get("content_type"))   # 可选：口吻提示（路 B 契约）
    _render_cover(prs, title, subtitle, style, cover_image=cover_image,
                  page_total=len(pages) + 1,   # 底部信息行「共 N 页」用
                  content_type=content_type)

    total = len(pages) + 1   # 物理总页数（含封面），页脚用

    # ---- 内容页
    consec_bullets = 0   # 连续 bullets 段计数：段内偶数张注入节奏缘带（精修 3）
    for i, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(blank)
        layout = _safe_str(page.get("layout"), "bullets").lower()
        render_fn = _LAYOUTS.get(layout, _render_bullets)   # 未知 layout 降级 bullets
        # 只按「声明的 bullets 版式」计数；被其他版式（含降级）打断即归零重数
        if layout == "bullets":
            consec_bullets += 1
        else:
            consec_bullets = 0
        try:
            if layout == "bullets":
                render_fn(slide, page, i, style, _alt=(consec_bullets % 2 == 0))
            elif layout == "closing":
                render_fn(slide, page, i, style, content_type=content_type)
            else:
                render_fn(slide, page, i, style)
        except Exception:
            _render_safe_fallback(slide, page, i, style)
        _add_footer(slide, i + 1, total, dark=(layout in DARK_LAYOUTS))
        _write_note(slide, page, i)
