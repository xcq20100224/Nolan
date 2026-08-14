# -*- coding: utf-8 -*-
"""
jarvis/test_theme.py —— 主题精修验收（真实渲染，不走 mock）

覆盖：
  1. 8 版式全覆盖 deck 渲染不抛异常，页数 = pages + 1（封面）；
  2. 每页备注物理非空；
  3. 字体契约：run 级 font.name 与 East Asian typeface 均为微软雅黑；
  4. 封面：胶囊、44pt 加字距主标、底部信息行「共 N 页」；
  5. 要点符号三态轮转色、页脚 10pt 浅暖灰「03 / 11」；
  6. 章节页水印序号 + PART 导引；closing 居中「谢谢聆听」；
  7. 有图封面：图 -> 30% 均匀蒙层 -> 渐变蒙层 -> 文字 z-order；
  8. 冒烟产物保存到 jarvis/files/theme_preview.pptx 供主控检查。

运行：cd jarvis && python test_theme.py（全部断言通过即全绿）
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_layouts import (   # noqa: E402
    render_deck, COLOR_DARK, COLOR_BG, COLOR_ACCENT, COLOR_WARMGREY,
    COLOR_TAN, COLOR_GOLD, COLOR_CAPTION, COLOR_HAIRLINE, COLOR_DEEP,
    FONT_NAME, SIZE_COVER_TITLE, SIZE_CAPTION, MARKER_ROTATION,
)

FILES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "files")
PREVIEW_PATH = os.path.join(FILES_DIR, "theme_preview.pptx")
TEST_IMG = os.path.join(FILES_DIR, "theme_test_cover.png")


# ---------------------------------------------------------------- 工具

def _make_test_image():
    """用 PIL 生成一张暖色 16:9 测试图（封面背景与 bullets 配图共用）。"""
    from PIL import Image
    os.makedirs(FILES_DIR, exist_ok=True)
    img = Image.new("RGB", (1600, 900))
    # 纵向暖色渐变：陶棕 -> 深棕，模拟真实摄影图而非纯色块
    top, bottom = (168, 123, 95), (59, 50, 44)
    px = img.load()
    for y in range(900):
        t = y / 899
        row = tuple(int(top[i] + (bottom[i] - top[i]) * t) for i in range(3))
        for x in range(1600):
            px[x, y] = row
    img.save(TEST_IMG)
    return TEST_IMG


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _reopen(prs, path):
    """存盘重开：按 OOXML 落盘后的真实状态断言，排除内存态假象。"""
    prs.save(path)
    return Presentation(path)


def _runs(slide):
    """收集一页上所有文本 run。"""
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            out.extend(p.runs)
    return out


def _fill_hex(shape):
    """形状实心填充色（无实色填充返回 ''）。"""
    try:
        srgb = shape._element.spPr.find(qn("a:solidFill"))
        if srgb is None:
            return ""
        clr = srgb.find(qn("a:srgbClr"))
        return str(clr.get("val")) if clr is not None else ""
    except Exception:
        return ""


def _full_deck(img):
    """8 版式全覆盖 + 连续 bullets 段 + 配图页 + chart 数据。"""
    return {
        "title": "工程审美到设计审美",
        "subtitle": "一套精品默认主题的精修实践",
        "pages": [
            {"layout": "toc", "page_title": "目录",
             "entries": ["现状诊断", "设计 tokens", "字体层级", "封面重生",
                         "装饰系统", "页脚精修"],
             "speaker_note": "目录页备注。"},
            {"layout": "section", "page_title": "现状诊断",
             "core_point": "工程审美不是丑，是缺少最后一公里的克制。",
             "speaker_note": "章节页备注。"},
            {"layout": "bullets", "page_title": "配色平、层级弱、装饰少",
             "bullets": ["配色平：只有深浅两档，缺少中间调",
                         "字体层级弱：标题正文注释拉不开差距",
                         "装饰少：没有胶囊、细线、水印这些排面小件"],
             "speaker_note": "要点页备注。"},
            {"layout": "bullets", "page_title": "连续段第二页（节奏缘带）",
             "bullets": ["偶数张右缘出现 12% 浅金缘带", "素-饰-素节奏防视觉疲劳"],
             "speaker_note": "要点页备注。"},
            {"layout": "bullets", "page_title": "左文右图配图页",
             "bullets": ["配图 8% 圆角裁切", "浅金衬底垫在图下", "图片失败静默降级"],
             "image": img, "speaker_note": "配图页备注。"},
            {"layout": "two_column", "page_title": "精修前后对比",
             "left": {"heading": "精修前", "points": ["大色带收尾", "页码 12pt 暖灰", "居中封面"]},
             "right": {"heading": "精修后", "points": ["细线加落款", "页码 10pt 浅暖灰", "左对齐编辑部封面"]},
             "speaker_note": "对比页备注。"},
            {"layout": "big_number", "page_title": "关键数字",
             "stats": [{"number": "6", "caption": "设计 tokens 主色"},
                        {"number": "4", "caption": "字体层级"},
                        {"number": "8", "caption": "精修版式"}],
             "speaker_note": "数字页备注。"},
            {"layout": "quote", "page_title": "设计观",
             "quote": "好的默认主题，让人说不出哪里好，但第一眼就信任。",
             "attribution": "精修师手记", "speaker_note": "金句页备注。"},
            {"layout": "chart", "page_title": "视觉收益排序",
             "chart": {"type": "bar", "title": "收益",
                        "categories": ["配色", "字体", "封面", "装饰", "页脚"],
                        "series": [{"name": "收益分", "values": [9, 8, 10, 7, 5]}]},
             "bullets": ["封面重生收益最大", "配色 tokens 是地基", "页脚精修成本最低"],
             "speaker_note": "图表页备注。"},
            {"layout": "closing", "page_title": "总结与行动",
             "bullets": ["tokens 常量区统一引用", "四级字体层级落地", "封面双构图上线"],
             "speaker_note": "结尾页备注。"},
        ],
    }


# ---------------------------------------------------------------- 测试项

def test_full_deck_render():
    """8 版式全覆盖渲染：页数、备注、字体契约、各版式视觉要点。"""
    img = _make_test_image()
    deck = _full_deck(img)
    prs = _new_prs()
    render_deck(prs, deck, style="设计分享")
    prs = _reopen(prs, PREVIEW_PATH)

    # 1) 页数 = 10 内容页 + 1 封面
    assert len(prs.slides) == 11, f"应为 11 页，实际 {len(prs.slides)}"

    # 2) 每页备注物理非空 + shape 数量合理（封面/内容页都不应是空壳）
    for i, slide in enumerate(prs.slides):
        note = slide.notes_slide.notes_text_frame.text.strip()
        assert note, f"第 {i + 1} 页备注为空"
        n_shapes = len(list(slide.shapes))
        assert n_shapes >= 3, f"第 {i + 1} 页只有 {n_shapes} 个 shape，疑似渲染残缺"

    # 3) 字体契约：所有 run 都是微软雅黑（含 East Asian typeface）
    checked = 0
    for slide in prs.slides:
        for r in _runs(slide):
            if not r.text.strip():
                continue
            assert r.font.name == FONT_NAME, f"run 字体名应为 {FONT_NAME}"
            rPr = r._r.find(qn("a:rPr"))
            ea = rPr.find(qn("a:ea")) if rPr is not None else None
            assert ea is not None and ea.get("typeface") == FONT_NAME, \
                "run 缺 East Asian typeface"
            checked += 1
    assert checked > 60, f"全 deck 至少应有 60 个文本 run，实际 {checked}"

    # 4) 封面（无图）：胶囊 + 44pt 加字距主标 + 底部信息行
    cover = prs.slides[0]
    cover_runs = _runs(cover)
    title_run = next(r for r in cover_runs if r.text == deck["title"])
    assert abs(title_run.font.size.pt - SIZE_COVER_TITLE) < 0.01, "封面主标应为 44pt"
    assert title_run.font.bold, "封面主标应加粗"
    spc = title_run._r.find(qn("a:rPr")).get("spc")
    assert spc == "200", f"封面主标字距应为 spc=200，实际 {spc}"
    capsule = [sh for sh in cover.shapes
               if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
               and _fill_hex(sh) == COLOR_ACCENT and sh.has_text_frame
               and sh.text_frame.text == "设计分享"]
    assert capsule, "封面应有赭红风格胶囊"
    assert any("共 11 页" in (r.text or "") for r in cover_runs), \
        "封面底部信息行应含总页数"

    # 5) 要点符号三态轮转色（slides[3] = pages[2]，3 个要点）
    b_runs = [r for r in _runs(prs.slides[3]) if r.text.strip() == "●"]
    assert len(b_runs) == 3, f"要点页应有 3 个要点符号，实际 {len(b_runs)}"
    for k, r in enumerate(b_runs):
        assert str(r.font.color.rgb) == MARKER_ROTATION[k % 3], \
            f"第 {k + 1} 个符号色应为 {MARKER_ROTATION[k % 3]}"

    # 6) 页脚精修：slides[3]（浅色 bullets 页，页序 3）页脚「04 / 11」，10pt 浅暖灰
    foot = [r for r in _runs(prs.slides[3]) if r.text.strip() == "04 / 11"]
    assert foot, "要点页应有页脚页码 04 / 11"
    assert abs(foot[0].font.size.pt - SIZE_CAPTION) < 0.01, "页码应为 10pt"
    assert str(foot[0].font.color.rgb) == COLOR_CAPTION, "页码应为浅暖灰"

    # 7) 章节页（slides[2] = pages[1]）：深底 + 水印序号 + PART 导引
    sec_runs = _runs(prs.slides[2])
    assert any(r.text == "02" for r in sec_runs), "章节页应有大序号 02"
    assert any(r.text == "PART 02" for r in sec_runs), "章节页应有 PART 02 导引"
    wm = next(r for r in sec_runs if r.text == "02")
    srgb = wm._r.find(qn("a:rPr")).find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgb.find(qn("a:alpha"))
    assert alpha is not None and alpha.get("val") == "15000", "水印序号 alpha 应为 15%"

    # 8) 页眉细线：浅色页存在浅灰 hairline 形状
    hairlines = [sh for sh in prs.slides[3].shapes
                 if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                 and _fill_hex(sh) == COLOR_HAIRLINE]
    assert hairlines, "浅色内容页应有页眉/页脚浅灰细线"

    # 9) 配图页（slides[5]）：圆角图 + 浅金衬底
    pics = [sh for sh in prs.slides[5].shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 1, "配图页应有 1 张图"
    assert 'prst="roundRect"' in pics[0]._element.xml, "配图应圆角裁切"
    assert any(_fill_hex(sh) == COLOR_GOLD for sh in prs.slides[5].shapes), \
        "配图页应有浅金衬底"

    # 10) chart 页（slides[9]）：原生图表 1 个
    charts = [sh for sh in prs.slides[9].shapes if sh.has_chart]
    assert len(charts) == 1, "chart 页应有 1 个原生图表"

    # 11) closing 页（slides[10]）：居中谢谢聆听 + 落款，无大色带
    cl_runs = _runs(prs.slides[10])
    assert any(r.text == "谢谢聆听" for r in cl_runs), "closing 应有居中谢谢聆听"
    assert any("期待行动" in (r.text or "") for r in cl_runs), "closing 应有落款"
    thanks = next(r for r in cl_runs if r.text == "谢谢聆听")
    assert thanks.font.bold, "谢谢聆听应加粗"
    band = [sh for sh in prs.slides[10].shapes
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and _fill_hex(sh) == COLOR_ACCENT
            and sh.width > Inches(10)]   # 大色带宽度接近全幅
    assert not band, "closing 不应再有大色带收尾"

    print("[PASS] test_full_deck_render: 11 页 / 备注 / 字体契约 / "
          "封面胶囊+字距 / 符号轮转 / 页脚 10pt / 章节水印 / 配图圆角 / "
          "chart / closing 三段式 全部正确")
    return prs


def test_cover_with_image():
    """有图封面：图 -> 30% 均匀蒙层（契约）-> 渐变蒙层 -> 文字 z-order。"""
    img = _make_test_image()
    deck = _full_deck(img)
    deck["cover_image"] = img
    prs = _new_prs()
    render_deck(prs, deck, style="设计分享")
    tmp = os.path.join(FILES_DIR, "theme_preview_cover_img.pptx")
    prs = _reopen(prs, tmp)

    slide = prs.slides[0]
    shapes = list(slide.shapes)
    pics = [(i, sh) for i, sh in enumerate(shapes)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics, "有图封面应有背景图"
    pic_idx = pics[0][0]

    overlay_idx = next((i for i, sh in enumerate(shapes)
                        if i > pic_idx and _fill_hex(sh) == COLOR_DARK), None)
    assert overlay_idx is not None, "背景图之上应有均匀蒙层"
    srgb = shapes[overlay_idx]._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgb.find(qn("a:alpha"))
    assert alpha is not None and alpha.get("val") == "30000", "均匀蒙层 alpha 应为 30000"

    grad = [sh for sh in shapes
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and sh._element.spPr.find(qn("a:gradFill")) is not None]
    assert grad, "应有底部渐变蒙层（上浅下深）"
    gs = grad[0]._element.spPr.find(qn("a:gradFill")).find(qn("a:gsLst"))
    stops = gs.findall(qn("a:gs"))
    assert len(stops) == 2, "渐变应为双停止点"
    top_a = stops[0].find(qn("a:srgbClr")).find(qn("a:alpha")).get("val")
    bot_a = stops[1].find(qn("a:srgbClr")).find(qn("a:alpha")).get("val")
    assert top_a == "0" and bot_a == "82000", \
        f"渐变应上 0% 下 82%，实际 {top_a}/{bot_a}"
    assert stops[0].find(qn("a:srgbClr")).get("val") == COLOR_DEEP, "渐变色应为深墨"

    title_idx = next(i for i, sh in enumerate(shapes)
                     if sh.has_text_frame and deck["title"] in sh.text_frame.text)
    assert pic_idx < overlay_idx < title_idx, "z-order 应为 图<蒙层<文字"

    os.remove(tmp)   # 有图变体只是断言用，不留产物
    print("[PASS] test_cover_with_image: 图/均匀蒙层(30000)/渐变(0->82%)/文字 "
          "z-order 与渐变停止点全部正确")


def test_unknown_layout_degrades():
    """未知 layout 降级 bullets，绝不抛异常；无效配图静默降级。"""
    deck = {"title": "容错", "subtitle": "",
            "pages": [{"layout": "不存在的版式", "page_title": "降级页",
                        "bullets": ["要点甲"], "image": "Z:/绝不存在的图.png",
                        "speaker_note": "备注。"}]}
    prs = _new_prs()
    render_deck(prs, deck)   # 不抛异常即过半
    prs = _reopen(prs, PREVIEW_PATH)   # 覆盖保存无妨，主预览由 test 1 最终重写
    slide = prs.slides[1]
    texts = "\n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "降级页" in texts and "要点甲" in texts
    assert not [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    print("[PASS] test_unknown_layout_degrades: 未知版式/无效配图 静默降级正确")


if __name__ == "__main__":
    prs = test_full_deck_render()
    test_cover_with_image()
    test_unknown_layout_degrades()
    # 主预览重渲染一次，保证最终落盘的是全覆盖 deck（test 3 曾覆盖保存）
    prs_final = _new_prs()
    render_deck(prs_final, _full_deck(_make_test_image()), style="设计分享")
    prs_final.save(PREVIEW_PATH)
    assert os.path.isfile(PREVIEW_PATH), "预览文件未落盘"
    print(f"\n[OK] test_theme 全绿；预览产物：{PREVIEW_PATH}")
