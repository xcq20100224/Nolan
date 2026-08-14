# -*- coding: utf-8 -*-
"""
Nolan · PPT 图示引擎（ppt_diagrams.py）—— 课件刚需图示：时间轴 / 流程步骤 / 金字塔

职责：在 ppt_layouts 八种版式之外，补三种「图示型」版式（对标 Gamma diagram 核心子集）。
本模块自包含：防御式 import ppt_layouts 的渲染工具与配色常量；ppt_layouts 缺席时
自动切换到内置的最小等价实现，任何情况下模块可导入、函数可调用、绝不抛异常。

公开契约（与 ppt_layouts 版式函数签名完全一致）：
    render_timeline(slide, page, page_no, style)   # page["events"] 2~6 个
    render_process(slide, page, page_no, style)    # page["steps"]  2~5 个
    render_pyramid(slide, page, page_no, style)    # page["levels"] 2~5 层

    page 键约定（三项结构一致）：
      page["events"|"steps"|"levels"] = [{"label": "标题(≤12字)", "desc": "说明(≤40字)"}, ...]

    模块尾部暴露注册表（主控接线用）：
      DIAGRAM_RENDERERS = {"timeline": render_timeline,
                           "process":  render_process,
                           "pyramid":  render_pyramid}

容错契约：缺字段给安全默认；events/steps/levels 缺失或为空时渲染占位节点，
绝不空页、绝不抛异常（单页异常本应由 render_deck 的 fallback 兜底，
但本模块自己也尽量吃掉一切异常）。
"""
from __future__ import annotations

# ================================================================ 防御式导入 ppt_layouts
# 与 ppt_maker/brain 同款风格：能 import 就复用它的渲染工具与配色常量，
# 缺席/失败时降级到本文件尾部的最小等价实现，保证模块任何时刻可用。
try:
    from ppt_layouts import (
        COLOR_DARK, COLOR_INK, COLOR_ACCENT, COLOR_BG, COLOR_WARMGREY,
        COLOR_GOLD, COLOR_TAN, FONT_NAME,
        SLIDE_W, SLIDE_H, MARGIN, CONTENT_W, CONTENT_LEFT, BODY_Y,
        SIZE_BODY_MID, SIZE_BODY_MIN, SIZE_NOTE,
        _rgb, _set_bg, _set_run_font, _set_fill_alpha,
        _add_textbox, _add_rect, _add_round_rect,
        _add_page_header, _safe_str,
    )
    _HAS_LAYOUTS = True
except Exception:                                   # pragma: no cover - 兜底路径
    _HAS_LAYOUTS = False

    # ---- 最小等价实现：常量与工具函数（与 ppt_layouts 同值同步，勿漂移）----
    COLOR_DARK = "3B322C"        # 深棕：标题
    COLOR_INK = "4A4440"         # 墨灰：正文
    COLOR_ACCENT = "C0604A"      # 赭红：强调
    COLOR_BG = "FAF7F2"          # 米白：底色
    COLOR_WARMGREY = "8A8578"    # 暖灰：辅助文字
    COLOR_GOLD = "D9C9A3"        # 浅金：装饰
    COLOR_TAN = "A87B5F"         # 陶棕：暖色系第三色
    FONT_NAME = "微软雅黑"

    SLIDE_W = 13.333             # 16:9 画布宽（英寸）
    SLIDE_H = 7.5                # 16:9 画布高
    MARGIN = 0.8                 # 页边距
    CONTENT_W = SLIDE_W - 2 * MARGIN
    CONTENT_LEFT = MARGIN
    BODY_Y = 1.7                 # 正文区顶

    SIZE_BODY_MID = 16
    SIZE_BODY_MIN = 14
    SIZE_NOTE = 12

    def _rgb(hex_str):
        from pptx.dml.color import RGBColor
        return RGBColor.from_string(hex_str)

    def _set_run_font(run, size_pt, color_hex, bold=False):
        """设置 run 字号/颜色/字体；中文字形显式写 East Asian typeface。"""
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        f = run.font
        f.size = Pt(size_pt)
        f.bold = bold
        f.color.rgb = _rgb(color_hex)
        f.name = FONT_NAME
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", FONT_NAME)

    def _set_fill_alpha(shape, alpha_pct):
        """实心填充加透明度（失败静默，不致命）。"""
        from pptx.oxml.ns import qn
        try:
            spPr = shape._element.spPr
            solid = spPr.find(qn("a:solidFill"))
            if solid is None:
                return
            srgb = solid.find(qn("a:srgbClr"))
            if srgb is None:
                return
            alpha = srgb.makeelement(qn("a:alpha"), {})
            alpha.set("val", str(int(alpha_pct * 1000)))
            srgb.append(alpha)
        except Exception:
            pass

    def _set_bg(slide, color_hex):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(color_hex)

    def _add_textbox(slide, left, top, width, height):
        from pptx.util import Inches
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        box.text_frame.word_wrap = True
        return box.text_frame

    def _add_rect(slide, left, top, width, height, color_hex):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
            Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(color_hex)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _add_round_rect(slide, left, top, width, height, color_hex, adj=0.08):
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
            Inches(width), Inches(height))
        try:
            shp.adjustments[0] = adj
        except Exception:
            pass
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(color_hex)
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _add_page_header(slide, page_title):
        """浅色内容页统一头部：左侧赭红色条 + 26pt 深棕标题。"""
        _add_rect(slide, MARGIN, 0.62, 0.09, 0.62, COLOR_ACCENT)
        tf = _add_textbox(slide, MARGIN + 0.28, 0.5, CONTENT_W - 0.28, 0.9)
        r = tf.paragraphs[0].add_run()
        r.text = page_title or "（无标题）"
        _set_run_font(r, 26, COLOR_DARK, bold=True)

    def _safe_str(v, default=""):
        s = str(v).strip() if v is not None else ""
        return s or default


# ================================================================ 共享：节点列表归一

# 各图示的节点数上限（与 ppt_maker 侧契约保持一致）
_MAX_ITEMS = {"timeline": 6, "process": 5, "pyramid": 5}

# 暖色系节点配色（填充色, 其上文字色）：深块浅字、浅块深字
_NODE_FILLS = [
    (COLOR_ACCENT, "FFFFFF"),    # 赭红 · 白字
    (COLOR_TAN, "FFFFFF"),       # 陶棕 · 白字
    (COLOR_WARMGREY, "FFFFFF"),  # 暖灰 · 白字
    (COLOR_GOLD, COLOR_DARK),    # 浅金 · 深棕字
    ("E3DACA", COLOR_DARK),      # 米灰 · 深棕字
    (COLOR_INK, "FFFFFF"),       # 墨灰 · 白字（第 6 节点兜底）
]

# 金字塔层级配色：顶层最深（最重）逐层变浅
_PYRAMID_FILLS = [
    (COLOR_ACCENT, "FFFFFF"),
    (COLOR_TAN, "FFFFFF"),
    (COLOR_WARMGREY, "FFFFFF"),
    (COLOR_GOLD, COLOR_DARK),
    ("E3DACA", COLOR_DARK),
]


def _norm_items(page, key, layout):
    """把 page[key] 归一成 [(label, desc), ...]：清洗、截断、去空壳。
    空列表/缺字段给占位节点，保证至少 2 个节点，绝不返回空、绝不抛异常。"""
    if not isinstance(page, dict):
        page = {}
    max_items = _MAX_ITEMS.get(layout, 5)
    out = []
    raw = page.get(key)
    if isinstance(raw, (list, tuple)):
        for it in raw[:max_items]:
            if not isinstance(it, dict):
                continue
            # 防御性截断：契约要求 label ≤12 字、desc ≤40 字，这里留余量
            label = _safe_str(it.get("label"))[:16]
            desc = _safe_str(it.get("desc"))[:60]
            if label or desc:
                out.append((label or f"第 {len(out) + 1} 项", desc))
    if not out:
        # 缺字段/空列表的安全降级：占位节点，绝不空页
        out = [("内容筹备中", "本页图示内容筹备中，细节以现场讲解为准。")]
    if len(out) == 1:
        out.append(("待补充", ""))
    return out


def _fill_node(slide, left, top, width, height, color_hex, shape_name="oval"):
    """画一个去边框去阴影的实心节点形状（圆/矩形/圆角矩形）。"""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    mso = {"oval": MSO_SHAPE.OVAL,
           "rect": MSO_SHAPE.RECTANGLE}.get(shape_name, MSO_SHAPE.OVAL)
    shp = slide.shapes.add_shape(
        mso, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _shape_text(shape, lines, size, color_hex, bold_first=True):
    """往形状自带文本框里写居中文字（lines: [(text, bold), ...]），垂直居中。"""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        _set_run_font(r, size, color_hex, bold=(bold and bold_first))


def _center_text(slide, left, top, width, height, text, size, color_hex, bold=False):
    """居中单段文本框（label/desc 用）。"""
    from pptx.enum.text import PP_ALIGN
    tf = _add_textbox(slide, left, top, width, height)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_run_font(r, size, color_hex, bold=bold)


def _render_minimal_fallback(slide, page):
    """图示渲染内部异常时的本模块自保：底色 + 页眉 + 一行占位，绝不空页。"""
    try:
        _set_bg(slide, COLOR_BG)
        title = _safe_str(page.get("page_title") if isinstance(page, dict) else "",
                          "（无标题）")
        _add_page_header(slide, title)
        tf = _add_textbox(slide, CONTENT_LEFT, BODY_Y, CONTENT_W, 1.0)
        r = tf.paragraphs[0].add_run()
        r.text = "本页图示内容筹备中，细节以现场讲解为准。"
        _set_run_font(r, SIZE_BODY_MID, COLOR_INK)
    except Exception:
        pass   # 自保也失败就交给 render_deck 的 fallback，这里不再抛


# ================================================================ 图示 1：timeline 时间轴

def _render_timeline_impl(slide, page, page_no, style):
    """横向中轴 + 圆点节点（内置序号徽章），label 在轴上方、desc 在节点下方。"""
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "（无标题）"))

    events = _norm_items(page, "events", "timeline")
    n = len(events)
    slot_w = CONTENT_W / n                       # 每个节点占的横向槽宽
    axis_y = 3.55                                # 中轴高度（正文区纵向居中）

    # 中轴：浅金细线横贯内容区
    _add_rect(slide, CONTENT_LEFT, axis_y, CONTENT_W, 0.03, COLOR_GOLD)

    for i, (label, desc) in enumerate(events):
        cx = CONTENT_LEFT + slot_w * (i + 0.5)   # 节点圆心 x
        fill, on_fill = _NODE_FILLS[i % len(_NODE_FILLS)]

        # 节点圆点 + 序号徽章（数字直接写进圆里，白字/深字随底色）
        d = 0.4
        node = _fill_node(slide, cx - d / 2, axis_y - d / 2 + 0.015, d, d, fill)
        _shape_text(node, [(str(i + 1), True)], 15, on_fill)

        # label：轴上方，加粗深棕，按槽宽居中
        _center_text(slide, cx - slot_w / 2 + 0.08, axis_y - 1.05,
                     slot_w - 0.16, 0.7, label, SIZE_BODY_MID, COLOR_DARK, bold=True)
        # desc：节点下方，小字暖灰
        if desc:
            _center_text(slide, cx - slot_w / 2 + 0.08, axis_y + 0.42,
                         slot_w - 0.16, 2.2, desc, SIZE_NOTE, COLOR_WARMGREY)


def render_timeline(slide, page, page_no, style):
    """时间轴图示页（公开入口，永不抛异常）。"""
    try:
        _render_timeline_impl(slide, page if isinstance(page, dict) else {},
                              page_no, style)
    except Exception:
        _render_minimal_fallback(slide, page)


# ================================================================ 图示 2：process 流程步骤

def _render_process_impl(slide, page, page_no, style):
    """横向 chevron 块链：首块为平头五边形、后续为箭头形，块间交叠出「接续感」。
    块内序号 + label（块色轮转暖色系），块下放 desc 小字。"""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "（无标题）"))

    steps = _norm_items(page, "steps", "process")
    n = len(steps)
    overlap = 0.28                               # 块间交叠量：chevron 咬住前一块
    block_w = (CONTENT_W + (n - 1) * overlap) / n
    block_h = 1.5
    y0 = 2.3

    for i, (label, desc) in enumerate(steps):
        x = CONTENT_LEFT + i * (block_w - overlap)
        fill, on_fill = _NODE_FILLS[i % len(_NODE_FILLS)]
        # 首块平头（homePlate 五边形），后续 chevron 箭头块
        mso = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        shp = slide.shapes.add_shape(
            mso, Inches(x), Inches(y0), Inches(block_w), Inches(block_h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
        shp.line.fill.background()
        shp.shadow.inherit = False
        try:
            tf = shp.text_frame
            tf.margin_right = Inches(0.35)   # 让开箭头尖
            tf.margin_left = Inches(0.15 if i == 0 else 0.35)             # 让开前块咬口
        except Exception:
            pass
        _shape_text(shp, [(f"{i + 1} · {label}", True)], SIZE_BODY_MID, on_fill)

        # desc：块下方小字，与该块水平对齐居中
        if desc:
            _center_text(slide, x + 0.15, y0 + block_h + 0.22,
                         block_w - 0.3, 1.9, desc, SIZE_NOTE, COLOR_INK)


def render_process(slide, page, page_no, style):
    """流程步骤图示页（公开入口，永不抛异常）。"""
    try:
        _render_process_impl(slide, page if isinstance(page, dict) else {},
                             page_no, style)
    except Exception:
        _render_minimal_fallback(slide, page)


# ================================================================ 图示 3：pyramid 金字塔

def _render_pyramid_impl(slide, page, page_no, style):
    """自顶向下宽度递增的圆角矩形堆叠（金字塔意形）。
    层内居中加粗 label（顶层色最深）；右侧对齐每层放 desc 小字 + 同层色小方块引子。"""
    _set_bg(slide, COLOR_BG)
    _add_page_header(slide, _safe_str(page.get("page_title"), "（无标题）"))

    levels = _norm_items(page, "levels", "pyramid")
    n = len(levels)

    # 金字塔区：靠左布置，右侧留 4 英寸宽的 desc 栏
    w_min, w_max = 2.6, 7.4                      # 顶层/底层宽度
    cx = CONTENT_LEFT + w_max / 2                # 金字塔中轴 x（底层左缘对齐页边距）
    desc_x = CONTENT_LEFT + w_max + 0.5          # desc 栏左缘
    desc_w = SLIDE_W - MARGIN - desc_x           # desc 栏宽（≈3.9 英寸）

    top_y = BODY_Y + 0.25
    gap = 0.09
    layer_h = (6.55 - top_y - (n - 1) * gap) / n

    for i, (label, desc) in enumerate(levels):
        # 自顶向下宽度线性递增
        w = w_min + (w_max - w_min) * (i / (n - 1) if n > 1 else 1.0)
        x = cx - w / 2
        y = top_y + i * (layer_h + gap)
        fill, on_fill = _PYRAMID_FILLS[i % len(_PYRAMID_FILLS)]

        layer = _add_round_rect(slide, x, y, w, layer_h, fill, adj=0.10)
        _shape_text(layer, [(label, True)], SIZE_BODY_MID, on_fill)

        # 右侧 desc：同层色小方块引子 + 小字说明，垂直对齐层中心
        if desc:
            cy = y + layer_h / 2
            _add_rect(slide, desc_x, cy - 0.07, 0.14, 0.14, fill)
            tf = _add_textbox(slide, desc_x + 0.3, cy - 0.55,
                              desc_w - 0.3, 1.1)
            r = tf.paragraphs[0].add_run()
            r.text = desc
            _set_run_font(r, SIZE_NOTE, COLOR_INK)


def render_pyramid(slide, page, page_no, style):
    """金字塔层级图示页（公开入口，永不抛异常）。"""
    try:
        _render_pyramid_impl(slide, page if isinstance(page, dict) else {},
                             page_no, style)
    except Exception:
        _render_minimal_fallback(slide, page)


# ================================================================ 注册表（主控接线用）

DIAGRAM_RENDERERS = {
    "timeline": render_timeline,
    "process": render_process,
    "pyramid": render_pyramid,
}
