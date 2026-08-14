# -*- coding: utf-8 -*-
"""
ppt_maker 两阶段 pipeline 单元测试：纯 mock，LLM 0 真实调用。
运行：python -m unittest test_ppt_maker -v   （在 jarvis/ 目录下）
落盘目录重定向到临时目录，不污染真实 files/。

mock 脚本约定：script 列表按调用顺序出队；元素是 str 则返回，是 Exception 则抛出。
调用顺序 = 1 次大纲 + 每页 1~3 次精写（重写才会追加；toc/section 页不耗调用）。

版式相关测试约定：
- setUp 强制 ppt_maker._render_deck = None（默认走内置兜底排版，与路 A 是否就位无关）；
- 需要验证 render_deck 接线的测试自行注入 _MockRender，tearDown 统一还原。
"""
import json
import re
import shutil
import sys
import tempfile
import unittest
import urllib.request
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

import ppt_maker
from pptx import Presentation


# ---------------------------------------------------------------- mock 素材

def _outline_json(n=4, title="人工智能简介"):
    """合法大纲：n 页，每页 page_title/core_point/keywords 齐全（无 layout -> bullets）。"""
    return json.dumps({
        "title": title,
        "pages": [
            {
                "page_title": f"章节{i}",
                "core_point": f"第{i}页核心论点：该环节的关键机制决定了整体效果，必须讲清原理与数据。",
                "keywords": [f"关键词{i}甲", f"关键词{i}乙"],
            }
            for i in range(1, n + 1)
        ],
    }, ensure_ascii=False)


def _long_bullet(tag):
    """60 字达标要点（30-60 字区间）。"""
    return (f"{tag}通过真实场景验证：头部试点项目把处理周期从三天压缩到四小时，"
            f"错误率下降约六成，验证了该机制在量产环境中的可行性与收益。")


def _page_json(tag="要点", n_bullets=4, bullet_len="long", note=True):
    """合法精写页：默认 4 条 60 字要点（总 240 字，过 220 严线）。"""
    if bullet_len == "long":
        bullets = [_long_bullet(f"{tag}{j}") for j in range(1, n_bullets + 1)]
    else:
        bullets = [f"{tag}{j}短句。" for j in range(1, n_bullets + 1)]
    data = {"bullets": bullets}
    if note:
        data["speaker_note"] = (
            "这一页先抛出听众最关心的问题，再逐条拆解机制与数据，"
            "每条讲完停顿半拍确认大家跟上，最后一句小结自然过渡到下一页。" * 2)
    return json.dumps(data, ensure_ascii=False)


_NOTE = ("这一页先抛出听众最关心的问题，再逐条拆解机制与数据，"
         "每条讲完停顿半拍确认大家跟上，最后一句小结自然过渡到下一页。" * 2)


def _tc_points(tag, n):
    """two_column 单栏要点：每条约 45 字，6 条约 270 字，过 220 严线。"""
    return [f"{tag}{j}：该维度实测数据显示头部玩家投入强度连续三年翻倍，单位成本下降约四成。"
            for j in range(1, n + 1)]


def _two_column_json(ln=3, rn=3, lh="中国", rh="海外"):
    """合法 two_column 精写页：默认左右各 3 点。"""
    return json.dumps({
        "left": {"heading": lh, "points": _tc_points("左", ln)},
        "right": {"heading": rh, "points": _tc_points("右", rn)},
        "speaker_note": _NOTE,
    }, ensure_ascii=False)


def _big_number_json(stats=None):
    """合法 big_number 精写页：默认 2 个数字。"""
    if stats is None:
        stats = [{"number": "42.7%", "caption": "2024 年国内市场渗透率，来源：行业白皮书"},
                 {"number": "1280 亿元", "caption": "2024 年市场规模，来源：头部券商研报"}]
    return json.dumps({"stats": stats, "speaker_note": _NOTE}, ensure_ascii=False)


def _chart_json(values=None, n_bullets=2):
    """合法 chart 精写页：默认 bar，3 类别 1 系列，2 条 60 字解读。"""
    if values is None:
        values = [100, 200, 300]
    return json.dumps({
        "chart": {"type": "bar", "title": "年度出货量",
                  "categories": ["2022", "2023", "2024"],
                  "series": [{"name": "出货量(万台)", "values": values}]},
        "bullets": [_long_bullet(f"解读{j}") for j in range(1, n_bullets + 1)],
        "speaker_note": _NOTE,
    }, ensure_ascii=False)


def _quote_json(text=None):
    """合法 quote 精写页：默认 30+ 字金句。"""
    if text is None:
        text = "技术本身没有方向，方向是人给的；每一次范式转移，奖励的都是先想清楚的人。"
    return json.dumps({"quote": text, "attribution": "—— 某行业观察者",
                       "speaker_note": _NOTE}, ensure_ascii=False)


def _layout_outline_json():
    """八种版式混合大纲：toc + bullets + two_column + big_number + chart + quote
    + 未知版式（应降级 bullets）+ closing。"""
    return json.dumps({
        "title": "新能源汽车产业分析",
        "subtitle": "从政策到市场的全景扫描",
        "pages": [
            {"layout": "toc", "page_title": "目录",
             "core_point": "全篇结构导览", "keywords": [],
             "entries": ["市场格局", "中外对比", "关键数字", "销量走势", "金句点题", "奇怪页", "总结行动"]},
            {"layout": "bullets", "page_title": "市场格局",
             "core_point": "市场集中度快速提升，头部三家拿走六成份额，腰部玩家加速出清。",
             "keywords": ["CR3", "出清"]},
            {"layout": "two_column", "page_title": "中外对比",
             "core_point": "中国靠供应链密度取胜，海外靠品牌溢价守成，路径分野已经清晰。",
             "keywords": ["供应链", "品牌溢价"],
             "left": {"heading": "中国"}, "right": {"heading": "海外"}},
            {"layout": "big_number", "page_title": "关键数字",
             "core_point": "三个数字足以看清这个行业：渗透率、市场规模、增速。",
             "keywords": ["渗透率"],
             "stats": [{"number": "42.7%", "caption": "渗透率草稿"}]},
            {"layout": "chart", "page_title": "销量走势",
             "core_point": "销量三年翻两番，增长斜率仍在变陡，未见平台期迹象。",
             "keywords": ["出货量"],
             "chart": {"type": "bar", "title": "销量", "categories": ["2022", "2023", "2024"],
                       "series": [{"name": "万辆", "values": [100, 200, 300]}]}},
            {"layout": "quote", "page_title": "金句点题",
             "core_point": "用一句行业金句收住论证，点明范式转移的本质。",
             "keywords": []},
            {"layout": "weird_layout", "page_title": "奇怪页",
             "core_point": "这一页版式名是乱写的，应被降级为普通要点页并正常精写。",
             "keywords": ["降级"]},
            {"layout": "closing", "page_title": "总结行动",
             "core_point": "总结三个判断并给出两类玩家各自的行动建议。",
             "keywords": ["行动建议"]},
        ],
    }, ensure_ascii=False)


class _MockLLM:
    """记录调用与 prompt，按脚本返回；脚本元素为 Exception 时抛出。"""

    def __init__(self, script):
        self.script = list(script)
        self.calls = []

    def __call__(self, prompt):
        self.calls.append(prompt)
        if not self.script:
            return None
        item = self.script.pop(0)
        if isinstance(item, Exception):
            raise item
        return item


class _MockRender:
    """mock 路 A 的 render_deck(prs, deck, style)：只记录入参，不造 slide。"""

    def __init__(self):
        self.calls = []

    def __call__(self, prs, deck, style):
        self.calls.append((prs, deck, style))


def _script_for(n_pages, page_json=None):
    """1 次大纲 + 每页 1 次精写 的标准脚本。"""
    return [_outline_json(n_pages)] + [page_json or _page_json(f"P{i}") for i in range(1, n_pages + 1)]


# ---------------------------------------------------------------- 路 B 生图 mock 素材

def _image_outline_json(n=4, layouts=None, title="人工智能简介"):
    """带配图 prompt 的大纲：cover_image_prompt + 每页 image_prompt。
    layouts 为 None 时全部 bullets（注意：末页会被归一化升级为 closing）。"""
    pages = []
    for i in range(1, n + 1):
        layout = (layouts[i - 1] if layouts else "bullets")
        item = {
            "layout": layout,
            "page_title": f"章节{i}",
            "core_point": f"第{i}页核心论点：该环节的关键机制决定了整体效果，必须讲清原理与数据。",
            "keywords": [f"关键词{i}甲", f"关键词{i}乙"],
            "image_prompt": f"warm scene of concept {i}, clay red accents",
        }
        if layout == "two_column":
            item["left"] = {"heading": "左栏"}
            item["right"] = {"heading": "右栏"}
        pages.append(item)
    return json.dumps({
        "title": title,
        "subtitle": "带配图的测试大纲",
        "cover_image_prompt": "vast abstract landscape, warm sunrise over geometric hills",
        "pages": pages,
    }, ensure_ascii=False)


class _FakeResp:
    """urlopen 打桩返回值：with 语义 + read() + headers。"""

    def __init__(self, body=b"", ctype="application/json"):
        self._body = body
        self.headers = {"Content-Type": ctype}

    def read(self):
        return self._body

    def __enter__(self):
        return self

    def __exit__(self, *a):
        return False


class _MockHTTP:
    """urllib.request.urlopen 打桩：按 URL 路由生图 API 与图片下载。
    api_fails / dl_fails：命中序号（1 起）的调用抛异常；
    dl_ctypes：第 n 次下载的 Content-Type（默认 image/png）。"""

    def __init__(self, api_fails=(), dl_fails=(), dl_ctypes=None):
        self.calls = []            # (url, api_body_bytes_or_None)
        self.api_fails = set(api_fails)
        self.dl_fails = set(dl_fails)
        self.dl_ctypes = dict(dl_ctypes or {})
        self._api_n = 0
        self._dl_n = 0

    def __call__(self, req, timeout=None):
        url = getattr(req, "full_url", str(req))
        if "/images/generations" in url:
            self._api_n += 1
            self.calls.append((url, getattr(req, "data", None)))
            if self._api_n in self.api_fails:
                raise RuntimeError(f"生图 API 第 {self._api_n} 次调用炸了")
            return _FakeResp(json.dumps(
                {"data": [{"url": f"https://img.mock/pic{self._api_n}.png"}]}
            ).encode("utf-8"))
        self._dl_n += 1
        self.calls.append((url, None))
        if self._dl_n in self.dl_fails:
            raise RuntimeError(f"第 {self._dl_n} 次下载超时")
        ctype = self.dl_ctypes.get(self._dl_n, "image/png")
        return _FakeResp(b"\x89PNG\r\n\x1a\n fake-image-bytes", ctype)

    def api_bodies(self):
        """全部生图 API 调用的请求体（解码成 dict）。"""
        return [json.loads(b.decode("utf-8")) for _u, b in self.calls if b]


def _body_text(prs, slide_idx):
    """抽某页正文文本框的全部文字。"""
    return "\n".join(
        sh.text_frame.text for sh in prs.slides[slide_idx].shapes if sh.has_text_frame)


# ---------------------------------------------------------------- 测试

class PptMakerTest(unittest.TestCase):
    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp(prefix="ppt_maker_test_"))
        self._orig_dir = ppt_maker.FILES_DIR
        ppt_maker.FILES_DIR = self.tmp
        # 默认强制走内置兜底排版：既有 15 个测试的行为与路 A 是否就位解耦
        self._orig_render = ppt_maker._render_deck
        ppt_maker._render_deck = None
        # 默认断掉生图链（路 B）：配置视为缺失 -> 静默跳过，既有测试零 HTTP。
        # 生图测试用 _enable_images 自行打桩假配置 + urlopen 路由。
        self._orig_img_cfg = ppt_maker._load_image_config
        ppt_maker._load_image_config = lambda: None
        # 默认断掉联网研究链（R1）：模块视为缺席 -> 研究材料空串，既有测试零 HTTP。
        # 研究注入测试用 _enable_research 自行打桩。
        self._orig_research = ppt_maker._research_topic
        ppt_maker._research_topic = None

    def tearDown(self):
        ppt_maker.FILES_DIR = self._orig_dir
        ppt_maker._render_deck = self._orig_render
        ppt_maker._load_image_config = self._orig_img_cfg
        ppt_maker._research_topic = self._orig_research
        shutil.rmtree(self.tmp, ignore_errors=True)

    # ---- 路 B 打桩工具 ----
    def _enable_images(self, http):
        """打桩生图链：假配置（base,key）+ urlopen 路由 mock。addCleanup 自动还原。"""
        orig_cfg = ppt_maker._load_image_config
        orig_open = urllib.request.urlopen
        ppt_maker._load_image_config = lambda: ("https://api.mock", "mock-key")
        urllib.request.urlopen = http

        def _restore():
            ppt_maker._load_image_config = orig_cfg
            urllib.request.urlopen = orig_open
        self.addCleanup(_restore)

    # 1. 标准两阶段 -> 真 pptx：页数、标题、notes 非空、调用次数 = 1 + N
    def test_standard_pipeline_builds_real_pptx(self):
        mock = _MockLLM(_script_for(4))
        r = ppt_maker.make_ppt("人工智能简介", pages=4, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(r["pages"], 5)               # 4 内容页 + 封面
        self.assertEqual(r["title"], "人工智能简介")
        self.assertTrue(r["file_name"].endswith(".pptx"))
        path = Path(r["path"])
        self.assertTrue(path.exists())
        self.assertEqual(path.parent, self.tmp)       # 落在（重定向后的）files 目录

        prs = Presentation(str(path))
        self.assertEqual(len(list(prs.slides)), 5)
        cover_text = _body_text(prs, 0)
        self.assertIn("人工智能简介", cover_text)
        for i, slide in enumerate(prs.slides):
            self.assertTrue(slide.has_notes_slide, f"slide {i} 缺 notes_slide")
            self.assertTrue(slide.notes_slide.notes_text_frame.text.strip(),
                            f"slide {i} notes 为空")
        self.assertAlmostEqual(prs.slide_width / prs.slide_height, 16 / 9, places=2)
        self.assertEqual(len(mock.calls), 5)          # 1 大纲 + 4 精写
        # 运行统计：每页 4 要点、240 字、0 重写
        self.assertEqual(len(ppt_maker.last_run["page_stats"]), 4)
        for st in ppt_maker.last_run["page_stats"]:
            self.assertEqual(st["bullets"], 4)
            self.assertGreaterEqual(st["chars"], 220)
            self.assertEqual(st["rewrites"], 0)
            self.assertFalse(st["fallback"])

    # 2. 大纲 JSON 带围栏/废话 -> 仍能解析
    def test_fenced_outline_parses(self):
        fenced = "好的，这是大纲：\n```json\n" + _outline_json(3) + "\n```\n希望对你有帮助。"
        mock = _MockLLM([fenced] + [_page_json(f"P{i}") for i in range(1, 4)])
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(r["pages"], 4)

    # 3. 大纲阶段垃圾 -> 修复重试 1 次仍垃圾 -> ok=False，不造文件
    def test_outline_garbage_fails_clean(self):
        mock = _MockLLM(["这不是JSON，随便聊聊吧", "还是垃圾 {{{"])
        r = ppt_maker.make_ppt("人工智能简介", llm_caller=mock)
        self.assertFalse(r["ok"])
        self.assertIn("JSON", r["error"])
        self.assertEqual(len(mock.calls), 2)          # 恰好修复重试一次
        self.assertEqual(list(self.tmp.glob("*.pptx")), [])

    # 3b. 大纲阶段 LLM 无响应（None）-> ok=False，不造空壳
    def test_none_outline_fails_clean(self):
        r = ppt_maker.make_ppt("人工智能简介", llm_caller=_MockLLM([None, None]))
        self.assertFalse(r["ok"])
        self.assertEqual(list(self.tmp.glob("*.pptx")), [])

    # 3c. 大纲首解失败、修复重试成功 -> ok
    def test_outline_retry_succeeds(self):
        mock = _MockLLM(["垃圾输出", _outline_json(3)]
                        + [_page_json(f"P{i}") for i in range(1, 4)])
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(len(mock.calls), 2 + 3)      # 大纲 2 次 + 3 页各 1 次

    # 4. 某页精写调用异常 -> 兜底页顶上，页数完整，整单不失败
    def test_page_exception_uses_fallback(self):
        script = [_outline_json(4),
                  _page_json("P1"),
                  RuntimeError("网络抖动"),           # 第 2 页精写炸掉
                  _page_json("P3"),
                  _page_json("P4")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("人工智能简介", pages=4, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(r["pages"], 5)               # 页数不缺
        prs = Presentation(str(path := Path(r["path"])))
        self.assertEqual(len(list(prs.slides)), 5)
        # 兜底页（物理第 3 页 = 封面 + 内容第 2 页）含大纲 core_point 内容
        body = _body_text(prs, 2)
        self.assertIn("核心论点", body)
        # 兜底页 notes 也非空
        self.assertTrue(prs.slides[2].notes_slide.notes_text_frame.text.strip())
        st = ppt_maker.last_run["page_stats"][1]
        self.assertTrue(st["fallback"])

    # 4b. 某页精写连吐垃圾（非 JSON）-> 3 次尝试耗尽后兜底页顶上
    def test_page_garbage_falls_back_after_attempts(self):
        script = [_outline_json(3),
                  "垃圾一", "垃圾二", "垃圾三",       # 第 1 页 3 次全废
                  _page_json("P2"),
                  _page_json("P3")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(r["pages"], 4)
        self.assertTrue(ppt_maker.last_run["page_stats"][0]["fallback"])
        self.assertEqual(len(mock.calls), 1 + 3 + 2)  # 大纲 + 第1页3次 + 后两页

    # 5. 质量闸触发重写：第一次 2 要点共 30 字，第二次给足 -> 用达标版
    def test_quality_gate_triggers_rewrite(self):
        thin = _page_json("薄", n_bullets=2, bullet_len="short")  # 2 条共约 12 字
        good = _page_json("厚")
        script = [_outline_json(3), thin, good, _page_json("P2"), _page_json("P3")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        # 重写真实发生：第 1 页调了 2 次（总 1 + 2 + 1 + 1 = 5 次）
        self.assertEqual(len(mock.calls), 5)
        self.assertEqual(ppt_maker.last_run["page_stats"][0]["rewrites"], 1)
        # 最终用的是达标版（厚版标记在物理第 2 页正文里，薄版不在）
        prs = Presentation(r["path"])
        body = _body_text(prs, 1)
        self.assertIn("厚1通过真实场景验证", body)
        self.assertNotIn("薄1", body)

    # 6. 重写 2 次仍不达标 -> 取历次最好（字数最多者），页仍生成
    def test_best_of_failed_rewrites(self):
        # 三版都只有 3 条（过不了要点数闸），但字数递增：第三版最好
        v1 = _page_json("甲", n_bullets=3, bullet_len="short")
        v2 = json.dumps({"bullets": [f"乙{j}这是一条中等长度的要点，带一些解释但不够充分。"
                                     for j in range(1, 4)],
                         "speaker_note": "备注。" * 30}, ensure_ascii=False)
        v3 = _page_json("丙", n_bullets=3, bullet_len="long")     # 3×60=180 字，仍差
        script = [_outline_json(3), v1, v2, v3, _page_json("P2"), _page_json("P3")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(len(mock.calls), 1 + 3 + 2)  # 第 1 页打满 3 次
        st = ppt_maker.last_run["page_stats"][0]
        self.assertEqual(st["rewrites"], 2)
        self.assertFalse(st["fallback"])
        # 最终采用字数最多的 v3（丙版）
        prs = Presentation(r["path"])
        body = _body_text(prs, 1)
        self.assertIn("丙1通过真实场景验证", body)
        self.assertNotIn("甲1", body)
        self.assertNotIn("乙1", body)

    # 7. 页数钳制 3-20（大纲 prompt 中的请求页数为证）
    def test_pages_clamped(self):
        mock_lo = _MockLLM(_script_for(3))
        r = ppt_maker.make_ppt("t", pages=1, llm_caller=mock_lo)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertIn("3 页", mock_lo.calls[0])       # 1 -> 钳到 3

        mock_hi = _MockLLM(_script_for(20))
        r = ppt_maker.make_ppt("t", pages=99, llm_caller=mock_hi)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertIn("20 页", mock_hi.calls[0])      # 99 -> 钳到 20
        self.assertEqual(r["pages"], 21)              # 20 内容页 + 封面

    # 8. 文件名安全性：非法字符被清洗，Windows 合法
    def test_filename_sanitized(self):
        nasty = '人工智能/简介:<>*?"| 版'
        r = ppt_maker.make_ppt(nasty, pages=3,
                               llm_caller=_MockLLM(_script_for(3)))
        self.assertTrue(r["ok"], r.get("error"))
        name = r["file_name"]
        self.assertNotRegex(name, r'[\\/:*?"<>|]')
        self.assertRegex(name, r"^[\w一-鿿-]+_\d{8}-\d{4}(-\d+)?\.pptx$")
        self.assertIn("人工智能", name)
        self.assertIn("简介", name)

    # 8b. 空主题 -> ok=False
    def test_empty_topic_fails(self):
        r = ppt_maker.make_ppt("   ", llm_caller=_MockLLM([]))
        self.assertFalse(r["ok"])

    # 8c. 大纲缺 pages 结构 -> ok=False
    def test_missing_outline_pages_fails(self):
        bad = json.dumps({"title": "只有标题没有页"}, ensure_ascii=False)
        r = ppt_maker.make_ppt("t", llm_caller=_MockLLM([bad, bad]))
        self.assertFalse(r["ok"])
        self.assertEqual(list(self.tmp.glob("*.pptx")), [])

    # 9. 精写页缺 speaker_note -> 兜底合成，notes 仍非空
    def test_missing_note_synthesized(self):
        script = [_outline_json(3)] + [_page_json(f"P{i}", note=False) for i in range(1, 4)]
        r = ppt_maker.make_ppt("t", pages=3, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        prs = Presentation(r["path"])
        for slide in prs.slides:
            self.assertTrue(slide.notes_slide.notes_text_frame.text.strip())

    # 10. 高密度内容触发字号自适应（不溢出：三档 18/16/14）
    def test_dense_content_shrinks_font(self):
        from pptx.util import Pt
        dense = _page_json("密", n_bullets=6, bullet_len="long")   # 6×60=360 字
        script = [_outline_json(3), dense, _page_json("P2"), _page_json("P3")]
        r = ppt_maker.make_ppt("t", pages=3, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        prs = Presentation(r["path"])
        # 第 1 内容页（物理 idx=1）：360 字应缩到 14pt 档
        sizes = set()
        for sh in prs.slides[1].shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for run in p.runs:
                        if run.text.startswith("•"):
                            sizes.add(run.font.size)
        self.assertEqual(sizes, {Pt(14)})
        # 第 2 内容页 240 字保持 16pt 档
        sizes2 = set()
        for sh in prs.slides[2].shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for run in p.runs:
                        if run.text.startswith("•"):
                            sizes2.add(run.font.size)
        self.assertEqual(sizes2, {Pt(16)})

    # ================================================================ 新增：版式管线

    # 11. 各版式大纲 -> 归一化 deck 符合路 A 契约（render_deck 捕获断言）
    def test_layouts_normalized_into_contract_deck(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        # big_number 精写给 5 个数字（应钳到 3）；chart 精写 values 给字符串数字且多 1 个（应转 float 并截齐）
        five_stats = _big_number_json(
            [{"number": f"{i}0%", "caption": f"第{i}个数字的口径与来源说明"} for i in range(1, 6)])
        str_chart = _chart_json(values=["12.5", "8.0", "20", "99"])   # 4 个值 vs 3 类别
        script = [_layout_outline_json(),
                  _page_json("格局"),        # bullets 页
                  _two_column_json(),        # two_column 页
                  five_stats,                # big_number 页
                  str_chart,                 # chart 页
                  _quote_json(),             # quote 页
                  _page_json("奇怪"),        # 未知版式降级页
                  _page_json("收尾")]        # closing 页
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("新能源汽车产业分析", pages=8, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(r["pages"], 9)                  # 8 大纲页（含 toc）+ 封面
        # toc 不耗调用：1 大纲 + 7 精写
        self.assertEqual(len(mock.calls), 8)

        self.assertEqual(len(render.calls), 1)           # render_deck 恰好被调一次
        prs, deck, style = render.calls[0]
        self.assertAlmostEqual(prs.slide_width / prs.slide_height, 16 / 9, places=2)
        self.assertEqual(style, "工作汇报")
        # deck 顶层契约
        self.assertEqual(deck["title"], "新能源汽车产业分析")
        self.assertEqual(deck["subtitle"], "从政策到市场的全景扫描")
        pages = deck["pages"]
        self.assertEqual(len(pages), 8)
        for pg in pages:
            self.assertIn("layout", pg)
            self.assertTrue(pg["speaker_note"].strip(), f"{pg['layout']} 页 speaker_note 为空")

        # toc：entries 与大纲一致
        self.assertEqual(pages[0]["layout"], "toc")
        self.assertEqual(pages[0]["entries"],
                         ["市场格局", "中外对比", "关键数字", "销量走势", "金句点题", "奇怪页", "总结行动"])
        # two_column：左右栏 heading + 各 3 点
        tc = pages[2]
        self.assertEqual(tc["layout"], "two_column")
        self.assertEqual(tc["left"]["heading"], "中国")
        self.assertEqual(tc["right"]["heading"], "海外")
        self.assertEqual(len(tc["left"]["points"]), 3)
        self.assertEqual(len(tc["right"]["points"]), 3)
        # big_number：5 个数字被钳到 3，number 全非空
        bn = pages[3]
        self.assertEqual(bn["layout"], "big_number")
        self.assertEqual(len(bn["stats"]), 3)
        for s in bn["stats"]:
            self.assertTrue(s["number"])
            self.assertTrue(s["caption"])
        # chart：字符串数字转 float、多余值截齐到 categories 等长
        ch = pages[4]
        self.assertEqual(ch["layout"], "chart")
        chart = ch["chart"]
        self.assertEqual(chart["type"], "bar")
        self.assertEqual(len(chart["series"]), 1)
        self.assertEqual(chart["series"][0]["values"], [12.5, 8.0, 20.0])
        for v in chart["series"][0]["values"]:
            self.assertIsInstance(v, float)
        self.assertEqual(len(chart["series"][0]["values"]), len(chart["categories"]))
        self.assertEqual(len(ch["bullets"]), 2)          # 2 条解读
        # quote：契约字段
        q = pages[5]
        self.assertEqual(q["layout"], "quote")
        self.assertGreaterEqual(len(q["quote"]), 20)
        self.assertTrue(q["attribution"])
        # 未知版式降级 bullets
        weird = pages[6]
        self.assertEqual(weird["layout"], "bullets")
        self.assertEqual(len(weird["bullets"]), 4)
        # 收尾页 closing
        self.assertEqual(pages[7]["layout"], "closing")

    # 12. ppt_layouts 缺席（_render_deck=None）-> 回退内置排版，版式页摊平渲染仍出文件
    def test_missing_ppt_layouts_falls_back_to_legacy(self):
        self.assertIsNone(ppt_maker._render_deck)        # setUp 已强制 None
        script = [_layout_outline_json(),
                  _page_json("格局"),
                  _two_column_json(),
                  _big_number_json(),
                  _chart_json(),
                  _quote_json(),
                  _page_json("奇怪"),
                  _page_json("收尾")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("新能源汽车产业分析", pages=8, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        path = Path(r["path"])
        self.assertTrue(path.exists())
        prs = Presentation(str(path))
        self.assertEqual(len(list(prs.slides)), 9)       # 封面 + 8 页
        # two_column 页（物理 idx=3）摊平后含栏题与要点
        body = _body_text(prs, 3)
        self.assertIn("【中国】", body)
        self.assertIn("【海外】", body)
        self.assertIn("左1", body)
        # big_number 页（物理 idx=4）含数字
        self.assertIn("42.7%", _body_text(prs, 4))
        # 每页 notes 非空
        for i, slide in enumerate(prs.slides):
            self.assertTrue(slide.notes_slide.notes_text_frame.text.strip(),
                            f"slide {i} notes 为空")

    # 13. two_column 质量闸：左右各 2 点不达标 -> 重写后各 3 点过闸
    def test_two_column_gate_triggers_rewrite(self):
        outline = json.dumps({
            "title": "中外模式对比", "pages": [
                {"layout": "two_column", "page_title": "中外对比",
                 "core_point": "中国靠供应链密度取胜，海外靠品牌溢价守成，分野清晰。",
                 "keywords": ["供应链"], "left": {"heading": "中国"}, "right": {"heading": "海外"}},
                {"layout": "bullets", "page_title": "总结",
                 "core_point": "两种路径各有死穴，胜负手在成本控制与品牌建设的组合拳。",
                 "keywords": ["组合拳"]},
            ]}, ensure_ascii=False)
        thin = _two_column_json(ln=2, rn=2)              # 左右各 2 点 -> 不过闸
        good = _two_column_json(ln=3, rn=3)
        script = [outline, thin, good, _page_json("P2")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("中外模式对比", pages=2, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(ppt_maker.last_run["page_stats"][0]["rewrites"], 1)
        self.assertEqual(ppt_maker.last_run["page_stats"][0]["layout"], "two_column")
        self.assertEqual(len(mock.calls), 1 + 2 + 1)     # 大纲 + 第1页2次 + 第2页

    # 14. big_number 质量闸：number 为空 -> 重写后补齐数字过闸
    def test_big_number_gate_triggers_rewrite(self):
        outline = json.dumps({
            "title": "关键数字", "pages": [
                {"layout": "big_number", "page_title": "关键数字",
                 "core_point": "三个数字看清行业：渗透率、市场规模、增速。",
                 "keywords": ["渗透率"], "stats": []},
                {"layout": "bullets", "page_title": "总结",
                 "core_point": "数字背后的结论是行业仍在高速扩张期，窗口未关。",
                 "keywords": ["窗口"]},
            ]}, ensure_ascii=False)
        empty_num = _big_number_json([{"number": "", "caption": "口径说明但数字为空"}])
        good = _big_number_json()
        script = [outline, empty_num, good, _page_json("P2")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("关键数字", pages=2, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        st = ppt_maker.last_run["page_stats"][0]
        self.assertEqual(st["rewrites"], 1)
        self.assertEqual(st["layout"], "big_number")
        self.assertEqual(len(mock.calls), 1 + 2 + 1)

    # 15. quote 质量闸：金句 <20 字 -> 重写后过闸
    def test_quote_gate_triggers_rewrite(self):
        outline = json.dumps({
            "title": "金句收束", "pages": [
                {"layout": "quote", "page_title": "金句点题",
                 "core_point": "用一句行业金句收住全篇论证，点明范式转移的本质。",
                 "keywords": []},
                {"layout": "bullets", "page_title": "总结",
                 "core_point": "收束全篇给出行动建议，落点要具体可执行。",
                 "keywords": ["行动"]},
            ]}, ensure_ascii=False)
        short = _quote_json("太短了。")                   # 4 字 -> 不过闸
        good = _quote_json()
        script = [outline, short, good, _page_json("P2")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("金句收束", pages=2, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        st = ppt_maker.last_run["page_stats"][0]
        self.assertEqual(st["rewrites"], 1)
        self.assertEqual(st["layout"], "quote")
        self.assertEqual(len(mock.calls), 1 + 2 + 1)

    # 16. 总页数 <12 时 section 被降级为 bullets（选版式规则的规范化兜底）
    def test_section_downgraded_when_deck_short(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        outline = json.dumps({
            "title": "短 deck", "pages": [
                {"layout": "section", "page_title": "第一章",
                 "core_point": "章节页在短 deck 里不允许，应降级为要点页。",
                 "keywords": ["降级"]},
                {"layout": "bullets", "page_title": "正文",
                 "core_point": "正文页正常精写，给出具体数据与机制解释。",
                 "keywords": ["正文"]},
                {"layout": "bullets", "page_title": "收尾",
                 "core_point": "收尾页升级为 closing，给出行动建议。",
                 "keywords": ["行动"]},
            ]}, ensure_ascii=False)
        script = [outline, _page_json("S1"), _page_json("P2"), _page_json("P3")]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("短 deck", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        deck = render.calls[0][1]
        self.assertEqual(deck["pages"][0]["layout"], "bullets")      # section 降级
        self.assertEqual(deck["pages"][2]["layout"], "closing")      # 尾页升级 closing
        # 降级后走 bullets 精写路径：3 页各 1 次调用
        self.assertEqual(len(mock.calls), 1 + 3)

    # 17. chart 归一化单元：字符串数字转 float、截齐补齐、非法 type 回 bar
    def test_norm_chart_coercion(self):
        chart = ppt_maker._norm_chart({
            "type": "AREA",                                   # 非法 -> bar
            "title": "t",
            "categories": ["a", "b", "c"],
            "series": [
                {"name": "s1", "values": ["1,200", "3.5%", "x", "9"]},  # 转 float + 截到 3
                {"name": "s2", "values": [1]},                          # 短了补 0.0
            ]})
        self.assertEqual(chart["type"], "bar")
        self.assertEqual(chart["series"][0]["values"], [1200.0, 3.5, 0.0])
        self.assertEqual(chart["series"][1]["values"], [1.0, 0.0, 0.0])
        for s in chart["series"]:
            for v in s["values"]:
                self.assertIsInstance(v, float)

    # 18. stats 归一化单元：钳 1-3 个、空壳丢弃
    def test_norm_stats_clamped(self):
        stats = ppt_maker._norm_stats(
            [{"number": f"{i}%", "caption": f"c{i}"} for i in range(5)]
            + [{"number": "", "caption": ""}, "不是字典"])
        self.assertEqual(len(stats), 3)
        self.assertEqual(stats[0], {"number": "0%", "caption": "c0"})

    # ================================================================ 新增：路 B 生图管线

    # 19. 配图 prompt 被收集、风格后缀被拼接；成功 -> 绝对路径 + 文件真实落盘
    def test_images_collected_suffixed_and_saved(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP(dl_ctypes={1: "image/jpeg"})   # 封面下载给 jpeg -> 后缀 .jpg
        self._enable_images(http)
        # 3 页全 bullets：末页归一化升级 closing -> 只有前 2 页是配图候选
        script = [_image_outline_json(3)] + [_page_json(f"P{i}") for i in range(1, 4)]
        mock = _MockLLM(script)
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))

        deck = render.calls[0][1]
        # API 调用 = 1 封面 + 2 内容页，prompt 含原始描述且拼接了统一风格后缀
        bodies = http.api_bodies()
        self.assertEqual(len(bodies), 3)
        self.assertIn("vast abstract landscape", bodies[0]["prompt"])
        self.assertIn("warm scene of concept 1", bodies[1]["prompt"])
        self.assertIn("warm scene of concept 2", bodies[2]["prompt"])
        for b in bodies:
            self.assertTrue(b["prompt"].endswith(ppt_maker.IMAGE_STYLE_SUFFIX))
            self.assertEqual(b["model"], "cogview-3")
        # 封面：绝对路径、.jpg 后缀（按 content-type）、文件真实落盘
        cover = deck["cover_image"]
        self.assertTrue(cover)
        self.assertTrue(Path(cover).is_absolute())
        self.assertTrue(cover.endswith(".jpg"))
        self.assertTrue(Path(cover).is_file())
        self.assertEqual(Path(cover).parent, self.tmp / "ppt_assets")
        # 内容页：前两页有图（.png），closing 页不给图（无 image 键）
        for pg in deck["pages"][:2]:
            self.assertEqual(pg["layout"], "bullets")
            self.assertTrue(pg["image"])
            self.assertTrue(Path(pg["image"]).is_absolute())
            self.assertTrue(pg["image"].endswith(".png"))
            self.assertTrue(Path(pg["image"]).is_file())
        self.assertEqual(deck["pages"][2]["layout"], "closing")
        self.assertNotIn("image", deck["pages"][2])
        self.assertEqual(ppt_maker.last_run["images"], 3)

    # 20. 生图 API 失败 -> 该张字段 None，其余照常，整单不失败
    def test_image_api_failure_degrades(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP(api_fails={2})       # 第 2 次 API（第 1 内容页）炸
        self._enable_images(http)
        script = [_image_outline_json(3)] + [_page_json(f"P{i}") for i in range(1, 4)]
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        deck = render.calls[0][1]
        self.assertTrue(deck["cover_image"])              # 封面成功
        self.assertIsNone(deck["pages"][0]["image"])      # 失败页 None
        self.assertTrue(deck["pages"][1]["image"])        # 后续页不受影响
        self.assertEqual(ppt_maker.last_run["images"], 2)
        # 失败页文件不落盘：ppt_assets 里只有 2 张
        self.assertEqual(len(list((self.tmp / "ppt_assets").glob("pptimg_*"))), 2)

    # 21. 下载失败 -> 字段 None，流程继续
    def test_image_download_failure_degrades(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP(dl_fails={1})        # 封面下载炸
        self._enable_images(http)
        script = [_image_outline_json(3)] + [_page_json(f"P{i}") for i in range(1, 4)]
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        deck = render.calls[0][1]
        self.assertIsNone(deck["cover_image"])            # 封面降级无图
        self.assertTrue(deck["pages"][0]["image"])        # 内容页照常
        self.assertTrue(deck["pages"][1]["image"])
        self.assertEqual(ppt_maker.last_run["images"], 2)

    # 22. 配图页数 >8 -> 钳制到 8（按页序保留前 8；N2 扩面后上限 4→8）
    def test_image_pages_clamped_to_eight(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP()
        self._enable_images(http)
        # 10 页全 bullets：末页升级 closing -> 9 个候选，钳到前 8
        script = [_image_outline_json(10)] + [_page_json(f"P{i}") for i in range(1, 11)]
        r = ppt_maker.make_ppt("人工智能简介", pages=10, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        deck = render.calls[0][1]
        # API 调用 = 1 封面 + 恰好 8 内容页
        self.assertEqual(len(http.api_bodies()), 9)
        for pg in deck["pages"][:8]:
            self.assertTrue(pg["image"], f"{pg['page_title']} 应有图")
        self.assertIsNone(deck["pages"][8]["image"])      # 第 9 页被钳掉
        self.assertEqual(deck["pages"][9]["layout"], "closing")
        self.assertEqual(ppt_maker.last_run["images"], 9)  # 封面 + 8 页

    # 23. 非 bullets 页的 image_prompt 被静默丢弃
    def test_non_bullets_image_prompt_dropped(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP()
        self._enable_images(http)
        outline = _image_outline_json(
            4, layouts=["two_column", "quote", "bullets", "closing"])
        script = [outline, _two_column_json(), _quote_json(),
                  _page_json("P3"), _page_json("P4")]
        r = ppt_maker.make_ppt("人工智能简介", pages=4, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        deck = render.calls[0][1]
        # 只有 bullets 页（第 3 页）配图：API = 1 封面 + 1 内容页
        self.assertEqual(len(http.api_bodies()), 2)
        self.assertEqual(deck["pages"][0]["layout"], "two_column")
        self.assertNotIn("image", deck["pages"][0])
        self.assertEqual(deck["pages"][1]["layout"], "quote")
        self.assertNotIn("image", deck["pages"][1])
        self.assertTrue(deck["pages"][2]["image"])
        self.assertNotIn("image", deck["pages"][3])       # closing 不配图

    # 24. with_images=False -> 零 HTTP 调用，契约键初始化为 None
    def test_with_images_false_zero_http(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP()
        self._enable_images(http)
        script = [_image_outline_json(3)] + [_page_json(f"P{i}") for i in range(1, 4)]
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=_MockLLM(script),
                               with_images=False)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(http.calls, [])                  # 零 HTTP
        deck = render.calls[0][1]
        self.assertIsNone(deck["cover_image"])
        for pg in deck["pages"][:2]:
            self.assertIsNone(pg["image"])
        self.assertEqual(ppt_maker.last_run["images"], 0)
        self.assertFalse((self.tmp / "ppt_assets").exists())

    # 25. 配置缺失（无 key）-> 整条生图链静默跳过、零 HTTP（setUp 默认即无配置）
    def test_missing_config_silent_skip(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP()
        orig_open = urllib.request.urlopen
        urllib.request.urlopen = http
        self.addCleanup(lambda: setattr(urllib.request, "urlopen", orig_open))
        script = [_image_outline_json(3)] + [_page_json(f"P{i}") for i in range(1, 4)]
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        self.assertEqual(http.calls, [])
        deck = render.calls[0][1]
        self.assertIsNone(deck["cover_image"])
        for pg in deck["pages"][:2]:
            self.assertIsNone(pg["image"])
        self.assertEqual(ppt_maker.last_run["images"], 0)

    # 26. 大纲不给 cover_image_prompt -> 封面无图但内容页配图照常
    def test_missing_cover_prompt_skips_cover_only(self):
        render = _MockRender()
        ppt_maker._render_deck = render
        http = _MockHTTP()
        self._enable_images(http)
        outline = json.loads(_image_outline_json(3))
        del outline["cover_image_prompt"]
        script = [json.dumps(outline, ensure_ascii=False)] + [_page_json(f"P{i}") for i in range(1, 4)]
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=_MockLLM(script))
        self.assertTrue(r["ok"], r.get("error"))
        deck = render.calls[0][1]
        self.assertIsNone(deck["cover_image"])
        self.assertEqual(len(http.api_bodies()), 2)       # 只有 2 个内容页
        self.assertTrue(deck["pages"][0]["image"])
        self.assertTrue(deck["pages"][1]["image"])

    # ================================================================ 新增：叙事升级（R2）

    # 27. 大纲 prompt 含结论式标题（action title）与叙事弧线规则
    def test_outline_prompt_has_action_title_rules(self):
        mock = _MockLLM(_script_for(3))
        r = ppt_maker.make_ppt("人工智能简介", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        outline_prompt = mock.calls[0]
        self.assertIn("断言", outline_prompt)                # 结论式标题规则关键词
        self.assertIn("action title", outline_prompt)
        self.assertIn("叙事弧线", outline_prompt)            # 情境→冲突→分析→行动
        self.assertIn("为什么现在该关心这件事", outline_prompt)

    # 28. page_title 截断放宽到 22 字（断言句天然更长）
    def test_page_title_truncated_at_22(self):
        long_title = "这是一段超过二十二字长度的断言式页面标题用来验证截断逻辑是否生效"
        outline = json.dumps({
            "title": "t",
            "pages": [{"page_title": long_title,
                       "core_point": "核心论点：验证标题截断行为是否符合预期。",
                       "keywords": ["截断"]}],
        }, ensure_ascii=False)
        data, err = ppt_maker._gen_outline("t", 3, "工作汇报", _MockLLM([outline]))
        self.assertIsNone(err)
        self.assertEqual(data["pages"][0]["page_title"], long_title[:22])
        self.assertEqual(len(data["pages"][0]["page_title"]), 22)

    # ================================================================ 新增：联网研究接入（R1 集成）

    # 29. 研究材料注入大纲与精写 prompt；空串时 prompt 不含研究段
    def test_research_injected_into_prompts(self):
        ppt_maker._research_topic = lambda topic, **kw: "2025年产销超1600万辆（中汽协，2026.1）"
        mock = _MockLLM(_script_for(3))
        r = ppt_maker.make_ppt("新能源车", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertIn("真实研究材料", mock.calls[0])          # 大纲 prompt 带研究段
        self.assertIn("1600万辆", mock.calls[0])              # 材料原文进去了
        self.assertIn("真实研究材料", mock.calls[1])          # 精写 prompt 同样带
        self.assertEqual(ppt_maker.last_run["research_chars"], len("2025年产销超1600万辆（中汽协，2026.1）"))

    # 30. 研究模块缺席/异常 -> 空串降级，prompt 无研究段，流程不受影响的
    def test_research_absent_degrades_cleanly(self):
        ppt_maker._research_topic = None                       # setUp 已置 None，显式重申语义
        mock = _MockLLM(_script_for(3))
        r = ppt_maker.make_ppt("新能源车", pages=3, llm_caller=mock)
        self.assertTrue(r["ok"], r.get("error"))
        self.assertNotIn("真实研究材料", mock.calls[0])
        self.assertEqual(ppt_maker.last_run["research_chars"], 0)


if __name__ == "__main__":
    unittest.main(verbosity=2)
