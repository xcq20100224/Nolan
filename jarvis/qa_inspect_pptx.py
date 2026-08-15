# -*- coding: utf-8 -*-
"""QA 体检：解析 Nolan 生成的 pptx，逐项量化体检指标。"""
import json
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

PPTX = sys.argv[1]
prs = Presentation(PPTX)
EMU_PER_IN = 914400
slide_w_in = prs.slide_width / EMU_PER_IN
slide_h_in = prs.slide_height / EMU_PER_IN

report = {
    "file": str(PPTX),
    "slide_size_in": [round(slide_w_in, 2), round(slide_h_in, 2)],
    "page_count": len(prs.slides),
    "pages": [],
    "images": [],
    "charts": 0,
    "notes_missing_pages": [],
    "notes_chars_total": 0,
    "font_sizes_seen": set(),
    "overflow_risks": [],
    "empty_text_pages": [],
}

zf = zipfile.ZipFile(PPTX)
names = zf.namelist()
chart_parts = [n for n in names if n.startswith("ppt/charts/chart")]
report["charts"] = len(chart_parts)
media = [n for n in names if n.startswith("ppt/media/")]
report["media_files"] = len(media)

for i, slide in enumerate(prs.slides, 1):
    page = {"page": i, "shapes": len(slide.shapes), "texts": [],
            "font_sizes": [], "pictures": 0, "has_notes": False,
            "notes_chars": 0, "total_chars": 0}
    title = ""
    for sh in slide.shapes:
        if sh.shape_type == 13:  # PICTURE
            page["pictures"] += 1
            report["images"].append({
                "page": i,
                "w_in": round(sh.width / EMU_PER_IN, 2),
                "h_in": round(sh.height / EMU_PER_IN, 2)})
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt:
                page["total_chars"] += len(txt)
                if not title:
                    title = txt.split("\n")[0][:40]
                for para in sh.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            pt = run.font.size.pt
                            page["font_sizes"].append(pt)
                            report["font_sizes_seen"].add(pt)
                        if run.text.strip():
                            page["texts"].append(run.text.strip()[:30])
        # 文字溢出风险：文本框字数 vs 框容量（按字号估算）
        if sh.has_text_frame:
            tf = sh.text_frame
            txt = tf.text.strip()
            if not txt:
                continue
            box_w = sh.width / EMU_PER_IN if sh.width else 0
            box_h = sh.height / EMU_PER_IN if sh.height else 0
            sizes = [r.font.size.pt for p in tf.paragraphs for r in p.runs if r.font.size]
            pt = max(sizes) if sizes else 18
            # 中文每字宽度≈字号pt/72 英寸；行高≈1.3×字号
            char_w = pt / 72.0
            line_h = pt * 1.35 / 72.0
            chars_per_line = max(1, int(box_w / char_w)) if box_w else 1
            lines_needed = 0
            for line in txt.split("\n"):
                lines_needed += max(1, -(-len(line) // chars_per_line))
            height_needed = lines_needed * line_h
            if box_h and height_needed > box_h * 1.05:
                report["overflow_risks"].append({
                    "page": i, "text_head": txt[:24], "font_pt": pt,
                    "box_h_in": round(box_h, 2), "need_h_in": round(height_needed, 2),
                    "chars": len(txt)})
    if slide.has_notes_slide:
        nt = slide.notes_slide.notes_text_frame.text.strip()
        if nt:
            page["has_notes"] = True
            page["notes_chars"] = len(nt)
            report["notes_chars_total"] += len(nt)
        else:
            report["notes_missing_pages"].append(i)
    else:
        report["notes_missing_pages"].append(i)
    if page["total_chars"] == 0 and page["pictures"] == 0:
        report["empty_text_pages"].append(i)
    page["title"] = title
    page["font_sizes"] = sorted(set(page["font_sizes"]))
    report["pages"].append(page)

report["font_sizes_seen"] = sorted(report["font_sizes_seen"])
print(json.dumps(report, ensure_ascii=False, indent=1))
